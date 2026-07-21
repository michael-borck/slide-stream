"""Tests for the `draft` command: document -> slide-deck Markdown."""

import pytest
from typer.testing import CliRunner

from slide_stream.cli import app
from slide_stream.draft import (
    DraftError,
    build_draft_prompt,
    clamp_source,
    clean_llm_markdown,
    extract_source_text,
    validate_deck_markdown,
)

runner = CliRunner()

_DECK_MD = "# Intro\n\n- Point one\n- Point two\n\n# Details\n\n- More\n"


@pytest.fixture(autouse=True)
def _no_home_config(monkeypatch):
    """Isolate from a real ~/.slidestream.yaml on the dev machine."""
    monkeypatch.setattr("slide_stream.config_loader.find_home_config", lambda: None)


# --- Text extraction ----------------------------------------------------------


def test_extract_txt(tmp_path):
    src = tmp_path / "notes.txt"
    src.write_text("Hello world.\nSecond line.")
    assert "Hello world." in extract_source_text(src)


def test_extract_markdown(tmp_path):
    src = tmp_path / "doc.md"
    src.write_text("# Heading\n\nBody text.")
    assert "Body text." in extract_source_text(src)


def test_extract_unsupported_suffix(tmp_path):
    src = tmp_path / "sheet.xlsx"
    src.write_text("x")
    with pytest.raises(DraftError, match="Unsupported source type"):
        extract_source_text(src)


def test_extract_pptx_roundtrips_titles_and_notes(tmp_path):
    """A .pptx source is flattened to titles/bullets/notes for the LLM."""
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # title + content
    slide.shapes.title.text = "Photosynthesis"  # type: ignore[union-attr]
    slide.placeholders[1].text = "Light reactions"  # type: ignore[attr-defined]
    slide.notes_slide.notes_text_frame.text = "Mention chlorophyll"  # type: ignore[union-attr]
    pptx = tmp_path / "in.pptx"
    prs.save(pptx)

    text = extract_source_text(pptx)
    assert "# Photosynthesis" in text
    assert "Light reactions" in text
    assert "Mention chlorophyll" in text


# --- Prompt / output helpers --------------------------------------------------


def test_prompt_pins_exact_slide_count():
    assert "exactly 7 slides" in build_draft_prompt("source", 7)


def test_prompt_lets_llm_choose_when_no_count():
    prompt = build_draft_prompt("source", None)
    assert "Choose a sensible number" in prompt


def test_clamp_source_truncates_long_input():
    text = "x" * 200_000
    clamped, truncated = clamp_source(text)
    assert truncated is True
    assert len(clamped) < len(text)


def test_clamp_source_leaves_short_input():
    clamped, truncated = clamp_source("short")
    assert (clamped, truncated) == ("short", False)


def test_clean_llm_markdown_strips_code_fence():
    fenced = "```markdown\n# Title\n\n- a\n```"
    assert clean_llm_markdown(fenced) == "# Title\n\n- a"


def test_validate_deck_rejects_empty():
    with pytest.raises(DraftError, match="no '# ' headings"):
        validate_deck_markdown("just prose, no headings")


def test_validate_deck_accepts_real_deck():
    slides = validate_deck_markdown(_DECK_MD)
    assert [s["title"] for s in slides] == ["Intro", "Details"]


# --- CLI ----------------------------------------------------------------------


def test_cli_draft_writes_markdown(tmp_path, monkeypatch, mocker):
    monkeypatch.chdir(tmp_path)
    src = tmp_path / "report.txt"
    src.write_text("A long report about widgets and gadgets.")

    mocker.patch("slide_stream.cli.get_llm_client", return_value=mocker.MagicMock())
    mocker.patch("slide_stream.cli.query_llm", return_value=_DECK_MD)

    result = runner.invoke(app, ["draft", str(src), "--llm-provider", "openai"])
    assert result.exit_code == 0, result.output
    out = tmp_path / "report.md"
    assert out.exists()
    assert "# Intro" in out.read_text()


def test_cli_draft_honours_output_arg_and_slides(tmp_path, monkeypatch, mocker):
    monkeypatch.chdir(tmp_path)
    src = tmp_path / "report.txt"
    src.write_text("content")

    mocker.patch("slide_stream.cli.get_llm_client", return_value=mocker.MagicMock())
    q = mocker.patch("slide_stream.cli.query_llm", return_value=_DECK_MD)

    result = runner.invoke(
        app, ["draft", str(src), "deck.md", "--slides", "5", "--llm-provider", "openai"]
    )
    assert result.exit_code == 0, result.output
    assert (tmp_path / "deck.md").exists()
    # The exact-count instruction reached the prompt.
    assert "exactly 5 slides" in q.call_args[0][2]


def test_cli_draft_requires_llm_provider(tmp_path, monkeypatch):
    monkeypatch.chdir(tmp_path)
    src = tmp_path / "report.txt"
    src.write_text("content")

    # No --llm-provider and default provider is 'none' -> clean error.
    result = runner.invoke(app, ["draft", str(src)])
    assert result.exit_code == 1
    assert "draft needs an LLM" in result.output


def test_cli_draft_refuses_overwrite_without_force(tmp_path, monkeypatch, mocker):
    monkeypatch.chdir(tmp_path)
    src = tmp_path / "report.txt"
    src.write_text("content")
    (tmp_path / "report.md").write_text("existing")

    mocker.patch("slide_stream.cli.get_llm_client", return_value=mocker.MagicMock())
    mocker.patch("slide_stream.cli.query_llm", return_value=_DECK_MD)

    result = runner.invoke(app, ["draft", str(src), "--llm-provider", "openai"])
    assert result.exit_code == 1
    assert "already exists" in result.output
    assert (tmp_path / "report.md").read_text() == "existing"


def test_cli_draft_missing_source(tmp_path, monkeypatch):
    monkeypatch.chdir(tmp_path)
    result = runner.invoke(app, ["draft", "nope.txt", "--llm-provider", "openai"])
    assert result.exit_code == 1
    assert "not found" in result.output
