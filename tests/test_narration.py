"""Tests for narration planning, vision narration, and CLI LLM options."""

import io

from PIL import Image
from pptx import Presentation
from pptx.util import Inches
from typer.testing import CliRunner

from slide_stream.cli import app
from slide_stream.llm import query_llm_with_image
from slide_stream.narration import (
    build_narration_prompt,
    narration_source,
    strip_stage_directions,
    target_words,
)
from slide_stream.powerpoint import parse_powerpoint

# --- Stage-direction stripping ------------------------------------------------


def test_strip_removes_square_bracket_directions():
    assert strip_stage_directions("Welcome [pause] everyone.") == "Welcome everyone."


def test_strip_removes_parenthetical_direction():
    assert strip_stage_directions("Say hello (pause here) then go.") == "Say hello then go."


def test_strip_preserves_bare_direction_words():
    # A prose word that merely contains a cue keyword must not be deleted.
    text = "We must emphasize safety at all times."
    assert strip_stage_directions(text) == text


def test_strip_tidies_punctuation_and_spacing():
    assert strip_stage_directions("Intro [pause], then begin.") == "Intro, then begin."


def test_strip_returns_unchanged_when_no_directions():
    text = "A perfectly ordinary sentence."
    assert strip_stage_directions(text) is text  # early-out, same object


def test_strip_handles_note_that_is_only_a_direction():
    assert strip_stage_directions("[click to advance]") == ""


def test_clean_narration_strips_directions_from_notes():
    """The default (no-LLM) narration builder must not voice bracketed cues."""
    from slide_stream.cli import _clean_narration

    spoken = _clean_narration("My Title", ["a bullet"], "Explain the chart [pause] slowly.")
    assert "[pause]" not in spoken
    assert "Explain the chart slowly." in spoken

# --- Source selection ---------------------------------------------------------


def test_notes_take_priority_over_content():
    slide = {"title": "T", "content": ["a"], "notes": "the real story", "images": [{}]}
    assert narration_source(slide) == "notes"


def test_content_used_when_no_notes():
    assert narration_source({"title": "T", "content": ["a"], "notes": ""}) == "content"


def test_image_used_when_no_text():
    slide = {"title": "T", "content": [], "notes": "", "images": [{"data": b"x"}]}
    assert narration_source(slide) == "image"


def test_title_only_fallback():
    assert narration_source({"title": "T", "content": [], "notes": ""}) == "title-only"


def test_placeholder_notes_do_not_count():
    assert narration_source({"title": "T", "content": ["a"], "notes": "  "}) == "content"


# --- Length targets -----------------------------------------------------------


def test_target_words_scales_with_duration():
    assert target_words(60, wpm=150) == 150
    assert target_words(30, wpm=150) == 75
    assert target_words(None) is None
    assert target_words(1) == 10  # floor


# --- Prompt construction --------------------------------------------------------


def test_notes_prompt_marks_notes_as_primary():
    slide = {"title": "Neurons", "content": ["Axons"], "notes": "Talk about dendrites"}
    prompt = build_narration_prompt(slide, "notes", 75)
    assert "Talk about dendrites" in prompt
    assert "PRIMARY SOURCE" in prompt
    assert "approximately 75 words" in prompt
    assert "30 seconds" in prompt


def test_content_prompt_forbids_reading_bullets():
    slide = {"title": "Neurons", "content": ["Axons", "Dendrites"], "notes": ""}
    prompt = build_narration_prompt(slide, "content", None)
    assert "- Axons" in prompt
    assert "do not recite the bullets" in prompt
    assert "complement the slide" in prompt


def test_image_prompt_relates_to_title():
    prompt = build_narration_prompt({"title": "Brain scan", "content": []}, "image", 75)
    assert "Brain scan" in prompt
    assert "image" in prompt.lower()


# --- PowerPoint image extraction ------------------------------------------------


def make_png_bytes(color=(200, 40, 40)):
    img = Image.new("RGB", (64, 48), color=color)
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def test_parse_powerpoint_extracts_images(tmp_path):
    """An image-only slide is kept and carries its picture bytes."""
    png = tmp_path / "pic.png"
    png.write_bytes(make_png_bytes())

    prs = Presentation()
    blank = prs.slide_layouts[6]  # blank layout: no title placeholder
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(str(png), Inches(1), Inches(1), Inches(3), Inches(2))
    pptx_path = tmp_path / "deck.pptx"
    prs.save(pptx_path)

    slides = parse_powerpoint(pptx_path)

    assert len(slides) == 1
    assert slides[0]["content"] == []
    assert len(slides[0]["images"]) == 1
    assert slides[0]["images"][0]["content_type"] == "image/png"
    assert slides[0]["images"][0]["data"].startswith(b"\x89PNG")


# --- Vision LLM plumbing ---------------------------------------------------------


def test_query_llm_with_image_claude_message_shape(mocker):
    from rich.console import Console

    fake_client = mocker.MagicMock()
    fake_block = mocker.MagicMock()
    fake_block.text = "A narration about the image."
    fake_client.messages.create.return_value.content = [fake_block]

    result = query_llm_with_image(
        fake_client, "claude", "Narrate this.", b"imagebytes", "image/png",
        Console(), model="claude-haiku-4-5",
    )

    assert result == "A narration about the image."
    kwargs = fake_client.messages.create.call_args[1]
    assert kwargs["model"] == "claude-haiku-4-5"
    content = kwargs["messages"][0]["content"]
    assert content[0]["type"] == "image"
    assert content[0]["source"]["media_type"] == "image/png"
    assert content[1]["type"] == "text"


def test_query_llm_with_image_rejects_non_vision_provider(mocker):
    from rich.console import Console

    result = query_llm_with_image(
        mocker.MagicMock(), "groq", "p", b"x", "image/png", Console()
    )
    assert result is None


# --- CLI options ------------------------------------------------------------------


def test_cli_llm_provider_flag_unknown_provider_fails(tmp_path, monkeypatch):
    monkeypatch.chdir(tmp_path)
    runner = CliRunner()
    md = tmp_path / "d.md"
    md.write_text("# One\n\n- a\n")

    result = runner.invoke(
        app, ["create", str(md), "out.mp4", "--llm-provider", "not-a-provider"]
    )
    assert result.exit_code == 1


def test_cli_llm_provider_claude_without_key_fails_cleanly(tmp_path, monkeypatch):
    monkeypatch.chdir(tmp_path)
    monkeypatch.delenv("ANTHROPIC_API_KEY", raising=False)
    runner = CliRunner()
    md = tmp_path / "d.md"
    md.write_text("# One\n\n- a\n")

    result = runner.invoke(
        app, ["create", str(md), "out.mp4", "--llm-provider", "claude"]
    )
    assert result.exit_code == 1


# --- Script file parsing --------------------------------------------------------


def test_parse_script_file_splits_on_dashes(tmp_path):
    from slide_stream.narration import parse_script_file

    script = tmp_path / "script.txt"
    script.write_text(
        "First slide narration.\nStill slide one.\n"
        "---\n"
        "Second slide narration.\n"
        "-----\n"
        "Third.\n"
    )
    blocks = parse_script_file(script)
    assert blocks == [
        "First slide narration.\nStill slide one.",
        "Second slide narration.",
        "Third.",
    ]


def test_parse_script_file_single_block_without_separators(tmp_path):
    from slide_stream.narration import parse_script_file

    script = tmp_path / "s.txt"
    script.write_text("Just one block.\n")
    assert parse_script_file(script) == ["Just one block."]


def test_parse_script_file_preserves_blank_blocks(tmp_path):
    """A blank block keeps slide alignment intact."""
    from slide_stream.narration import parse_script_file

    script = tmp_path / "s.txt"
    script.write_text("One.\n---\n---\nThree.\n")
    assert parse_script_file(script) == ["One.", "", "Three."]


def test_parse_script_file_trailing_separator_ignored(tmp_path):
    """A trailing --- must not append a spurious empty final block (which
    would trigger a bogus slide-count-mismatch warning)."""
    from slide_stream.narration import parse_script_file

    script = tmp_path / "s.txt"
    script.write_text("One.\n---\nTwo.\n---\n")
    assert parse_script_file(script) == ["One.", "Two."]


def test_parse_script_file_trailing_separator_keeps_interior_blanks(tmp_path):
    """Dropping the trailing empty block must not disturb interior blanks."""
    from slide_stream.narration import parse_script_file

    script = tmp_path / "s.txt"
    script.write_text("One.\n---\n---\nThree.\n---\n")
    assert parse_script_file(script) == ["One.", "", "Three."]


def test_hyphenated_text_is_not_a_separator(tmp_path):
    """A line with dashes plus other text is content, not a separator."""
    from slide_stream.narration import parse_script_file

    script = tmp_path / "s.txt"
    script.write_text("Well-being matters.\n-- but this too\n")
    assert parse_script_file(script) == ["Well-being matters.\n-- but this too"]


def test_cli_script_file_is_used_verbatim(tmp_path, mocker, monkeypatch):
    """--script narration is spoken verbatim (no LLM, no key needed)."""
    import wave

    monkeypatch.chdir(tmp_path)
    runner = CliRunner()

    md = tmp_path / "deck.md"
    md.write_text("# One\n\n- a\n\n# Two\n\n- b\n")
    script = tmp_path / "script.txt"
    script.write_text("Custom narration for slide one.\n---\nAnd slide two.\n")
    cfg = tmp_path / "config.yaml"
    cfg.write_text(
        "settings:\n  video:\n    resolution: [320, 180]\n    fps: 8\n"
        "    default_slide_duration: 1.0\n"
    )

    spoken = []

    def fake_save(filename):
        spoken.append(True)
        with wave.open(str(filename), "wb") as w:
            w.setnchannels(1)
            w.setsampwidth(2)
            w.setframerate(22050)
            w.writeframes(b"\x00\x00" * 8000)

    fake_tts = mocker.MagicMock()
    fake_tts.save.side_effect = fake_save
    gtts = mocker.patch("gtts.gTTS", return_value=fake_tts)

    result = runner.invoke(
        app,
        ["create", str(md), str(tmp_path / "out.mp4"), "--config", str(cfg), "--script", str(script)],
    )

    assert result.exit_code == 0, result.output
    # gTTS was called with the verbatim script text, not slide bullets.
    spoken_texts = [c.kwargs.get("text", "") for c in gtts.call_args_list]
    assert "Custom narration for slide one." in spoken_texts
    assert "And slide two." in spoken_texts


def test_cli_tts_overrides_reach_the_provider(tmp_path, monkeypatch):
    """--tts-provider/--tts-base-url/--voice select a Chatterbox server; the
    provider must be initialised from them (verified via the console output)."""
    monkeypatch.chdir(tmp_path)
    runner = CliRunner()
    md = tmp_path / "d.md"
    md.write_text("# One\n\n- a\n")

    # chatterbox with a base_url is 'available'; --strict makes an unreachable
    # server abort at synthesis rather than silently using gtts — proving the
    # override took effect (the default has no base_url, so it never reaches a
    # server).
    result = runner.invoke(
        app,
        [
            "create", str(md), str(tmp_path / "out.mp4"),
            "--tts-provider", "chatterbox",
            "--tts-base-url", "http://127.0.0.1:1/v1",
            "--voice", "Emily.wav",
            "--strict",
        ],
    )
    assert result.exit_code == 1
    assert "chatterbox" in result.output.lower()
