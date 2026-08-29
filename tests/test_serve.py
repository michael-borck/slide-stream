"""Tests for the web UI (slide-stream serve)."""

import copy
import io
import time
from io import BytesIO
from pathlib import Path

import pytest
import yaml

pytest.importorskip("fastapi")
from fastapi.testclient import TestClient  # noqa: E402

from slide_stream import serve  # noqa: E402
from slide_stream.config_loader import DEFAULT_CONFIG  # noqa: E402


def _tiny_png() -> bytes:
    """A real (decodable) 4x4 PNG for photo uploads."""
    from PIL import Image

    buf = BytesIO()
    Image.new("RGB", (4, 4), "red").save(buf, format="PNG")
    return buf.getvalue()


@pytest.fixture
def base_config():
    return copy.deepcopy(DEFAULT_CONFIG)


@pytest.fixture(autouse=True)
def _clear_jobs():
    serve._JOBS.clear()
    yield
    serve._JOBS.clear()


def test_index_served(base_config):
    client = TestClient(serve.create_app(config=base_config, token="secret"))
    r = client.get("/")
    assert r.status_code == 200
    assert "SlideStream" in r.text


def test_api_requires_token(base_config):
    client = TestClient(serve.create_app(config=base_config, token="secret"))
    # /api/config is public so the UI can bootstrap...
    cfg = client.get("/api/config")
    assert cfg.status_code == 200
    assert cfg.json()["auth_required"] is True
    # ...but protected endpoints reject a missing token.
    assert client.get("/api/jobs/whatever").status_code == 401
    ok = client.get("/api/jobs/whatever", headers={"Authorization": "Bearer secret"})
    assert ok.status_code == 404  # authorized, just no such job


def test_no_token_means_open(base_config):
    client = TestClient(serve.create_app(config=base_config, token=None))
    assert client.get("/api/config").status_code == 200
    assert client.get("/api/jobs/whatever").status_code == 404


def test_demo_flag_exposed(base_config):
    client = TestClient(serve.create_app(config=base_config, token=None, demo=True))
    assert client.get("/api/config").json()["demo"] is True
    off = TestClient(serve.create_app(config=base_config, token=None, demo=False))
    assert off.get("/api/config").json()["demo"] is False


def test_create_job_rejects_bad_deck_type(base_config):
    client = TestClient(serve.create_app(config=base_config, token=None))
    r = client.post("/api/jobs", files={"deck": ("notes.docx", b"hi", "text/plain")})
    assert r.status_code == 400


def test_job_lifecycle_and_download(base_config, tmp_path, monkeypatch):
    """Submit a deck; the (mocked) render produces a video that downloads."""

    def fake_run_job(job, deck_path, job_yaml, voice_path, photo_path,
                     mode="video", notes=None):
        # Simulate a successful render writing output.mp4.
        assert job.workdir is not None
        out = job.workdir / "output.mp4"
        out.write_bytes(b"FAKEMP4")
        job.status = "done"
        job.output_path = out
        job.log = "rendered"
        # Ephemeral inputs are wiped by the real runner; emulate it.
        for p in (voice_path, photo_path):
            if p is not None:
                p.unlink(missing_ok=True)

    monkeypatch.setattr(serve, "_run_job", fake_run_job)
    # Run submitted work synchronously so the test is deterministic.
    monkeypatch.setattr(
        serve.ThreadPoolExecutor, "submit",
        lambda self, fn, *a, **k: fn(*a, **k),
    )

    client = TestClient(serve.create_app(config=base_config, token=None))
    r = client.post(
        "/api/jobs",
        files={
            "deck": ("deck.md", b"# One\n\n- a\n", "text/markdown"),
            "voice": ("me.wav", b"AUDIO", "audio/wav"),
            "photo": ("me.png", _tiny_png(), "image/png"),
        },
        data={"avatar": "true", "narration_seconds": "30"},
    )
    assert r.status_code == 200
    job_id = r.json()["job_id"]
    token = r.json()["token"]
    assert token  # per-job download token minted at creation

    status = client.get(f"/api/jobs/{job_id}").json()
    assert status["status"] == "done"
    assert status["token"] == token  # non-demo status echoes it for the UI

    # The download needs the per-job token, not just the job id.
    assert client.get(f"/api/jobs/{job_id}/result").status_code == 401
    result = client.get(f"/api/jobs/{job_id}/result", params={"t": token})
    assert result.status_code == 200
    assert result.content == b"FAKEMP4"


def test_api_check_returns_doctor_report(base_config):
    """The preflight endpoint parses the deck + config and returns findings."""
    client = TestClient(serve.create_app(config=base_config, token=None))
    r = client.post(
        "/api/check",
        files={"deck": ("deck.md", b"# One\n\n- a\n\n# Two\n\n- b\n", "text/markdown")},
        data={"output": "video"},
    )
    assert r.status_code == 200
    body = r.json()
    assert "findings" in body and "estimates" in body
    assert any("2 slide(s)" in f["message"] for f in body["findings"])
    assert any("Video length" in e for e in body["estimates"])


def test_pptx_output_runs_enrich(base_config, monkeypatch):
    """Output=pptx routes the job through enrich (zip), not create (mp4)."""
    seen = {}

    def fake_run_job(job, deck_path, job_yaml, voice_path, photo_path,
                     mode="video", notes=None):
        seen["mode"] = mode
        seen["notes"] = notes
        assert job.workdir is not None
        out = job.workdir / "enriched.zip"
        out.write_bytes(b"PK\x03\x04zip")
        job.status = "done"
        job.output_path = out
        job.media_type = "application/zip"
        job.download_name = "slidestream-deck.zip"

    monkeypatch.setattr(serve, "_run_job", fake_run_job)
    monkeypatch.setattr(
        serve.ThreadPoolExecutor, "submit", lambda self, fn, *a, **k: fn(*a, **k)
    )
    client = TestClient(serve.create_app(config=base_config, token=None))
    r = client.post(
        "/api/jobs",
        files={"deck": ("deck.md", b"# One\n\n- a\n", "text/markdown")},
        data={"output": "pptx", "notes": "all"},
    )
    assert r.status_code == 200
    assert seen == {"mode": "pptx", "notes": "all"}
    token = r.json()["token"]
    job_id = r.json()["job_id"]
    result = client.get(f"/api/jobs/{job_id}/result", params={"t": token})
    assert result.status_code == 200
    assert result.content.startswith(b"PK")  # a zip, not an mp4


def test_pptx_rejects_bad_notes_mode(base_config):
    client = TestClient(serve.create_app(config=base_config, token=None))
    r = client.post(
        "/api/jobs",
        files={"deck": ("deck.md", b"# One\n\n- a\n", "text/markdown")},
        data={"output": "pptx", "notes": "bogus"},
    )
    assert r.status_code == 400


def test_job_config_ephemeral_voice_and_photo(base_config, tmp_path):
    """The per-job config wires the uploaded voice/photo into providers."""
    workdir = tmp_path / "job"
    workdir.mkdir()
    voice = workdir / "voice.wav"
    voice.write_bytes(b"a")
    photo = workdir / "photo.png"
    photo.write_bytes(b"b")

    job_yaml = serve._build_job_config(
        base_config, workdir,
        {"avatar": True, "narration_seconds": "25", "image_provider": "swarmui"},
        voice, photo,
    )
    import yaml

    cfg = yaml.safe_load(job_yaml.read_text())
    assert cfg["providers"]["tts"]["voice_sample"] == str(voice)
    assert cfg["providers"]["avatar"]["source_image"] == str(photo)
    assert cfg["providers"]["images"]["provider"] == "swarmui"
    assert cfg["settings"]["narration"]["target_seconds"] == 25.0


def test_job_config_photo_not_animated_becomes_static_self(base_config, tmp_path):
    """Photo + animate off = a still of themselves, not 'no head'."""
    workdir = tmp_path / "job"
    workdir.mkdir()
    photo = workdir / "photo.png"
    photo.write_bytes(b"b")
    job_yaml = serve._build_job_config(
        base_config, workdir, {"avatar": False}, None, photo
    )
    import yaml

    cfg = yaml.safe_load(job_yaml.read_text())
    assert cfg["providers"]["avatar"]["provider"] == "static"
    assert cfg["providers"]["avatar"]["source"] == str(photo)


def test_job_config_no_presenter_means_no_head(base_config, tmp_path):
    workdir = tmp_path / "job"
    workdir.mkdir()
    job_yaml = serve._build_job_config(base_config, workdir, {"avatar": True}, None, None)
    import yaml

    cfg = yaml.safe_load(job_yaml.read_text())
    assert cfg["providers"]["avatar"]["provider"] == "none"


def test_job_config_mascot_animate_toggle(base_config, tmp_path):
    """Mascot: without a usable animation engine, animate falls back to a
    static mascot — the web UI never produces the puppet mouth-flap."""
    import yaml

    workdir = tmp_path / "job"
    workdir.mkdir()
    on = yaml.safe_load(serve._build_job_config(
        base_config, workdir, {"avatar": True, "avatar_name": "teddy"}, None, None
    ).read_text())
    assert on["providers"]["avatar"]["provider"] == "static"
    assert on["providers"]["avatar"]["source"] == "teddy"
    off = yaml.safe_load(serve._build_job_config(
        base_config, workdir, {"avatar": False, "avatar_name": "teddy"}, None, None
    ).read_text())
    assert off["providers"]["avatar"]["provider"] == "static"


def test_job_config_mascot_animates_via_wan_s2v_when_configured(base_config, tmp_path):
    """With a detector-free engine configured, an animated mascot lip-syncs
    for real (wan-s2v) instead of falling back to the puppet mouth-flap."""
    import yaml

    base_config["providers"]["avatar"] = {
        "provider": "wan-s2v", "base_url": "https://comfy.example.org",
        "api_key": "tok",
    }
    workdir = tmp_path / "job"
    workdir.mkdir()
    cfg = yaml.safe_load(serve._build_job_config(
        base_config, workdir, {"avatar": True, "avatar_name": "owl"}, None, None
    ).read_text())
    av = cfg["providers"]["avatar"]
    assert av["provider"] == "wan-s2v"
    assert av["source"] == "owl"
    assert av["base_url"] == "https://comfy.example.org"  # server connection kept
    # Animate off still holds the mascot as a static image, engine or not.
    off = yaml.safe_load(serve._build_job_config(
        base_config, workdir, {"avatar": False, "avatar_name": "owl"}, None, None
    ).read_text())
    assert off["providers"]["avatar"]["provider"] == "static"


def test_job_config_mascot_falls_back_to_static_for_human_only_engine(base_config, tmp_path):
    """A human-only engine (sadtalker) can't animate a mascot, so an
    animated mascot degrades to a still image rather than the puppet
    mouth-flap (or a doomed SadTalker render)."""
    import yaml

    base_config["providers"]["avatar"] = {
        "provider": "sadtalker", "base_url": "https://comfy.example.org",
    }
    workdir = tmp_path / "job"
    workdir.mkdir()
    cfg = yaml.safe_load(serve._build_job_config(
        base_config, workdir, {"avatar": True, "avatar_name": "teddy"}, None, None
    ).read_text())
    assert cfg["providers"]["avatar"]["provider"] == "static"


def test_job_config_video_always_animates(base_config, tmp_path):
    """A video presenter uses the video engine even with animate off."""
    import yaml

    workdir = tmp_path / "job"
    workdir.mkdir()
    clip = workdir / "me.mp4"
    clip.write_bytes(b"v")
    cfg = yaml.safe_load(serve._build_job_config(
        base_config, workdir, {"avatar": False}, None, clip
    ).read_text())
    assert cfg["providers"]["avatar"].get("provider") != "static"
    assert cfg["providers"]["avatar"]["source_video"] == str(clip)


def test_job_config_video_routes_away_from_wan_s2v(base_config, tmp_path):
    """With the server engine wan-s2v (a still-image model), an uploaded
    video presenter routes through the auto provider (video -> Wav2Lip)."""
    import yaml

    base_config["providers"]["avatar"] = {
        "provider": "wan-s2v", "base_url": "https://comfy.example.org",
    }
    workdir = tmp_path / "job"
    workdir.mkdir()
    clip = workdir / "me.mp4"
    clip.write_bytes(b"v")
    cfg = yaml.safe_load(serve._build_job_config(
        base_config, workdir, {"avatar": True}, None, clip
    ).read_text())
    av = cfg["providers"]["avatar"]
    assert av["provider"] == "comfyui"
    assert av["source_video"] == str(clip)
    assert av["base_url"] == "https://comfy.example.org"  # same server kept


def test_download_uses_per_job_token(base_config, tmp_path):
    """A browser download link (no header) authenticates via a per-job ?t=."""
    client = TestClient(serve.create_app(config=base_config, token="secret"))
    # Seed a finished job directly.
    out = tmp_path / "o.mp4"
    out.write_bytes(b"VID")
    job = serve.Job(id="j1", status="done", workdir=tmp_path, output_path=out,
                    download_token="dltok")
    serve._JOBS["j1"] = job

    # No auth at all -> 401.
    assert client.get("/api/jobs/j1/result").status_code == 401
    # The long-lived instance token must never ride in a URL: rejected.
    assert client.get("/api/jobs/j1/result", params={"t": "secret"}).status_code == 401
    # The per-job download token -> ok.
    r = client.get("/api/jobs/j1/result", params={"t": "dltok"})
    assert r.status_code == 200
    assert r.content == b"VID"
    # API clients can still use the Authorization header.
    hdr = client.get("/api/jobs/j1/result",
                     headers={"Authorization": "Bearer secret"})
    assert hdr.status_code == 200


# --- demo mode: friction-free but guard-railed --------------------------------


@pytest.fixture(autouse=True)
def _clear_demo_hits():
    serve._DEMO_HITS.clear()
    serve._DEMO_DRAFT_HITS.clear()
    yield
    serve._DEMO_HITS.clear()
    serve._DEMO_DRAFT_HITS.clear()


def test_demo_mode_needs_no_token(base_config):
    """Demo mode is open even when a token is configured."""
    client = TestClient(serve.create_app(config=base_config, token="secret", demo=True))
    cfg = client.get("/api/config").json()
    assert cfg["auth_required"] is False
    assert cfg["demo"] is True
    assert cfg["limits"] == {
        "max_slides": serve.DEMO_MAX_SLIDES,
        "jobs_per_hour": serve.DEMO_JOBS_PER_HOUR,
    }
    # protected endpoint reachable without a token
    assert client.get("/api/jobs/whatever").status_code == 404


# --- teaser: text formats, voice picker, avatar art ----------------------------


def test_build_topic_prompt_contains_topic_and_count():
    from slide_stream.draft import build_topic_prompt

    prompt = build_topic_prompt("Why sleep matters", 5)
    assert "Why sleep matters" in prompt
    assert "exactly 5 slides" in prompt
    assert "'# ' heading" in prompt  # output contract for the parser


def test_topic_draft_endpoint(base_config, monkeypatch):
    captured: dict[str, object] = {}

    def fake_draft(topic, slides, provider, model, base_url,
                   api_key=None, think=None):
        captured.update(topic=topic, slides=slides)
        return "# A\n\n- x\n\n# B\n\n- y\n"

    monkeypatch.setattr(serve, "_do_topic_draft", fake_draft)
    client = TestClient(serve.create_app(config=_llm_config(), token=None))
    r = client.post(
        "/api/draft-deck",
        json={"topic": "Why sleep matters for learning", "slides": 2},
    )
    assert r.status_code == 200, r.text
    assert "# A" in r.json()["markdown"]
    assert captured["topic"] == "Why sleep matters for learning"
    assert captured["slides"] == 2


def test_topic_draft_demo_caps_slides_and_defaults_to_five(monkeypatch):
    captured: dict[str, object] = {}

    def fake_draft(topic, slides, provider, model, base_url,
                   api_key=None, think=None):
        captured["slides"] = slides
        return "# A\n\n- x\n"

    monkeypatch.setattr(serve, "_do_topic_draft", fake_draft)
    client = TestClient(serve.create_app(config=_llm_config(), demo=True))
    r = client.post("/api/draft-deck", json={"topic": "dogs", "slides": 9})
    assert r.status_code == 200
    assert r.json()["slides"] == serve.DEMO_MAX_SLIDES
    assert captured["slides"] == serve.DEMO_MAX_SLIDES


def test_topic_draft_requires_topic(base_config):
    client = TestClient(serve.create_app(config=_llm_config(), token=None))
    assert client.post("/api/draft-deck", json={"topic": "  "}).status_code == 400


def test_topic_draft_requires_llm():
    client = TestClient(serve.create_app(config=copy.deepcopy(DEFAULT_CONFIG), token=None))
    r = client.post("/api/draft-deck", json={"topic": "dogs at the park"})
    assert r.status_code == 400
    assert "LLM" in r.json()["detail"]


def test_topic_draft_demo_rate_limited(monkeypatch):
    monkeypatch.setattr(
        serve, "_do_topic_draft", lambda *a, **k: "# A\n\n- x\n"
    )
    client = TestClient(serve.create_app(config=_llm_config(), demo=True))
    for _ in range(serve.DEMO_DRAFTS_PER_HOUR):
        assert client.post("/api/draft-deck", json={"topic": "cats"}).status_code == 200
    over = client.post("/api/draft-deck", json={"topic": "cats"})
    assert over.status_code == 429
    assert "drafts per hour" in over.json()["detail"]


def test_demo_accepts_txt_with_slides(base_config, mocker, monkeypatch):
    """A plain .txt structured with '# ' headings is parsed as a deck."""
    seen: dict[str, str] = {}

    def fake_run_job(job, deck_path, job_yaml, voice_path, photo_path,
                     mode="video", notes=None):
        assert job.workdir is not None
        out = job.workdir / "output.mp4"
        out.write_bytes(b"V")
        job.status = "done"
        job.output_path = out
        seen["deck_name"] = deck_path.name

    mocker.patch.object(serve, "_run_job", side_effect=fake_run_job)
    monkeypatch.setattr(
        serve.ThreadPoolExecutor, "submit", lambda self, fn, *a, **k: fn(*a, **k)
    )
    client = TestClient(serve.create_app(config=base_config, demo=True))
    r = client.post("/api/jobs", files={
        "deck": ("notes.txt", b"# First slide\n- a\n\n# Second\n- b\n", "text/plain"),
    })
    assert r.status_code == 200, r.text
    assert seen["deck_name"] == "deck.md"  # stored under the canonical name


def test_demo_rejects_text_without_slide_structure(base_config, mocker):
    client = TestClient(serve.create_app(config=base_config, demo=True))
    r = client.post("/api/jobs", files={
        "deck": ("notes.txt", b"just some prose,\nno headings at all\n", "text/plain"),
    })
    assert r.status_code == 400
    assert "No slides found" in r.json()["detail"]
    assert "'# '" in r.json()["detail"]


def test_demo_accepts_qmd(base_config, mocker, monkeypatch):
    """A Quarto (.qmd) file with front matter parses; front matter is dropped."""
    seen: dict[str, int] = {}

    def fake_run_job(job, deck_path, job_yaml, voice_path, photo_path,
                     mode="video", notes=None):
        assert job.workdir is not None
        out = job.workdir / "output.mp4"
        out.write_bytes(b"V")
        job.status = "done"
        job.output_path = out
        seen["slides"] = len(serve._parse_deck_slides(deck_path))

    mocker.patch.object(serve, "_run_job", side_effect=fake_run_job)
    monkeypatch.setattr(
        serve.ThreadPoolExecutor, "submit", lambda self, fn, *a, **k: fn(*a, **k)
    )
    client = TestClient(serve.create_app(config=base_config, demo=True))
    qmd = b"---\ntitle: Deck\n---\n\n# One\n\n- a\n\n# Two\n\n- b\n"
    r = client.post("/api/jobs", files={"deck": ("deck.qmd", qmd, "text/markdown")})
    assert r.status_code == 200, r.text
    assert seen["slides"] == 2


def test_voice_name_maps_to_configured_stock_voice(base_config, tmp_path):
    """A teaser voice choice overrides stored voice keys and resolves the
    display name back to the configured file ('Emily' -> 'Emily.wav')."""
    base = copy.deepcopy(base_config)
    base["providers"]["tts"].update({
        "provider": "voicebox",
        "profile_id": "stored-id",
        "voice": "old-stock.wav",
        "reference_text": "stale",
        "voice_choices": ["Emily.wav", "James.wav"],
    })
    workdir = tmp_path / "job"
    workdir.mkdir()
    cfg = yaml.safe_load(serve._build_job_config(
        base, workdir, {"avatar": True, "voice_name": "Emily"}, None, None
    ).read_text())
    tts = cfg["providers"]["tts"]
    assert tts["voice"] == "Emily.wav"
    assert "profile_id" not in tts and "reference_text" not in tts


def test_api_config_exposes_voices(base_config):
    base_config["providers"]["tts"]["voice_choices"] = ["Emily.wav", "James.wav"]
    client = TestClient(serve.create_app(config=base_config, demo=True))
    assert client.get("/api/config").json()["voices"] == ["Emily", "James"]


def test_avatar_image_endpoint(base_config):
    from slide_stream.avatars import avatar_names

    client = TestClient(serve.create_app(config=base_config, demo=True))
    name = avatar_names()[0]
    ok = client.get(f"/api/avatars/{name}")
    assert ok.status_code == 200
    assert ok.headers["content-type"].startswith("image/")
    assert client.get("/api/avatars/not-a-mascot").status_code == 404


def test_demo_defaults_presenter_to_first_and_last(base_config, mocker):
    """Teaser economics: without explicit placement, demo presenters appear
    on the opening and closing slides only (wan-s2v is minutes per clip)."""
    real = serve._build_job_config
    captured: dict[str, object] = {}

    def spy(base, workdir, options, *args, **kwargs):
        captured["slides"] = options.get("avatar_slides")
        return real(base, workdir, options, *args, **kwargs)

    mocker.patch.object(serve, "_build_job_config", side_effect=spy)
    mocker.patch.object(serve, "_run_job")
    client = TestClient(serve.create_app(config=base_config, demo=True))
    r = client.post("/api/jobs", files={"deck": ("deck.md", b"# One\n- a\n")},
                    data={"avatar_name": "teddy"})
    assert r.status_code == 200, r.text
    assert captured["slides"] == "first,last"

    # Explicit placement always wins.
    captured.clear()
    r = client.post("/api/jobs", files={"deck": ("deck.md", b"# One\n- a\n")},
                    data={"avatar_name": "teddy", "avatar_slides": "first"})
    assert r.status_code == 200
    assert captured["slides"] == "first"


def test_demo_drops_voice_and_photo_uploads(base_config, mocker, monkeypatch):
    """The teaser accepts no biometrics: voice/photo uploads are never
    written to disk or passed to the render, even from API clients."""
    seen: dict[str, object] = {}

    def fake_run_job(job, deck_path, job_yaml, voice_path, photo_path,
                     mode="video", notes=None):
        assert job.workdir is not None
        out = job.workdir / "output.mp4"
        out.write_bytes(b"V")
        job.status = "done"
        job.output_path = out
        seen["voice"] = voice_path
        seen["photo"] = photo_path

    mocker.patch.object(serve, "_run_job", side_effect=fake_run_job)
    monkeypatch.setattr(
        serve.ThreadPoolExecutor, "submit", lambda self, fn, *a, **k: fn(*a, **k)
    )
    client = TestClient(serve.create_app(config=base_config, demo=True))
    r = client.post("/api/jobs", files={
        "deck": ("deck.md", b"# One\n- a\n"),
        "voice": ("me.wav", b"AUDIO", "audio/wav"),
        "photo": ("me.png", _tiny_png(), "image/png"),
    })
    assert r.status_code == 200
    assert seen["voice"] is None
    assert seen["photo"] is None
    # Nothing biometric ever landed in the job workdir.
    workdir = serve._JOBS[r.json()["job_id"]].workdir
    assert workdir is not None
    assert not list(workdir.glob("voice*"))
    assert not list(workdir.glob("photo*"))


def test_demo_slide_cap_trims_to_five(base_config, mocker, monkeypatch):
    """A big deck is no longer rejected: the demo trims it to the first 5."""
    seen: dict[str, int] = {}

    def fake_run_job(job, deck_path, job_yaml, voice_path, photo_path,
                     mode="video", notes=None):
        assert job.workdir is not None
        out = job.workdir / "output.mp4"
        out.write_bytes(b"V")
        job.status = "done"
        job.output_path = out
        seen["slides"] = len(serve._parse_deck_slides(deck_path))

    mocker.patch.object(serve, "_run_job", side_effect=fake_run_job)
    monkeypatch.setattr(
        serve.ThreadPoolExecutor, "submit", lambda self, fn, *a, **k: fn(*a, **k)
    )
    client = TestClient(serve.create_app(config=base_config, demo=True))
    big_deck = "".join(f"# Slide {i}\n- point\n" for i in range(10))
    r = client.post("/api/jobs", files={"deck": ("deck.md", big_deck.encode())})
    assert r.status_code == 200
    body = r.json()
    assert body["notice"]
    assert "first 5" in body["notice"]
    assert "desktop app" in body["notice"]
    assert seen["slides"] == serve.DEMO_MAX_SLIDES

    # A small deck is accepted untouched, with no notice.
    ok = client.post("/api/jobs", files={"deck": ("deck.md", b"# One\n- a\n")})
    assert ok.status_code == 200
    assert ok.json()["notice"] is None


def test_demo_pptx_cap_trims_to_five(base_config, mocker, monkeypatch):
    """An oversized .pptx is trimmed to DEMO_MAX_SLIDES slides in place."""
    from pptx import Presentation

    def fake_run_job(job, deck_path, job_yaml, voice_path, photo_path,
                     mode="video", notes=None):
        assert job.workdir is not None
        out = job.workdir / "output.mp4"
        out.write_bytes(b"V")
        job.status = "done"
        job.output_path = out

    mocker.patch.object(serve, "_run_job", side_effect=fake_run_job)
    monkeypatch.setattr(
        serve.ThreadPoolExecutor, "submit", lambda self, fn, *a, **k: fn(*a, **k)
    )
    prs = Presentation()
    for i in range(8):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Slide {i}"  # type: ignore[union-attr]
    buf = BytesIO()
    prs.save(buf)

    client = TestClient(serve.create_app(config=base_config, demo=True))
    r = client.post(
        "/api/jobs",
        files={"deck": ("deck.pptx", buf.getvalue(),
                        "application/vnd.openxmlformats-officedocument."
                        "presentationml.presentation")},
    )
    assert r.status_code == 200, r.text
    assert "first 5" in r.json()["notice"]
    # The deck handed to the render has exactly DEMO_MAX_SLIDES slides.
    jobs = [j for j in serve._JOBS.values() if j.workdir]
    assert jobs
    workdir = jobs[-1].workdir
    assert workdir is not None
    pptx_files = list(workdir.glob("*.pptx"))
    assert pptx_files
    assert serve._count_slides(pptx_files[0]) == serve.DEMO_MAX_SLIDES


def test_truncate_markdown_deck_heading_style():
    from slide_stream.parser import parse_markdown

    md = "".join(f"# S{i}\n\n- point {i}\n\n" for i in range(8))
    out = serve._truncate_markdown_deck(md, 3)
    assert len(parse_markdown(out)) == 3
    assert "# S0" in out and "# S2" in out and "# S3" not in out


def test_truncate_markdown_deck_separator_style():
    from slide_stream.parser import parse_markdown

    blocks = [f"# S{i}\npoint {i}" for i in range(6)]
    md = "\n---\n".join(blocks) + "\n"
    out = serve._truncate_markdown_deck(md, 2)
    slides = parse_markdown(out)
    assert len(slides) == 2
    assert slides[0]["title"] == "S0" and slides[1]["title"] == "S1"


def test_truncate_markdown_deck_keeps_front_matter():
    from slide_stream.parser import parse_markdown

    md = "---\ntitle: Deck\ntheme: dark\n---\n# A\n\n# B\n\n# C\n"
    out = serve._truncate_markdown_deck(md, 2)
    assert out.startswith("---\ntitle: Deck")
    assert len(parse_markdown(out)) == 2


def test_truncate_markdown_deck_matches_parser_on_fences():
    """Boundary rules intentionally mirror parse_markdown — a '#' line inside
    an unclosed-looking fence counts (and cuts) the same way in both, so a
    trimmed deck never parses to more slides than were kept."""
    from slide_stream.parser import parse_markdown

    md = "# A\n```text\n# fake\n```\n\n# B\n\n# C\n"
    out = serve._truncate_markdown_deck(md, 3)
    assert "# B" in out and "# C" not in out
    assert len(parse_markdown(out)) <= 3


def test_truncate_markdown_deck_short_deck_unchanged():
    from slide_stream.parser import parse_markdown

    md = "# A\n\n# B\n"
    out = serve._truncate_markdown_deck(md, 5)
    assert len(parse_markdown(out)) == 2


def test_demo_rate_limit(base_config, mocker):
    client = TestClient(serve.create_app(config=base_config, demo=True))
    mocker.patch.object(serve, "_run_job")
    deck = {"deck": ("deck.md", b"# One\n- a\n")}
    for _ in range(serve.DEMO_JOBS_PER_HOUR):
        assert client.post("/api/jobs", files=deck).status_code == 200
    over = client.post("/api/jobs", files=deck)
    assert over.status_code == 429
    assert "per hour" in over.json()["detail"]


def test_demo_rate_window_expires():
    now = 1_000_000.0
    ip = "1.2.3.4"
    for _ in range(serve.DEMO_JOBS_PER_HOUR):
        assert serve._demo_rate_ok(ip, now=now)
    assert not serve._demo_rate_ok(ip, now=now + 10)
    # an hour later the window has rolled over
    assert serve._demo_rate_ok(ip, now=now + 3601)


def test_non_demo_still_requires_token(base_config):
    client = TestClient(serve.create_app(config=base_config, token="secret", demo=False))
    assert client.get("/api/jobs/whatever").status_code == 401


# --- local (desktop) mode ------------------------------------------------------


def test_local_mode_open_and_no_demo_limits(base_config):
    """Desktop mode: no token, demo limits off even if both are set."""
    client = TestClient(serve.create_app(
        config=base_config, token="secret", demo=True, local=True))
    cfg = client.get("/api/config").json()
    assert cfg["local"] is True
    assert cfg["auth_required"] is False
    assert cfg["demo"] is False
    assert client.get("/api/jobs/whatever").status_code == 404  # open


def test_settings_hidden_unless_local(base_config):
    client = TestClient(serve.create_app(config=base_config, local=False))
    assert client.get("/api/settings").status_code == 404


def test_settings_roundtrip(base_config, tmp_path, monkeypatch):
    from pathlib import Path as _P

    monkeypatch.setattr(_P, "home", staticmethod(lambda: tmp_path))
    # Local mode only accepts state-changing requests addressed to localhost.
    client = TestClient(serve.create_app(config=base_config, local=True),
                        base_url="http://127.0.0.1")

    s = client.get("/api/settings").json()
    assert s["yaml"] == ""  # nothing saved yet
    assert "providers:" in s["template"]

    ok = client.put("/api/settings",
                    json={"yaml": "providers:\n  tts:\n    provider: gtts\n"})
    assert ok.status_code == 200
    saved = tmp_path / ".slidestream.yaml"
    assert saved.exists()
    # Owner-only: the file typically holds API keys.
    assert (saved.stat().st_mode & 0o777) == 0o600
    assert "gtts" in client.get("/api/settings").json()["yaml"]


def test_settings_rejects_bad_yaml(base_config, tmp_path, monkeypatch):
    from pathlib import Path as _P

    monkeypatch.setattr(_P, "home", staticmethod(lambda: tmp_path))
    client = TestClient(serve.create_app(config=base_config, local=True),
                        base_url="http://127.0.0.1")
    bad = client.put("/api/settings", json={"yaml": "a: [unclosed"})
    assert bad.status_code == 400
    assert not (tmp_path / ".slidestream.yaml").exists()
    notdict = client.put("/api/settings", json={"yaml": "- just\n- a list\n"})
    assert notdict.status_code == 400


# --- upload limits ---------------------------------------------------------


def test_deck_size_limit_413(base_config, monkeypatch, mocker):
    monkeypatch.setattr(serve, "MAX_DECK_BYTES", 1024)
    mocker.patch.object(serve, "_run_job")
    client = TestClient(serve.create_app(config=base_config, token=None))
    r = client.post("/api/jobs",
                    files={"deck": ("deck.md", b"# S\n" + b"x" * 4096)})
    assert r.status_code == 413
    assert "Deck" in r.json()["detail"]


def test_voice_size_limit_413(base_config, monkeypatch, mocker):
    monkeypatch.setattr(serve, "MAX_VOICE_BYTES", 1024)
    mocker.patch.object(serve, "_run_job")
    client = TestClient(serve.create_app(config=base_config, token=None))
    r = client.post("/api/jobs", files={
        "deck": ("deck.md", b"# One\n- a\n"),
        "voice": ("me.wav", b"x" * 4096, "audio/wav"),
    })
    assert r.status_code == 413
    assert "Voice" in r.json()["detail"]


def test_photo_upload_must_be_a_real_image(base_config, mocker):
    mocker.patch.object(serve, "_run_job")
    client = TestClient(serve.create_app(config=base_config, token=None))
    r = client.post("/api/jobs", files={
        "deck": ("deck.md", b"# One\n- a\n"),
        "photo": ("me.png", b"NOTANIMAGE", "image/png"),
    })
    assert r.status_code == 400


def test_photo_dimension_cap(base_config, monkeypatch, mocker):
    monkeypatch.setattr(serve, "MAX_IMAGE_DIM", 2)  # the test PNG is 4x4
    mocker.patch.object(serve, "_run_job")
    client = TestClient(serve.create_app(config=base_config, token=None))
    r = client.post("/api/jobs", files={
        "deck": ("deck.md", b"# One\n- a\n"),
        "photo": ("me.png", _tiny_png(), "image/png"),
    })
    assert r.status_code == 400
    assert "maximum" in r.json()["detail"]


# --- X-Forwarded-For handling ------------------------------------------------


def test_xff_ignored_without_trusted_proxy(base_config, mocker, monkeypatch):
    monkeypatch.delenv("SLIDESTREAM_TRUSTED_PROXY", raising=False)
    mocker.patch.object(serve, "_run_job")
    client = TestClient(serve.create_app(config=base_config, demo=True))
    deck = {"deck": ("deck.md", b"# One\n- a\n")}
    for i in range(serve.DEMO_JOBS_PER_HOUR):
        r = client.post("/api/jobs", files=deck,
                        headers={"X-Forwarded-For": f"10.0.0.{i}"})
        assert r.status_code == 200
    # A freshly spoofed XFF must not buy another job.
    over = client.post("/api/jobs", files=deck,
                       headers={"X-Forwarded-For": "10.99.99.99"})
    assert over.status_code == 429


def test_xff_rightmost_honored_behind_trusted_proxy(base_config, mocker,
                                                    monkeypatch):
    monkeypatch.setenv("SLIDESTREAM_TRUSTED_PROXY", "1")
    mocker.patch.object(serve, "_run_job")
    client = TestClient(serve.create_app(config=base_config, demo=True))
    deck = {"deck": ("deck.md", b"# One\n- a\n")}
    # The proxy-appended (rightmost) value keys the limit; the client-chosen
    # left entries are ignored.
    for i in range(serve.DEMO_JOBS_PER_HOUR):
        r = client.post("/api/jobs", files=deck,
                        headers={"X-Forwarded-For": f"10.0.0.{i}, 9.9.9.9"})
        assert r.status_code == 200
    over = client.post("/api/jobs", files=deck,
                       headers={"X-Forwarded-For": "10.0.0.99, 9.9.9.9"})
    assert over.status_code == 429
    # A genuinely different client (per the proxy) gets its own bucket.
    ok = client.post("/api/jobs", files=deck,
                     headers={"X-Forwarded-For": "10.0.0.1, 8.8.8.8"})
    assert ok.status_code == 200


# --- demo hardening -----------------------------------------------------------


def test_demo_unparseable_deck_fails_closed(base_config, mocker):
    mocker.patch.object(serve, "_run_job")
    mocker.patch.object(serve, "_count_slides", return_value=None)
    client = TestClient(serve.create_app(config=base_config, demo=True))
    r = client.post("/api/jobs", files={"deck": ("deck.md", b"# One\n- a\n")})
    assert r.status_code == 400
    assert "parse" in r.json()["detail"].lower()


def test_demo_status_hides_log_and_token(base_config):
    client = TestClient(serve.create_app(config=base_config, demo=True))
    serve._JOBS["j1"] = serve.Job(id="j1", status="running",
                                  log="Traceback: /secret/server/path",
                                  download_token="dltok")
    s = client.get("/api/jobs/j1").json()
    assert s["status"] == "running"  # coarse progress the UI shows
    assert s["log"] == ""
    assert "token" not in s  # else anyone with the UUID could download


def test_status_returns_log_and_token_when_authed(base_config):
    client = TestClient(serve.create_app(config=base_config, token="secret"))
    serve._JOBS["j1"] = serve.Job(id="j1", status="running",
                                  log="rendering slide 2",
                                  download_token="dltok")
    s = client.get("/api/jobs/j1",
                   headers={"Authorization": "Bearer secret"}).json()
    assert s["log"] == "rendering slide 2"
    assert s["token"] == "dltok"


def test_demo_download_needs_job_token_and_wipes_workdir(base_config, tmp_path):
    client = TestClient(serve.create_app(config=base_config, demo=True))
    workdir = tmp_path / "job"
    workdir.mkdir()
    out = workdir / "output.mp4"
    out.write_bytes(b"VID")
    serve._JOBS["j1"] = serve.Job(id="j1", status="done", workdir=workdir,
                                  output_path=out, download_token="dltok")
    # The UUID alone is no longer enough in demo mode.
    assert client.get("/api/jobs/j1/result").status_code == 401
    r = client.get("/api/jobs/j1/result", params={"t": "dltok"})
    assert r.status_code == 200
    assert r.content == b"VID"
    # "Nothing stored": one download, then the whole job is gone.
    assert not workdir.exists()
    assert "j1" not in serve._JOBS


# --- local (desktop) mode CSRF/DNS-rebinding guards ---------------------------


def test_local_mode_rejects_foreign_host(base_config):
    app = serve.create_app(config=base_config, local=True)
    rebound = TestClient(app, base_url="http://evil.example.com")
    assert rebound.put("/api/settings", json={"yaml": ""}).status_code == 403
    assert rebound.post("/api/jobs", files={
        "deck": ("deck.md", b"# One\n- a\n")}).status_code == 403


def test_local_mode_rejects_foreign_origin(base_config, tmp_path, monkeypatch):
    from pathlib import Path as _P

    monkeypatch.setattr(_P, "home", staticmethod(lambda: tmp_path))
    client = TestClient(serve.create_app(config=base_config, local=True),
                        base_url="http://localhost")
    evil = client.put("/api/settings", json={"yaml": "a: 1\n"},
                      headers={"Origin": "https://evil.example.com"})
    assert evil.status_code == 403
    # Local pages and the Tauri shell's webview stay allowed.
    ok = client.put("/api/settings", json={"yaml": "a: 1\n"},
                    headers={"Origin": "http://127.0.0.1:8080"})
    assert ok.status_code == 200
    tauri = client.put("/api/settings", json={"yaml": "a: 1\n"},
                       headers={"Origin": "tauri://localhost"})
    assert tauri.status_code == 200


def test_local_origin_helper():
    assert serve._local_origin_ok("http://localhost:8080")
    assert serve._local_origin_ok("http://127.0.0.1")
    assert serve._local_origin_ok("tauri://localhost")
    assert serve._local_origin_ok("http://tauri.localhost")
    assert not serve._local_origin_ok("https://evil.example.com")
    assert not serve._local_origin_ok("null")
    assert not serve._local_origin_ok("http://localhost.evil.com")


# --- job hygiene: secrets and artifacts do not persist -------------------------


def test_job_yaml_written_owner_only(base_config, tmp_path):
    workdir = tmp_path / "job"
    workdir.mkdir()
    job_yaml = serve._build_job_config(base_config, workdir,
                                       {"avatar": True}, None, None)
    assert (job_yaml.stat().st_mode & 0o777) == 0o600


def test_voice_upload_overrides_stored_profile(base_config, tmp_path):
    """An uploaded clip must beat any inherited stored-voice keys."""
    base = copy.deepcopy(base_config)
    base["providers"]["tts"].update({
        "provider": "voicebox", "profile_id": "stored-id",
        "voice": "narrator", "reference_text": "server clip transcript",
    })
    workdir = tmp_path / "job"
    workdir.mkdir()
    voice = workdir / "voice.wav"
    voice.write_bytes(b"a")
    cfg = yaml.safe_load(serve._build_job_config(
        base, workdir, {"avatar": True}, voice, None).read_text())
    tts = cfg["providers"]["tts"]
    assert tts["voice_sample"] == str(voice)
    assert "profile_id" not in tts
    assert "voice" not in tts
    assert "reference_text" not in tts


def test_render_artifacts_wiped_after_job(base_config, tmp_path, monkeypatch):
    """After a render only output.mp4 remains: no deck, job.yaml, or tmp/."""
    workdir = tmp_path / "job"
    workdir.mkdir()
    deck = workdir / "deck.md"
    deck.write_text("# One\n- a\n")
    voice = workdir / "voice.wav"
    voice.write_bytes(b"a")
    job_yaml = serve._build_job_config(base_config, workdir,
                                       {"avatar": True}, voice, None)
    (workdir / "tmp").mkdir()

    class FakePopen:
        def __init__(self, cmd, **kwargs):
            self.cmd = cmd
            self.returncode = 0
            self.stdout = io.StringIO("rendered ok\n")
            Path(cmd[cmd.index("create") + 2]).write_bytes(b"MP4")

        def wait(self, timeout=None):
            return 0

        def kill(self):
            pass

    monkeypatch.setattr(serve.subprocess, "Popen", FakePopen)
    job = serve.Job(id="j1", workdir=workdir, created_at=time.time())
    serve._run_job(job, deck, job_yaml, voice, None)

    assert job.status == "done"
    assert (workdir / "output.mp4").exists()
    assert not job_yaml.exists()  # held live API keys
    assert not deck.exists()
    assert not voice.exists()
    assert not (workdir / "tmp").exists()


def test_reaper_evicts_expired_jobs(tmp_path):
    workdir = tmp_path / "job"
    workdir.mkdir()
    now = time.time()
    serve._JOBS["old"] = serve.Job(
        id="old", status="done", workdir=workdir,
        created_at=now - serve.JOB_TTL_SECONDS - 10)
    serve._JOBS["new"] = serve.Job(id="new", status="done", created_at=now)
    serve._JOBS["run"] = serve.Job(
        id="run", status="running",
        created_at=now - serve.JOB_TTL_SECONDS - 10)
    serve._reap_expired_jobs(now=now)
    assert "old" not in serve._JOBS
    assert not workdir.exists()
    assert "new" in serve._JOBS
    assert "run" in serve._JOBS  # running jobs get the render-timeout grace


def test_settings_template_mentions_voicebox():
    assert "voicebox" in serve.SETTINGS_TEMPLATE
    for key in ("base_url", "profile_id", "voice_sample", "engine"):
        assert key in serve.SETTINGS_TEMPLATE


# --- live progress: streamed logs, heartbeat, demo-safe progress ---------------


def test_stream_process_updates_log_and_heartbeat(base_config, tmp_path, monkeypatch):
    """Output arrives line-by-line while the process runs, refreshing
    job.log and job.updated_at — the render looks alive from the outside."""
    workdir = tmp_path / "job"
    workdir.mkdir()
    deck = workdir / "deck.md"
    deck.write_text("# One\n- a\n")
    output = workdir / "output.mp4"
    output.write_bytes(b"MP4")

    class FakePopen:
        def __init__(self, cmd, **kwargs):
            self.returncode = 0
            self.stdout = io.StringIO(
                "Slide 1/2: One\n  - Generated audio with gTTS\nSlide 2/2: Two\n"
            )

        def wait(self, timeout=None):
            return 0

        def kill(self):
            pass

    monkeypatch.setattr(serve.subprocess, "Popen", FakePopen)
    job = serve.Job(id="j1", workdir=workdir, created_at=time.time())
    code = serve._stream_process(["create"], job, timeout=60)

    assert code == 0
    assert "Slide 2/2" in job.log
    assert job.updated_at > 0
    progress = serve._job_progress(job)
    assert progress is not None
    assert progress["slide"] == 2
    assert progress["slides"] == 2
    assert progress["stage"] == "generating voice"


def test_job_progress_stitching_stage():
    job = serve.Job(id="j", created_at=time.time())
    job.log = "Slide 1/3: A\nSlide 2/3: B\nSlide 3/3: C\n2. Combining Video Fragments...\n"
    progress = serve._job_progress(job)
    assert progress == {
        "slide": 3, "slides": 3, "stage": "stitching the final video",
    }


def test_job_progress_heartbeat_line():
    """A long GPU render prints heartbeats; the stage reads 'rendering the
    presenter' even mid-slide."""
    job = serve.Job(id="j", created_at=time.time())
    job.log = "Slide 4/10: Deep dive\n  - wan-s2v talking head: rendering… 190s elapsed\n"
    progress = serve._job_progress(job)
    assert progress is not None
    assert progress["slide"] == 4
    assert progress["stage"] == "rendering the presenter"


def test_job_progress_none_for_empty_or_unrecognised():
    job = serve.Job(id="j", created_at=time.time())
    assert serve._job_progress(job) is None
    job.log = "hello world\nnothing useful here\n"
    assert serve._job_progress(job) is None


def test_status_exposes_progress_and_idle_but_not_raw_log_in_demo(base_config):
    """Demo clients get derived progress + heartbeat, never raw log text."""
    serve._JOBS["j1"] = serve.Job(
        id="j1", status="running", log="Slide 1/5: A\n",
        created_at=time.time() - 30, updated_at=time.time() - 4,
        download_token="dltok",
    )
    client = TestClient(serve.create_app(config=base_config, demo=True))
    s = client.get("/api/jobs/j1").json()
    assert s["progress"] == {"slide": 1, "slides": 5}
    assert s["idle"] == 4
    assert s["log"] == ""  # raw text still withheld on the open demo
    assert "token" not in s


def test_status_idle_uses_created_at_before_output(base_config):
    """No output yet -> the heartbeat counts from job creation."""
    serve._JOBS["j1"] = serve.Job(
        id="j1", status="running", log="",
        created_at=time.time() - 12,
    )
    client = TestClient(serve.create_app(config=base_config, token=None))
    s = client.get("/api/jobs/j1").json()
    assert s["idle"] == 12
    assert "progress" not in s  # nothing recognisable printed yet


# --- demo-safe step trace -------------------------------------------------------


def test_job_trace_parses_providers_and_fallbacks():
    job = serve.Job(id="j", created_at=time.time())
    job.log = (
        "Found 2 slides.\n"
        "Slide 1/2: Coffee\n"
        "  - Generated SwarmUI image: coffee beans\n"
        "  - Generated audio with voicebox\n"
        "  - Generated wan-s2v talking head (slide 1)\n"
        "Slide 2/2: Tea\n"
        "  - SwarmUI error: timeout. Using text fallback.\n"
        "  - Generated audio with voicebox\n"
        "  - wan-s2v avatar error: GPU OOM\n"
        "2. Combining Video Fragments...\n"
    )
    trace = serve._job_trace(job)
    assert trace == [
        "slide 1: image via Swarmui",
        "slide 1: voice via voicebox",
        "slide 1: presenter animated (wan-s2v)",
        "slide 2: image failed — used a text card",
        "slide 2: voice via voicebox",
        "slide 2: presenter failed (wan-s2v)",
        "stitching slides together",
    ]


def test_job_warnings_include_text_fallback_count():
    job = serve.Job(id="j", created_at=time.time())
    job.log = (
        "Slide 1/3: A\n  - SwarmUI error: x. Using text fallback.\n"
        "Slide 2/3: B\n  - SwarmUI error: y. Using text fallback.\n"
        "Slide 3/3: C\n  - Generated Imagen image: z\n"
    )
    warnings = serve._job_warnings(job)
    assert any("2 of 3 slides used a plain text card" in w for w in warnings)


def test_status_exposes_trace_in_demo(base_config):
    serve._JOBS["j1"] = serve.Job(
        id="j1", status="running",
        log="Slide 1/2: A\n  - Generated SwarmUI image: x\n",
        created_at=time.time(), updated_at=time.time(),
    )
    client = TestClient(serve.create_app(config=base_config, demo=True))
    s = client.get("/api/jobs/j1").json()
    assert s["trace"] == ["slide 1: image via Swarmui"]
    assert s["log"] == ""  # raw text still withheld


def test_status_exposes_warnings_in_demo(base_config):
    serve._JOBS["j1"] = serve.Job(
        id="j1", status="done",
        log="Slide 1/1: A\n  - SwarmUI error: down. Using text fallback.\n",
        created_at=time.time(),
    )
    client = TestClient(serve.create_app(config=base_config, demo=True))
    s = client.get("/api/jobs/j1").json()
    assert any("text card" in w for w in s["warnings"])


# --- Project workflow (draft -> edit -> enrich -> render) ---------------------


@pytest.fixture(autouse=True)
def _clear_projects():
    serve._PROJECTS.clear()
    yield
    serve._PROJECTS.clear()


def _llm_config():
    cfg = copy.deepcopy(DEFAULT_CONFIG)
    cfg["providers"]["llm"]["provider"] = "openai"
    return cfg


def test_create_empty_project_and_state():
    client = TestClient(serve.create_app(config=_llm_config(), token=None))
    r = client.post("/api/projects")
    assert r.status_code == 200, r.text
    body = r.json()
    pid, token = body["project_id"], body["token"]
    assert token
    assert body["state"]["has_deck"] is False

    # State needs the per-project token, not just the id.
    assert client.get(f"/api/projects/{pid}").status_code == 401
    ok = client.get(f"/api/projects/{pid}", headers={"X-Project-Token": token})
    assert ok.status_code == 200
    assert ok.json()["has_deck"] is False


def test_create_project_with_deck_upload():
    client = TestClient(serve.create_app(config=_llm_config(), token=None))
    r = client.post(
        "/api/projects",
        files={"deck": ("deck.md", b"# One\n\n- a\n# Two\n\n- b\n", "text/markdown")},
    )
    assert r.status_code == 200, r.text
    state = r.json()["state"]
    assert state["has_deck"] is True
    assert state["slide_count"] == 2
    assert state["deck_format"] == "md"


def test_project_draft_writes_deck(monkeypatch):
    monkeypatch.setattr(
        serve, "_do_draft",
        lambda *a, **k: "# Intro\n\n- point\n# Details\n\n- more\n",
    )
    client = TestClient(serve.create_app(config=_llm_config(), token=None))
    pid = client.post("/api/projects").json()["project_id"]
    token = serve._PROJECTS[pid].token

    r = client.post(
        f"/api/projects/{pid}/draft",
        files={"source": ("report.pdf", b"%PDF-fake", "application/pdf")},
        data={"slides": "2"},
        headers={"X-Project-Token": token},
    )
    assert r.status_code == 200, r.text
    assert "# Intro" in r.json()["markdown"]
    assert r.json()["state"]["slide_count"] == 2


def test_project_draft_requires_llm_provider():
    # DEFAULT_CONFIG has llm provider 'none' -> draft must refuse cleanly.
    client = TestClient(serve.create_app(config=copy.deepcopy(DEFAULT_CONFIG), token=None))
    pid = client.post("/api/projects").json()["project_id"]
    token = serve._PROJECTS[pid].token
    r = client.post(
        f"/api/projects/{pid}/draft",
        files={"source": ("a.txt", b"hi", "text/plain")},
        headers={"X-Project-Token": token},
    )
    assert r.status_code == 400
    assert "LLM" in r.json()["detail"]


def test_project_save_deck_validates():
    client = TestClient(serve.create_app(config=_llm_config(), token=None))
    pid = client.post("/api/projects").json()["project_id"]
    token = serve._PROJECTS[pid].token

    # Prose with no headings is not a deck.
    bad = client.put(
        f"/api/projects/{pid}/deck",
        json={"markdown": "just prose"},
        headers={"X-Project-Token": token},
    )
    assert bad.status_code == 400

    good = client.put(
        f"/api/projects/{pid}/deck",
        json={"markdown": "# A\n\n- x\n# B\n\n- y\n"},
        headers={"X-Project-Token": token},
    )
    assert good.status_code == 200
    assert good.json()["state"]["slide_count"] == 2


def test_project_render_spawns_job(monkeypatch):
    def fake_run_job(job, deck_path, job_yaml, voice_path, photo_path,
                     mode="video", notes=None):
        assert job.workdir is not None
        assert deck_path.exists()  # the project deck was copied in
        out = job.workdir / "output.mp4"
        out.write_bytes(b"VID")
        job.status = "done"
        job.output_path = out

    monkeypatch.setattr(serve, "_run_job", fake_run_job)
    monkeypatch.setattr(
        serve.ThreadPoolExecutor, "submit", lambda self, fn, *a, **k: fn(*a, **k)
    )
    client = TestClient(serve.create_app(config=_llm_config(), token=None))
    pid = client.post(
        "/api/projects",
        files={"deck": ("deck.md", b"# One\n\n- a\n", "text/markdown")},
    ).json()["project_id"]
    token = serve._PROJECTS[pid].token

    r = client.post(
        f"/api/projects/{pid}/render",
        data={"avatar": "false"},
        headers={"X-Project-Token": token},
    )
    assert r.status_code == 200, r.text
    job_id, jtok = r.json()["job_id"], r.json()["token"]
    result = client.get(f"/api/jobs/{job_id}/result", params={"t": jtok})
    assert result.status_code == 200
    assert result.content == b"VID"
    # The canonical project deck survives the render (was copied, not consumed).
    assert (serve._PROJECTS[pid].workdir / "deck.md").exists()


def test_project_enrich_spawns_enrich_job(monkeypatch):
    captured = {}

    def fake_enrich(job, project, deck_path, job_yaml, notes):
        captured["notes"] = notes
        captured["project_id"] = project.id
        assert job.workdir is not None
        out = job.workdir / "enriched.zip"
        out.write_bytes(b"ZIP")
        job.status = "done"
        job.output_path = out

    monkeypatch.setattr(serve, "_run_project_enrich", fake_enrich)
    monkeypatch.setattr(
        serve.ThreadPoolExecutor, "submit", lambda self, fn, *a, **k: fn(*a, **k)
    )
    client = TestClient(serve.create_app(config=_llm_config(), token=None))
    pid = client.post(
        "/api/projects",
        files={"deck": ("deck.md", b"# One\n\n- a\n", "text/markdown")},
    ).json()["project_id"]
    token = serve._PROJECTS[pid].token

    r = client.post(
        f"/api/projects/{pid}/enrich",
        data={"image_provider": "text", "notes": "all"},
        headers={"X-Project-Token": token},
    )
    assert r.status_code == 200, r.text
    assert captured["notes"] == "all"
    assert captured["project_id"] == pid


def test_project_workflow_blocked_in_demo():
    client = TestClient(serve.create_app(config=_llm_config(), demo=True))
    assert client.post("/api/projects").status_code == 403


def test_reap_expired_projects(tmp_path):
    p = serve.Project(id="old", workdir=tmp_path, created_at=1.0, token="t")
    serve._PROJECTS["old"] = p
    marker = tmp_path / "deck.md"
    marker.write_text("# A\n")
    serve._reap_expired_projects(now=1.0 + serve.PROJECT_TTL_SECONDS + 1)
    assert "old" not in serve._PROJECTS


def test_run_project_enrich_populates_project_images(tmp_path, monkeypatch):
    proj_dir = tmp_path / "proj"
    proj_dir.mkdir()
    job_dir = tmp_path / "job"
    job_dir.mkdir()
    proj = serve.Project(id="p", workdir=proj_dir, created_at=time.time(), token="t")
    job = serve.Job(id="j", workdir=job_dir, created_at=time.time())
    deck = job_dir / "deck.md"
    deck.write_text("# A\n")
    jy = job_dir / "job.yaml"
    jy.write_text("{}")

    class FakePopen:
        def __init__(self, cmd, **kwargs):
            self.returncode = 0
            self.stdout = io.StringIO("enriched\n")
            images = job_dir / "enriched" / "images"
            images.mkdir(parents=True)
            (images / "slide_1.png").write_bytes(b"PNG")
            (job_dir / "enriched.zip").write_bytes(b"ZIP")

        def wait(self, timeout=None):
            return 0

        def kill(self):
            pass

    monkeypatch.setattr(serve.subprocess, "Popen", FakePopen)
    serve._run_project_enrich(job, proj, deck, jy, None)
    assert job.status == "done"
    assert (proj_dir / "images" / "slide_1.png").read_bytes() == b"PNG"


def test_project_image_thumbnail_endpoint():
    client = TestClient(serve.create_app(config=_llm_config(), token=None))
    pid = client.post("/api/projects").json()["project_id"]
    project = serve._PROJECTS[pid]
    (project.workdir / "images").mkdir()
    (project.workdir / "images" / "slide_1.png").write_bytes(b"\x89PNGdata")

    # Wrong/absent token is rejected; the ?t= token works (for <img> tags).
    assert client.get(f"/api/projects/{pid}/images/slide_1.png").status_code == 401
    r = client.get(f"/api/projects/{pid}/images/slide_1.png", params={"t": project.token})
    assert r.status_code == 200
    assert r.content == b"\x89PNGdata"
    # Path traversal is neutralised.
    assert client.get(f"/api/projects/{pid}/images/nope.png",
                      params={"t": project.token}).status_code == 404


def test_project_render_includes_enriched_images(monkeypatch):
    seen = {}

    def fake_run_job(job, deck_path, job_yaml, voice_path, photo_path,
                     mode="video", notes=None):
        assert job.workdir is not None
        seen["has_images"] = (job.workdir / "images" / "slide_1.png").exists()
        out = job.workdir / "output.mp4"
        out.write_bytes(b"VID")
        job.status = "done"
        job.output_path = out

    monkeypatch.setattr(serve, "_run_job", fake_run_job)
    monkeypatch.setattr(
        serve.ThreadPoolExecutor, "submit", lambda self, fn, *a, **k: fn(*a, **k)
    )
    client = TestClient(serve.create_app(config=_llm_config(), token=None))
    pid = client.post(
        "/api/projects",
        files={"deck": ("deck.md", b"# One\n\n![](images/slide_1.png)\n\n- a\n", "text/markdown")},
    ).json()["project_id"]
    project = serve._PROJECTS[pid]
    (project.workdir / "images").mkdir()
    (project.workdir / "images" / "slide_1.png").write_bytes(b"PNG")

    r = client.post(f"/api/projects/{pid}/render", data={"avatar": "false"},
                    headers={"X-Project-Token": project.token})
    assert r.status_code == 200, r.text
    assert seen["has_images"] is True  # enriched images travelled with the deck


def test_build_job_config_applies_avatar_placement_and_transition(tmp_path):
    base = copy.deepcopy(DEFAULT_CONFIG)
    options = {
        "avatar_slides": "first,last",
        "reuse_avatar": True,
        "transition_seconds": "0.4",
    }
    p = serve._build_job_config(base, tmp_path, options, None, None)
    cfg = yaml.safe_load(p.read_text(encoding="utf-8"))
    assert cfg["settings"]["avatar"]["slides"] == "first,last"
    assert cfg["settings"]["avatar"]["reuse_clip"] is True
    assert cfg["settings"]["video"]["transition_seconds"] == 0.4


def test_project_render_passes_placement_and_transition(monkeypatch):
    from pathlib import Path

    captured = {}

    def fake_run_job(job, deck_path, job_yaml, voice, photo, mode="video", notes=None):
        captured["cfg"] = yaml.safe_load(Path(job_yaml).read_text(encoding="utf-8"))
        assert job.workdir is not None
        out = job.workdir / "output.mp4"
        out.write_bytes(b"V")
        job.status = "done"
        job.output_path = out

    monkeypatch.setattr(serve, "_run_job", fake_run_job)
    monkeypatch.setattr(
        serve.ThreadPoolExecutor, "submit", lambda self, fn, *a, **k: fn(*a, **k)
    )
    client = TestClient(serve.create_app(config=_llm_config(), token=None))
    pid = client.post(
        "/api/projects",
        files={"deck": ("deck.md", b"# A\n\n- a\n", "text/markdown")},
    ).json()["project_id"]
    token = serve._PROJECTS[pid].token

    r = client.post(
        f"/api/projects/{pid}/render",
        data={"avatar": "false", "avatar_slides": "first",
              "reuse_avatar": "true", "transition": "0.5"},
        headers={"X-Project-Token": token},
    )
    assert r.status_code == 200, r.text
    cfg = captured["cfg"]
    assert cfg["settings"]["avatar"]["slides"] == "first"
    assert cfg["settings"]["avatar"]["reuse_clip"] is True
    assert cfg["settings"]["video"]["transition_seconds"] == 0.5
