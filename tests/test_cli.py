"""Tests for the CLI interface."""

# python-pptx ships incomplete type stubs; the PowerPoint fixture below builds
# a .pptx file via the runtime API, so suppress the stub-gap diagnostics here.
# pyright: reportArgumentType=false, reportAttributeAccessIssue=false, reportOptionalMemberAccess=false

import tempfile
import wave
from pathlib import Path

import pytest
from pptx import Presentation
from typer.testing import CliRunner

from slide_stream.cli import app

# Small/fast video settings so the real-encode tests stay quick.
FAST_VIDEO_CONFIG = (
    "settings:\n"
    "  video:\n"
    "    resolution: [320, 180]\n"
    "    fps: 8\n"
    "    default_slide_duration: 1.0\n"
    "    slide_duration_padding: 0.2\n"
)


def write_silent_wav(filename, seconds=0.3):
    """A real, decodable audio file (WAV content behind any extension)."""
    with wave.open(str(filename), "wb") as w:
        w.setnchannels(1)
        w.setsampwidth(2)
        w.setframerate(22050)
        w.writeframes(b"\x00\x00" * int(22050 * seconds))


@pytest.fixture
def runner():
    """Create a CLI test runner."""
    return CliRunner()


@pytest.fixture
def sample_markdown():
    """Create sample markdown content."""
    return """# Test Slide

- First point
- Second point

# Another Slide

- More content
- Final point
"""


@pytest.fixture
def sample_powerpoint():
    """Create sample PowerPoint file."""
    temp_file = Path(tempfile.mkdtemp()) / "test.pptx"

    # Create a presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout

    title = slide.shapes.title
    title.text = "Test PowerPoint Slide"

    content = slide.placeholders[1]
    content.text = "First bullet point\nSecond bullet point"

    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "These are test speaker notes."

    prs.save(temp_file)
    return temp_file


def test_cli_help(runner):
    """Test CLI help command."""
    result = runner.invoke(app, ["--help"])
    assert result.exit_code == 0
    assert "Create a video from a Markdown (.md) or PowerPoint (.pptx) file" in result.stdout


def test_cli_version(runner):
    """Test CLI version command."""
    from slide_stream import __version__

    result = runner.invoke(app, ["--version"])
    assert result.exit_code == 0
    assert "SlideStream" in result.stdout
    assert __version__ in result.stdout


def test_cli_no_args_shows_help(runner):
    """Bare invocation prints help and exits 0 instead of erroring."""
    result = runner.invoke(app, [])
    assert result.exit_code == 0
    assert "Usage" in result.output


def test_cli_init_refuses_overwrite_without_force(runner, tmp_path):
    """init must not clobber an existing file unless --force is given."""
    target = tmp_path / "slidestream.yaml"
    target.write_text("existing: content\n")

    result = runner.invoke(app, ["init", str(target)])
    assert result.exit_code != 0
    assert target.read_text() == "existing: content\n"

    # --force overwrites the existing file.
    result = runner.invoke(app, ["init", str(target), "--force"])
    assert result.exit_code == 0
    assert "SlideStream Configuration" in target.read_text()


def test_cli_create_missing_input(runner):
    """Test CLI with missing input file."""
    result = runner.invoke(app, ["create"])
    assert result.exit_code != 0


def test_cli_create_basic(runner, sample_markdown, tmp_path, mocker, monkeypatch):
    """create renders a real video from Markdown, with TTS mocked out."""
    monkeypatch.chdir(tmp_path)
    md_file = tmp_path / "slides.md"
    md_file.write_text(sample_markdown)
    config_file = tmp_path / "config.yaml"
    config_file.write_text(FAST_VIDEO_CONFIG)

    fake_tts = mocker.MagicMock()
    fake_tts.save.side_effect = write_silent_wav
    mocker.patch("gtts.gTTS", return_value=fake_tts)

    output = tmp_path / "out.mp4"
    result = runner.invoke(
        app,
        ["create", str(md_file), str(output), "--config", str(config_file)],
    )

    assert result.exit_code == 0, result.output
    assert output.exists()


def test_cli_create_powerpoint(runner, sample_powerpoint, tmp_path, mocker, monkeypatch):
    """create renders a real video from PowerPoint, with TTS mocked out."""
    monkeypatch.chdir(tmp_path)
    config_file = tmp_path / "config.yaml"
    config_file.write_text(FAST_VIDEO_CONFIG)

    fake_tts = mocker.MagicMock()
    fake_tts.save.side_effect = write_silent_wav
    mocker.patch("gtts.gTTS", return_value=fake_tts)

    output = tmp_path / "out.mp4"
    try:
        result = runner.invoke(
            app,
            [
                "create",
                str(sample_powerpoint),
                str(output),
                "--config",
                str(config_file),
            ],
        )

        assert result.exit_code == 0, result.output
        assert output.exists()
    finally:
        sample_powerpoint.unlink(missing_ok=True)


def test_cli_unsupported_file_type(runner):
    """Test CLI with unsupported file type."""
    with tempfile.NamedTemporaryFile(mode='w', suffix='.txt', delete=False) as f:
        f.write("Some text content")
        f.flush()

        result = runner.invoke(app, [
            "create", f.name, "test_output.mp4"
        ])

        assert result.exit_code != 0
        assert "Unsupported file type" in result.stdout

        # Clean up
        Path(f.name).unlink(missing_ok=True)


def test_cli_nonexistent_file(runner):
    """Test CLI with non-existent input file."""
    result = runner.invoke(app, [
        "create", "nonexistent.md", "test_output.mp4"
    ])

    assert result.exit_code != 0
    assert "Input file not found" in result.stdout


def test_cli_strict_aborts_on_unavailable_tts_provider(
    runner, sample_markdown, tmp_path, monkeypatch
):
    """--strict fails fast when the configured TTS provider is unusable,
    instead of silently rendering the whole video with gTTS."""
    monkeypatch.delenv("ELEVENLABS_API_KEY", raising=False)
    monkeypatch.chdir(tmp_path)

    md_file = tmp_path / "slides.md"
    md_file.write_text(sample_markdown)
    config_file = tmp_path / "config.yaml"
    config_file.write_text(
        "providers:\n  tts:\n    provider: elevenlabs\n"
    )

    result = runner.invoke(
        app,
        [
            "create",
            str(md_file),
            str(tmp_path / "out.mp4"),
            "--config",
            str(config_file),
            "--strict",
        ],
    )
    assert result.exit_code == 1
    assert not (tmp_path / "out.mp4").exists()


def test_cli_strict_aborts_on_llm_narration_failure(
    runner, sample_markdown, tmp_path, mocker, monkeypatch
):
    """--strict fails fast when LLM narration returns nothing, instead of
    silently falling back to raw-bullet narration."""
    # An unconfigured voicebox (the default TTS) quietly maps to gTTS; make
    # sure a developer's env vars can't turn it into a strict-mode error
    # before the LLM branch under test is even reached.
    monkeypatch.delenv("VOICEBOX_BASE_URL", raising=False)
    monkeypatch.delenv("VOICEBOX_TOKEN", raising=False)
    monkeypatch.chdir(tmp_path)
    md_file = tmp_path / "slides.md"
    md_file.write_text(sample_markdown)

    mocker.patch(
        "slide_stream.cli.get_llm_client", return_value=mocker.MagicMock()
    )
    # query_llm swallows all backend errors into None.
    mocker.patch("slide_stream.cli.query_llm", return_value=None)

    output = tmp_path / "out.mp4"
    result = runner.invoke(
        app,
        [
            "create",
            str(md_file),
            str(output),
            "--llm-provider",
            "claude",
            "--strict",
        ],
    )

    assert result.exit_code == 1
    assert "LLM narration failed" in result.output
    assert not output.exists()


def test_cli_scan_skips_images_whose_description_fails(runner, tmp_path, mocker):
    """A failed vision call must not rename the file to a slug of its own
    name; the image is skipped and reported instead."""
    (tmp_path / "a.png").write_bytes(b"img")
    (tmp_path / "b.png").write_bytes(b"img")

    mocker.patch(
        "slide_stream.scan.get_llm_client", return_value=mocker.MagicMock()
    )
    mocker.patch(
        "slide_stream.scan.query_llm_with_image",
        side_effect=["golden retriever in park", None],
    )

    result = runner.invoke(app, ["scan", str(tmp_path), "--apply"])

    assert result.exit_code == 0, result.output
    assert (tmp_path / "golden-retriever-in-park.png").exists()
    assert (tmp_path / "b.png").exists()  # untouched, not renamed to "b.png" slug
    assert "b.png" in result.output


def test_cli_scan_exits_nonzero_when_all_descriptions_fail(
    runner, tmp_path, mocker
):
    (tmp_path / "a.png").write_bytes(b"img")

    mocker.patch(
        "slide_stream.scan.get_llm_client", return_value=mocker.MagicMock()
    )
    mocker.patch("slide_stream.scan.query_llm_with_image", return_value=None)

    result = runner.invoke(app, ["scan", str(tmp_path), "--apply"])

    assert result.exit_code == 1
    assert (tmp_path / "a.png").exists()  # nothing was renamed


def test_cli_scan_dry_run_matches_apply(runner, tmp_path, mocker):
    """Dry-run resolves collisions against the planned renames (names taken
    and vacated), so its preview matches what --apply actually does."""
    (tmp_path / "a.png").write_bytes(b"img")
    (tmp_path / "b.png").write_bytes(b"img")

    mocker.patch(
        "slide_stream.scan.get_llm_client", return_value=mocker.MagicMock()
    )
    # a.png -> zebra.png (vacating "a.png"), b.png -> a.png (the vacated
    # name). Two invocations (dry-run then apply) consume the side effects.
    mocker.patch(
        "slide_stream.scan.query_llm_with_image",
        side_effect=["zebra", "a", "zebra", "a"],
    )

    dry = runner.invoke(app, ["scan", str(tmp_path)])
    assert dry.exit_code == 0, dry.output
    # Without vacated-name tracking the preview would show "a-1.png".
    assert "a-1.png" not in dry.output

    applied = runner.invoke(app, ["scan", str(tmp_path), "--apply"])
    assert applied.exit_code == 0, applied.output
    assert (tmp_path / "zebra.png").exists()
    assert (tmp_path / "a.png").exists()  # b.png took the vacated name
    assert not (tmp_path / "b.png").exists()
