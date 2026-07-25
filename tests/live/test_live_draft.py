"""Live draft / rewrite: turn a document into a slide deck via the real LLM.

    uv run pytest tests/live/test_live_draft.py --run-live -q
"""

from typer.testing import CliRunner

from slide_stream.cli import app
from slide_stream.parser import parse_markdown

_SOURCE = (
    "The water cycle moves water continuously through the environment. "
    "Evaporation lifts water from oceans and lakes into the air as vapour. "
    "The vapour cools and condenses into clouds, then falls as precipitation "
    "— rain, snow, or hail. Runoff and groundwater carry it back to rivers and "
    "oceans, and the cycle repeats. The sun drives the whole system, and plants "
    "add moisture through transpiration."
)


def test_live_draft_from_text(live_config, tmp_path):
    """A plain document becomes a valid multi-slide deck."""
    src = tmp_path / "source.txt"
    src.write_text(_SOURCE)
    out = tmp_path / "deck.md"

    result = CliRunner().invoke(
        app, ["draft", str(src), str(out), "--slides", "3", "--llm-provider", "claude"]
    )
    assert result.exit_code == 0, result.output
    slides = parse_markdown(out.read_text(encoding="utf-8"))
    assert len(slides) >= 2, f"draft produced too few slides: {len(slides)}"
    assert all(s["title"] for s in slides), "a slide has no title"
    print(f"\nDrafted {len(slides)} slides from a text document")


def test_live_rewrite_from_pptx(live_config, tmp_path):
    """An existing PowerPoint is restructured into a clean Markdown deck."""
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Photosynthesis"  # type: ignore[union-attr]
    slide.placeholders[1].text = "Light reactions\nCalvin cycle\nChlorophyll"  # type: ignore[attr-defined]
    src = tmp_path / "in.pptx"
    prs.save(src)
    out = tmp_path / "deck.md"

    result = CliRunner().invoke(
        app, ["draft", str(src), str(out), "--llm-provider", "claude"]
    )
    assert result.exit_code == 0, result.output
    slides = parse_markdown(out.read_text(encoding="utf-8"))
    assert slides, "rewrite produced no slides"
    print(f"\nRewrote a .pptx into a {len(slides)}-slide markdown deck")
