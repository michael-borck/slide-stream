"""Live enrich: slides -> an editable deck (Markdown + images + PowerPoint),
using the real configured image provider. This is the non-video output track.

Skips when the image provider resolves to 'text' (no live image gen).
    uv run pytest tests/live/test_live_enrich.py --run-live -q
"""

import numpy as np
import pytest
from PIL import Image

from slide_stream.enrich import enrich_deck
from slide_stream.llm import get_llm_client
from slide_stream.providers.factory import ProviderFactory


def test_live_enrich_writes_deck_images_and_pptx(live_config, tmp_path):
    provider = ProviderFactory.create_image_provider(live_config)
    if provider.name == "text":
        pytest.skip("image provider resolved to 'text' — no live image gen configured")

    slides = [
        {"title": "A red apple on a wooden table", "content": ["fresh fruit"]},
        {"title": "A blue ceramic coffee mug", "content": ["steam", "morning"]},
    ]
    out = tmp_path / "enriched"
    enrich_deck(slides, provider, out, "deck", also_pptx=True)

    # Editable Markdown deck + one image per slide.
    assert (out / "deck.md").exists(), "no enriched markdown written"
    images = sorted((out / "images").glob("slide_*.png"))
    assert len(images) == 2, f"expected 2 images, got {len(images)}"

    # Images are real generations, not flat text-card fallbacks.
    for p in images:
        arr = np.asarray(Image.open(p).convert("RGB"))
        assert float(arr.std()) > 12, f"{p.name} looks like a flat text card"

    # PowerPoint opens and has one slide per entry.
    from pptx import Presentation

    prs = Presentation(str(out / "deck.pptx"))
    assert len(prs.slides) == 2, "pptx slide count mismatch"
    print(f"\nEnrich exercised: {provider.name} -> deck.md + {len(images)} images + deck.pptx")


def test_live_enrich_writes_ai_notes(live_config, tmp_path):
    """`enrich --notes all` writes real AI presenter notes into the .pptx."""
    llm = live_config["providers"].get("llm", {})
    if llm.get("provider", "none") == "none":
        pytest.skip("no LLM configured for presenter-note generation")

    # The image provider is incidental here (notes are the subject); any works.
    provider = ProviderFactory.create_image_provider(live_config)
    narration = live_config["settings"].get("narration", {})
    ctx = {
        "client": get_llm_client(llm["provider"], base_url=llm.get("base_url")),
        "provider": llm["provider"],
        "model": llm.get("model"),
        "target_seconds": narration.get("target_seconds"),
        "wpm": narration.get("wpm", 150),
    }
    slides = [{"title": "Photosynthesis", "content": ["plants turn light into energy"]}]
    out = tmp_path / "out"
    enrich_deck(slides, provider, out, "deck", also_pptx=True, notes_mode="all", llm=ctx)

    from pptx import Presentation

    slide = Presentation(str(out / "deck.pptx")).slides[0]
    note = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""  # type: ignore[union-attr]
    assert len(note.split()) >= 5, f"AI notes too short: {note!r}"
    print(f"\nAI notes generated via {llm['provider']} ({len(note.split())} words)")
