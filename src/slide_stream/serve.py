# FastAPI route handlers are registered via decorators, so pyright's
# unused-function check is a false positive for them.
# pyright: reportUnusedFunction=false
"""Web UI for slide-stream (optional ``[serve]`` extra).

`slide-stream serve` starts a small FastAPI app: upload a deck (+ optional
voice sample and photo), render it as a background job, and download the video.
Token-authenticated so it can run locally or on a VPS.

Design notes:
- The server is **stateless about biometric data**: an uploaded voice sample /
  photo is used only for that render and deleted afterwards. The lecturer's
  browser remembers them (IndexedDB) so they need not re-pick each job — the
  data stays on their laptop, never stored on the server at rest.
- Each render runs as a subprocess (``python -m slide_stream create``) so a
  crash can't take down the server and ffmpeg/moviepy memory is reclaimed.
"""

import copy
import os
import re
import secrets
import shutil
import subprocess
import sys
import tempfile
import threading
import time
import uuid
from collections.abc import Awaitable, Callable
from concurrent.futures import ThreadPoolExecutor
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any
from urllib.parse import urlsplit

import yaml

from . import __version__
from .config_loader import load_config
from .parser import _H1_RE, _YAML_KEY_RE, _use_separator_style


@dataclass
class Job:
    id: str
    status: str = "queued"  # queued | running | done | error
    log: str = ""
    error: str = ""
    workdir: Path | None = None
    output_path: Path | None = None
    created_at: float = field(default_factory=lambda: 0.0)
    # When log output last arrived (heartbeat for long renders: a GPU job can
    # legitimately go quiet for minutes, and the UI says "still working").
    updated_at: float = field(default_factory=lambda: 0.0)
    # Per-job download secret: knowing the job UUID alone must not be enough
    # to fetch the video, and the long-lived instance token never goes in a URL.
    download_token: str = ""
    # Output media (video render vs enriched-deck zip) for the download response.
    media_type: str = "video/mp4"
    download_name: str = "slidestream.mp4"


# In-memory job registry (single-process v1). job_id -> Job.
_JOBS: dict[str, Job] = {}
_LOCK = threading.Lock()


@dataclass
class Project:
    """A workflow session: a server-side workdir holding the deck (and, once
    enriched, its images) so the draft -> enrich -> render steps chain without
    re-uploading between each. Unlike a Job, a Project persists across requests;
    it is reaped on a TTL like jobs. Voice/photo inputs are NEVER stored on a
    project — those stay ephemeral, attached per render job."""
    id: str
    workdir: Path
    created_at: float = field(default_factory=lambda: 0.0)
    # Per-project secret: the project id alone must not authorize edits/reads.
    token: str = ""


# In-memory project registry (single-process v1). project_id -> Project.
_PROJECTS: dict[str, Project] = {}


def _shutdown() -> None:  # pragma: no cover - kills the process
    os._exit(0)


# Demo-mode guardrails: friction-free (no token) but bounded.
DEMO_MAX_SLIDES = 5
DEMO_JOBS_PER_HOUR = 3
_DEMO_HITS: dict[str, list[float]] = {}  # client ip -> job timestamps


def _env_int(name: str, default: int) -> int:
    """Integer from the environment, falling back on missing/garbage values."""
    try:
        return int(os.getenv(name, "") or default)
    except ValueError:
        return default


# Upload caps (bytes), overridable via SLIDESTREAM_MAX_{DECK,VOICE,PHOTO}_MB.
MAX_DECK_BYTES = _env_int("SLIDESTREAM_MAX_DECK_MB", 30) * 1024 * 1024
MAX_VOICE_BYTES = _env_int("SLIDESTREAM_MAX_VOICE_MB", 30) * 1024 * 1024
MAX_PHOTO_BYTES = _env_int("SLIDESTREAM_MAX_PHOTO_MB", 15) * 1024 * 1024
# Largest photo we will hand to Pillow/the avatar engines (per side).
MAX_IMAGE_DIM = 8000
# How long a job (and its workdir) may live before it is reaped.
JOB_TTL_SECONDS = _env_int("SLIDESTREAM_JOB_TTL_MIN", 60) * 60


def _reap_expired_jobs(now: float | None = None) -> None:
    """Evict jobs past the TTL and delete their workdirs.

    Called lazily from the job endpoints so 'nothing stored' holds without a
    dedicated reaper thread. Jobs that are still queued/running get an extra
    hour of grace (the render subprocess timeout) so a live render's files
    are never deleted out from under it.
    """
    t = now if now is not None else time.time()
    expired: list[Job] = []
    with _LOCK:
        for job in list(_JOBS.values()):
            grace = 0 if job.status in ("done", "error") else 3600
            if job.created_at and t - job.created_at > JOB_TTL_SECONDS + grace:
                expired.append(job)
                del _JOBS[job.id]
    for job in expired:
        if job.workdir is not None:
            shutil.rmtree(job.workdir, ignore_errors=True)


# Projects live at least as long as jobs; a workflow spans several requests.
PROJECT_TTL_SECONDS = _env_int("SLIDESTREAM_PROJECT_TTL_MIN", 240) * 60


def _reap_expired_projects(now: float | None = None) -> None:
    """Evict projects past the TTL and delete their workdirs (lazy, like jobs)."""
    t = now if now is not None else time.time()
    expired: list[Project] = []
    with _LOCK:
        for project in list(_PROJECTS.values()):
            if project.created_at and t - project.created_at > PROJECT_TTL_SECONDS:
                expired.append(project)
                del _PROJECTS[project.id]
    for project in expired:
        shutil.rmtree(project.workdir, ignore_errors=True)


def _project_deck(project: Project) -> Path | None:
    """The project's canonical deck file (.md preferred), or None if unset."""
    for name in ("deck.md", "deck.pptx"):
        candidate = project.workdir / name
        if candidate.exists():
            return candidate
    return None


def _project_state(project: Project) -> dict[str, Any]:
    """A JSON-able snapshot of what the project currently holds."""
    deck = _project_deck(project)
    images_dir = project.workdir / "images"
    images = (
        sorted(p.name for p in images_dir.iterdir() if p.is_file())
        if images_dir.is_dir()
        else []
    )
    slide_count = None
    if deck is not None:
        try:
            slide_count = len(_parse_deck_slides(deck))
        except Exception:
            slide_count = None
    return {
        "project_id": project.id,
        "has_deck": deck is not None,
        "deck_format": deck.suffix.lstrip(".") if deck else None,
        "slide_count": slide_count,
        "images": images,
    }


# Origins the local (desktop) server accepts state-changing requests from:
# this machine's own pages, or the Tauri shell's webview.
_LOCAL_HOSTS = ("localhost", "127.0.0.1", "::1")
_TAURI_ORIGINS = ("tauri://localhost", "http://tauri.localhost",
                  "https://tauri.localhost")


def _local_origin_ok(origin: str) -> bool:
    """True if an Origin header value belongs to this machine (or Tauri)."""
    if origin.lower() in _TAURI_ORIGINS:
        return True
    try:
        parts = urlsplit(origin)
    except ValueError:
        return False
    return parts.scheme in ("http", "https") and (
        (parts.hostname or "") in _LOCAL_HOSTS
    )


def _validate_photo_upload(path: Path) -> str | None:
    """Error message if an uploaded image is undecodable or hostile, else None.

    Pillow's MAX_IMAGE_PIXELS bomb guard stays active; on top of it we bound
    the dimensions so downstream avatar engines get something reasonable.
    """
    from PIL import Image, UnidentifiedImageError

    try:
        with Image.open(path) as im:
            width, height = im.size
            if width > MAX_IMAGE_DIM or height > MAX_IMAGE_DIM:
                return (
                    f"Image is {width}x{height}; the maximum is "
                    f"{MAX_IMAGE_DIM}x{MAX_IMAGE_DIM}"
                )
            im.verify()
    except Image.DecompressionBombError:
        return "Image is too large to decode safely"
    except (UnidentifiedImageError, OSError, ValueError):
        return "Could not decode the image"
    return None


def _demo_rate_ok(ip: str, now: float | None = None) -> bool:
    """True if this IP may start another demo job (and record the hit)."""
    t = now if now is not None else time.time()
    with _LOCK:
        hits = [h for h in _DEMO_HITS.get(ip, []) if t - h < 3600]
        if len(hits) >= DEMO_JOBS_PER_HOUR:
            _DEMO_HITS[ip] = hits
            return False
        hits.append(t)
        _DEMO_HITS[ip] = hits
        return True


def _count_slides(deck_path: Path) -> int | None:
    """Best-effort slide count for the demo cap; None if unparseable."""
    try:
        if deck_path.suffix.lower() == ".pptx":
            from pptx import Presentation  # type: ignore[import-untyped]

            return len(Presentation(str(deck_path)).slides)
        from .parser import parse_markdown

        return len(parse_markdown(deck_path.read_text(encoding="utf-8")))
    except Exception:
        return None


def _parse_deck_slides(deck_path: Path) -> list[dict[str, Any]]:
    """Parse a deck to slide dicts for the doctor preflight (.md or .pptx)."""
    if deck_path.suffix.lower() == ".pptx":
        from .powerpoint import parse_powerpoint

        return parse_powerpoint(deck_path)
    from .parser import parse_markdown

    return parse_markdown(deck_path.read_text(encoding="utf-8"))


# --- Demo deck trimming -------------------------------------------------------
# The open demo renders the FIRST N slides of any deck instead of rejecting
# bigger ones; these helpers cut a deck down in place.


def _split_front_matter(text: str) -> tuple[str, str]:
    """(front matter incl. fences, body). Mirrors parser._strip_front_matter's
    rules so a trimmed deck keeps its YAML header verbatim."""
    lines = text.splitlines()
    if not lines or lines[0].strip() != "---":
        return "", text
    for i in range(1, len(lines)):
        if lines[i].strip() == "---":
            block = [ln for ln in lines[1:i] if ln.strip()]
            if not block:
                return "", text
            for ln in block:
                if ln.lstrip().startswith("#"):
                    return "", text  # a heading: this is a slide
                if not _YAML_KEY_RE.match(ln) and not ln[0].isspace():
                    return "", text  # bullets/prose: a real first slide
            return "\n".join(lines[: i + 1]) + "\n", "\n".join(lines[i + 1 :])
    return "", text


def _truncate_markdown_deck(text: str, max_slides: int) -> str:
    """The first ``max_slides`` slides of a Markdown deck (verbatim).

    Deliberately uses the SAME slide-boundary rules as ``parser.parse_markdown``
    (H1 headings, or ``---`` separators when the deck reads as Marp-style) so
    what we count is exactly what we cut — including its quirks, e.g. a
    ``#`` line inside a code fence still opens a slide in both.
    """
    normalized = text.replace("\r\n", "\n").replace("\r", "\n")
    front, body = _split_front_matter(normalized)
    lines = body.split("\n")

    def _close(text_tail: list[str]) -> str:
        trimmed = ("\n".join(text_tail)).rstrip()
        return (front + trimmed + "\n") if trimmed else (front or "")

    if _use_separator_style(body):
        # Cut at the separator line that opens slide max_slides+1.
        count, has_content, cut = 0, False, None
        for idx, ln in enumerate(lines):
            if ln.strip() == "---":
                if has_content:
                    count += 1
                    if count >= max_slides and cut is None:
                        cut = idx
                    has_content = False
                continue
            if ln.strip():
                has_content = True
        return normalized.rstrip() + "\n" if cut is None else _close(lines[:cut])

    # Heading style: cut right before H1 number max_slides+1.
    count, cut = 0, None
    for idx, ln in enumerate(lines):
        if _H1_RE.match(ln.strip()):
            count += 1
            if count > max_slides:
                cut = idx
                break
    return normalized.rstrip() + "\n" if cut is None else _close(lines[:cut])


def _truncate_pptx_deck(path: Path, max_slides: int) -> bool:
    """Drop slides past ``max_slides`` in place (python-pptx has no public
    delete API; removing sldId entries is the established approach)."""
    try:
        from pptx import Presentation  # type: ignore[import-untyped]

        prs = Presentation(str(path))
        sld_ids = prs.slides._sldIdLst  # noqa: SLF001
        for sld_id in list(sld_ids)[max_slides:]:
            sld_ids.remove(sld_id)
        prs.save(str(path))
    except Exception:
        return False
    return True


def _trim_demo_deck(deck_path: Path, max_slides: int) -> bool:
    """Best-effort in-place trim to the first ``max_slides`` slides."""
    try:
        if deck_path.suffix.lower() == ".pptx":
            return _truncate_pptx_deck(deck_path, max_slides)
        text = deck_path.read_text(encoding="utf-8")
        deck_path.write_text(
            _truncate_markdown_deck(text, max_slides), encoding="utf-8"
        )
    except Exception:
        return False
    return True


# Remote engines that animate a stylized mascot (no human face to detect):
# Wan2.2-S2V has no detector at all, and D-ID handles stylized faces. The
# others (sadtalker/wav2lip/comfyui-auto) start with a human face detector and
# fail on a mascot, so a mascot presenter falls back to the no-GPU puppet there.
_MASCOT_ANIMATE_ENGINES = ("wan-s2v", "d-id")


def _server_animation_engine(base: dict[str, Any]) -> str | None:
    """The server-configured animated-avatar engine, if one is usable.

    Returns the base config's avatar provider name when it selects a remote
    engine that has its connection details set (a ComfyUI ``base_url``, or a
    D-ID ``api_key``), else None. Used to decide whether an animated presenter
    can lip-sync for real or must fall back to the no-GPU puppet mouth-flap.
    """
    av = base.get("providers", {}).get("avatar", {})
    provider = av.get("provider")
    if provider in ("wan-s2v", "sadtalker", "wav2lip", "comfyui"):
        return provider if av.get("base_url") else None
    if provider == "d-id":
        return provider if av.get("api_key") else None
    return None


def _build_job_config(base: dict[str, Any], workdir: Path, options: dict[str, Any],
                      voice_path: Path | None, photo_path: Path | None) -> Path:
    """Write a per-job config YAML: server base + this job's overrides."""
    cfg = copy.deepcopy(base)
    cfg.setdefault("providers", {}).setdefault("tts", {})
    cfg["providers"].setdefault("images", {})
    cfg["providers"].setdefault("avatar", {})
    cfg.setdefault("settings", {})

    # Keep renders self-contained inside the job dir.
    cfg["settings"]["temp_dir"] = str(workdir / "tmp")

    if options.get("narration_seconds"):
        cfg["settings"].setdefault("narration", {})["target_seconds"] = float(
            options["narration_seconds"]
        )
    if options.get("image_provider"):
        cfg["providers"]["images"]["provider"] = options["image_provider"]
    if options.get("accent"):
        cfg["providers"]["tts"]["accent"] = options["accent"]
    # Presenter placement / slide transitions (from the wizard).
    if options.get("avatar_slides"):
        cfg["settings"].setdefault("avatar", {})["slides"] = options["avatar_slides"]
    # Whole-deck presenter reuse (settings.avatar.reuse_clip): one talking-head
    # render looped over every slide instead of one per slide. The web UI no
    # longer offers it — a slide deck is a sequence of discrete units, and the
    # per-slide pipeline (fragment per slide, stitched with transition_seconds)
    # keeps lip-sync in phase with each slide's narration. Still honoured here
    # for YAML configs (settings.avatar.reuse_clip), the CLI (--reuse-avatar),
    # and API clients that post reuse_avatar=true.
    if options.get("reuse_avatar") is not None:
        cfg["settings"].setdefault("avatar", {})["reuse_clip"] = bool(
            options["reuse_avatar"]
        )
    if options.get("transition_seconds"):
        cfg["settings"].setdefault("video", {})["transition_seconds"] = float(
            options["transition_seconds"]
        )

    # A per-job voice sample turns on ephemeral cloning. Chatterbox and
    # Voicebox both accept just the clip: Chatterbox needs no transcript, and
    # Voicebox transcribes it server-side when reference_text is absent.
    # The upload must win: drop inherited stored-voice keys that providers
    # would otherwise prefer (voicebox picks profile_id over voice_sample),
    # and the server config's reference_text describes ITS clip, not this one.
    if voice_path is not None:
        tts = cfg["providers"]["tts"]
        for key in ("profile_id", "voice", "reference_text"):
            tts.pop(key, None)
        tts["voice_sample"] = str(voice_path)

    # A teaser voice choice ("Emily" from providers.tts.voice_choices) picks a
    # STOCK server voice: override any stored voice keys, same as an uploaded
    # sample would. The stem maps back to the configured filename ("Emily" ->
    # "Emily.wav") so operators list plain display names in voice_choices.
    if options.get("voice_name"):
        tts = cfg["providers"]["tts"]
        wanted = str(options["voice_name"])
        choice = next(
            (
                str(c)
                for c in (tts.get("voice_choices") or [])
                if Path(str(c)).stem == wanted
            ),
            wanted,
        )
        for key in ("profile_id", "voice", "reference_text"):
            tts.pop(key, None)
        tts["voice"] = choice

    # Presenter: a built-in mascot wins over an uploaded file. The 'animate'
    # toggle then picks the engine per source:
    #   mascot  + animate -> the server's detector-free engine (wan-s2v/d-id)
    #                        for real AI lip-sync if configured, else static;
    #                        the web UI never falls back to the puppet
    #                        mouth-flap — an animated presenter must be a real
    #                        character model (wan-s2v animates mascots AND
    #                        humans); animate off -> static mascot
    #   photo   + animate -> server's engine (wan-s2v/sadtalker/d-id/comfyui);
    #                        else static photo (a still of themselves)
    #   video             -> always the video engine (a clip is inherently
    #                        animated; wav2lip/comfyui)
    #   nothing           -> no head
    # (The puppet provider stays available to CLI/YAML users via
    # providers.avatar.provider: puppet — it is only the web fallback that is
    # gone.)
    av = cfg["providers"]["avatar"]
    animate = options.get("avatar", True)
    engine = _server_animation_engine(base)
    if options.get("avatar_name"):
        name = options["avatar_name"]
        av["source"] = name
        if animate and engine in _MASCOT_ANIMATE_ENGINES:
            # A real engine can lip-sync the mascot from the narration audio.
            av["provider"] = engine
            if engine == "d-id":
                av["source_image"] = name
        else:
            av["provider"] = "static"
    elif photo_path is not None:
        from .providers.avatar import _source_kind

        av["source"] = str(photo_path)
        if _source_kind(str(photo_path)) == "video":
            av["source_video"] = str(photo_path)
            # A video clip needs the video engine (wan-s2v consumes a still
            # image only), so route it through the auto provider on the same
            # ComfyUI server: photo -> SadTalker, video -> Wav2Lip.
            if engine == "wan-s2v":
                av["provider"] = "comfyui"
        elif animate:
            av["source_image"] = str(photo_path)
        else:
            av["provider"] = "static"
    else:
        av["provider"] = "none"

    # Owner-only from the first byte: the expanded config holds live API keys.
    job_yaml = workdir / "job.yaml"
    fd = os.open(job_yaml, os.O_WRONLY | os.O_CREAT | os.O_TRUNC, 0o600)
    with os.fdopen(fd, "w", encoding="utf-8") as f:
        f.write(yaml.safe_dump(cfg))
    return job_yaml


def _stream_process(command: list[str], job: Job, timeout: float) -> int:
    """Run a subprocess, streaming its output into ``job.log`` as it arrives.

    ``subprocess.run(capture_output=True)`` buffers everything until exit,
    which makes a multi-minute GPU render look dead. A reader thread pumps
    lines into the bounded log and refreshes ``updated_at`` (the heartbeat
    the status endpoint reports as "last output Ns ago").
    """
    proc = subprocess.Popen(
        command, stdout=subprocess.PIPE, stderr=subprocess.STDOUT, text=True
    )
    started = time.monotonic()

    def _pump() -> None:
        stream = proc.stdout
        if stream is None:  # pragma: no cover - PIPE is always set above
            return
        for line in stream:
            with _LOCK:
                job.log = (job.log + line)[-16000:]
                job.updated_at = time.time()

    reader = threading.Thread(target=_pump, daemon=True)
    reader.start()
    try:
        return proc.wait(timeout=max(1.0, timeout - (time.monotonic() - started)))
    except subprocess.TimeoutExpired:
        proc.kill()
        proc.wait()
        raise
    finally:
        reader.join(timeout=5)


# Demo-safe progress derived from the render log. Raw log lines never reach
# the open demo (paths/exception text); these whitelisted, derived labels do.
_SLIDE_LINE_RE = re.compile(r"Slide (\d+)/(\d+)")
_STAGE_PATTERNS: tuple[tuple[str, str], ...] = (
    ("combining", "stitching the final video"),
    ("talking head", "rendering the presenter"),
    ("rendering…", "rendering the presenter"),
    ("avatar", "rendering the presenter"),
    ("audio", "generating voice"),
    ("speaker notes", "writing speaker notes"),
    ("narration", "writing narration"),
    ("image", "generating slide image"),
)


def _job_progress(job: Job) -> dict[str, Any] | None:
    """{'slide': n, 'slides': total, 'stage': label} parsed from job.log, or
    None when nothing recognisable has been printed yet."""
    log = job.log
    if not log:
        return None
    out: dict[str, Any] = {}
    last = None
    for m in _SLIDE_LINE_RE.finditer(log):
        last = m
    if last:
        out["slide"], out["slides"] = int(last.group(1)), int(last.group(2))
    for line in reversed(log.splitlines()[-60:]):
        low = line.lower()
        for needle, label in _STAGE_PATTERNS:
            if needle in low:
                out["stage"] = label
                break
        if "stage" in out:
            break
    return out or None


def _job_warnings(job: Job) -> list[str]:
    """Demo-safe warnings derived from the log — no raw text. Today: the
    talking-head engine failed, so some slides render without their presenter
    (slide + narration are unaffected); and image generation fell back to
    text cards."""
    log = job.log
    if not log:
        return []
    warnings: list[str] = []
    if "avatar error" in log or "avatar generation failed" in log:
        warnings.append(
            "The presenter could not be animated on some slides (the "
            "animation engine failed or timed out) — those slides render "
            "with narration only."
        )
    fallbacks = log.lower().count("using text fallback")
    if fallbacks:
        m = None
        for m_ in _SLIDE_LINE_RE.finditer(log):
            m = m_
        total = f" of {m.group(2)}" if m else ""
        warnings.append(
            f"{fallbacks}{total} slides used a plain text card because "
            "image generation failed — check the image provider."
        )
    text_only = log.lower().count("left text-only")
    if text_only:
        warnings.append(
            f"{text_only} slides had no image (generation failed) — they "
            "remain text-only in the enhanced deck."
        )
    return warnings


_TRACE_IMAGE_OK_RE = re.compile(r"Generated ([\w-]+) image")
_TRACE_AUDIO_RE = re.compile(r"Generated audio with ([^(\n]+)")
_TRACE_HEAD_OK_RE = re.compile(r"Generated (\S+) talking head \(slide (\d+)\)")
_TRACE_HEAD_ERR_RE = re.compile(r"(\S+) avatar error")


def _job_trace(job: Job, limit: int = 40) -> list[str]:
    """A demo-safe, step-by-step trace parsed from the render log: which
    provider produced each slide's image/voice/presenter, and where anything
    fell back. Only whitelisted, derived text — raw lines never leave the
    server."""
    log = job.log
    if not log:
        return []
    events: list[str] = []
    slide: int | None = None
    for raw in log.splitlines():
        line = raw.strip()
        low = line.lower()
        m = _SLIDE_LINE_RE.search(line)
        if m and low.startswith("slide "):
            slide = int(m.group(1))
            continue
        if "using text fallback" in low and slide:
            events.append(f"slide {slide}: image failed — used a text card")
            continue
        if "left text-only" in low and slide:
            events.append(f"slide {slide}: image failed — left text-only")
            continue
        m = _TRACE_IMAGE_OK_RE.search(line)
        if m and slide:
            events.append(f"slide {slide}: image via {m.group(1).capitalize()}")
            continue
        if "wrote ai speaker notes" in low and slide:
            events.append(f"slide {slide}: notes written")
            continue
        m = _TRACE_AUDIO_RE.search(line)
        if m and slide:
            events.append(f"slide {slide}: voice via {m.group(1).strip()}")
            continue
        m = _TRACE_HEAD_OK_RE.search(line)
        if m:
            events.append(f"slide {m.group(2)}: presenter animated ({m.group(1)})")
            continue
        m = _TRACE_HEAD_ERR_RE.search(line)
        if m and slide:
            events.append(f"slide {slide}: presenter failed ({m.group(1)})")
            continue
        if "combining video fragments" in low:
            events.append("stitching slides together")
    return events[-limit:]


def _run_job(job: Job, deck_path: Path, job_yaml: Path,
             voice_path: Path | None, photo_path: Path | None,
             mode: str = "video", notes: str | None = None) -> None:
    """Run one job in a subprocess, then wipe the biometric inputs.

    ``mode`` 'video' renders an MP4 (``create``); 'pptx' produces an enriched
    deck zip (``enrich --pptx --zip``), optionally with AI presenter notes.
    """
    assert job.workdir is not None
    if mode == "pptx":
        out_dir = job.workdir / "enriched"
        output = job.workdir / "enriched.zip"  # enrich --zip writes <dir>.zip
        command = [sys.executable, "-m", "slide_stream", "enrich",
                   str(deck_path), str(out_dir), "--config", str(job_yaml),
                   "--pptx", "--zip"]
        if notes in ("fill", "all"):
            command += ["--notes", notes]
        media_type, download_name = "application/zip", "slidestream-deck.zip"
    else:
        output = job.workdir / "output.mp4"
        command = [sys.executable, "-m", "slide_stream", "create",
                   str(deck_path), str(output), "--config", str(job_yaml)]
        media_type, download_name = "video/mp4", "slidestream.mp4"

    with _LOCK:
        job.status = "running"
        job.media_type = media_type
        job.download_name = download_name
    try:
        returncode = _stream_process(command, job, 3600)
        with _LOCK:
            if returncode == 0 and output.exists():
                job.status = "done"
                job.output_path = output
            else:
                job.status = "error"
                job.error = (f"{'enrich' if mode == 'pptx' else 'render'} "
                             f"exited {returncode}")
    except subprocess.TimeoutExpired:
        with _LOCK:
            job.status = "error"
            job.error = "render timed out"
    except Exception as e:  # pragma: no cover - defensive
        with _LOCK:
            job.status = "error"
            job.error = str(e)
    finally:
        # Ephemeral: inputs and the key-bearing job config never persist past
        # the render. Only output.mp4 remains, until it is downloaded (demo)
        # or the TTL reaper removes the whole workdir.
        for p in (voice_path, photo_path, deck_path, job_yaml):
            if p is not None:
                Path(p).unlink(missing_ok=True)
        shutil.rmtree(job.workdir / "tmp", ignore_errors=True)


def _run_video_pipeline(job: Job, deck_path: Path, job_yaml: Path,
                        voice_path: Path | None, photo_path: Path | None) -> None:
    """Demo mascot-video job: enrich FIRST (each slide's image generated once,
    failures isolated per slide), then create the video FROM the enhanced
    deck — which reuses those images instead of regenerating art mid-video."""
    assert job.workdir is not None
    with _LOCK:
        job.status = "running"
    try:
        enhanced = job.workdir / "enhanced"
        code = _stream_process(
            [sys.executable, "-m", "slide_stream", "enrich",
             str(deck_path), str(enhanced), "--config", str(job_yaml)],
            job, 1800,
        )
        enriched_md = enhanced / f"{deck_path.stem}.md"
        if code != 0 or not enriched_md.exists():
            with _LOCK:
                job.status = "error"
                job.error = (
                    f"enrich exited {code}" if code else "enrich produced no deck"
                )
            return
        video = job.workdir / "output.mp4"
        code = _stream_process(
            [sys.executable, "-m", "slide_stream", "create",
             str(enriched_md), str(video), "--config", str(job_yaml)],
            job, 3600,
        )
        with _LOCK:
            if code == 0 and video.exists():
                job.status = "done"
                job.output_path = video
            else:
                job.status = "error"
                job.error = f"create exited {code}"
    except subprocess.TimeoutExpired:
        with _LOCK:
            job.status = "error"
            job.error = "render timed out"
    except Exception as e:  # pragma: no cover - defensive
        with _LOCK:
            job.status = "error"
            job.error = str(e)
    finally:
        for p in (voice_path, photo_path, deck_path, job_yaml):
            if p is not None:
                Path(p).unlink(missing_ok=True)
        shutil.rmtree(job.workdir / "tmp", ignore_errors=True)


def _run_project_enrich(job: Job, project: Project, deck_path: Path,
                        job_yaml: Path, notes: str | None) -> None:
    """Enrich a project's deck into a downloadable zip AND copy the generated
    images back into the project's ``images/`` so the UI can preview them and a
    later render can reuse them."""
    assert job.workdir is not None
    out_dir = job.workdir / "enriched"
    output = job.workdir / "enriched.zip"  # enrich --zip writes <dir>.zip
    command = [sys.executable, "-m", "slide_stream", "enrich",
               str(deck_path), str(out_dir), "--config", str(job_yaml),
               "--pptx", "--zip"]
    if notes in ("fill", "all"):
        command += ["--notes", notes]
    with _LOCK:
        job.status = "running"
        job.media_type = "application/zip"
        job.download_name = "slidestream-deck.zip"
    try:
        returncode = _stream_process(command, job, 3600)
        if returncode == 0 and output.exists():
            src_images = out_dir / "images"
            if src_images.is_dir():
                proj_images = project.workdir / "images"
                proj_images.mkdir(exist_ok=True)
                for img in src_images.iterdir():
                    if img.is_file():
                        shutil.copyfile(img, proj_images / img.name)
            with _LOCK:
                job.status = "done"
                job.output_path = output
        else:
            with _LOCK:
                job.status = "error"
                job.error = f"enrich exited {returncode}"
    except subprocess.TimeoutExpired:
        with _LOCK:
            job.status = "error"
            job.error = "enrich timed out"
    except Exception as e:  # pragma: no cover - defensive
        with _LOCK:
            job.status = "error"
            job.error = str(e)
    finally:
        for p in (deck_path, job_yaml):
            Path(p).unlink(missing_ok=True)


def _do_draft(source_path: Path, slides: int | None, provider: str,
              model: str | None, base_url: str | None,
              api_key: str | None = None, think: bool | None = None) -> str:
    """Extract a document and draft deck Markdown from it (blocking: offload to
    a threadpool). Raises DraftError / ValueError with a user-facing message."""
    from .draft import (
        DraftError,
        build_draft_prompt,
        clamp_source,
        clean_llm_markdown,
        extract_source_text,
        validate_deck_markdown,
    )
    from .llm import get_llm_client, query_llm

    source_text = extract_source_text(source_path)
    if not source_text.strip():
        raise DraftError(
            "No extractable text was found in the document "
            "(a scanned/image-only PDF, perhaps?)."
        )
    source_text, _ = clamp_source(source_text)

    import io

    from rich.console import Console

    client = get_llm_client(provider, base_url=base_url, api_key=api_key)
    # Capture the LLM layer's progress/error prints so a failure can name the
    # actual cause (e.g. a bad API key) in the 400 the UI shows.
    quiet_console = Console(file=io.StringIO(), record=True)
    result = query_llm(
        client, provider, build_draft_prompt(source_text, slides),
        quiet_console, model, think=think,
    )
    if not result:
        raise DraftError(_llm_failure_hint(quiet_console))
    deck_markdown = clean_llm_markdown(result)
    validate_deck_markdown(deck_markdown)  # raises DraftError if unusable
    return deck_markdown.rstrip() + "\n"


def _do_topic_draft(topic: str, slides: int | None, provider: str,
                    model: str | None, base_url: str | None,
                    api_key: str | None = None, think: bool | None = None) -> str:
    """Draft a deck from a typed idea/topic (blocking: offload to a
    threadpool). Raises DraftError / ValueError with a user-facing message."""
    import io

    from rich.console import Console

    from .draft import (
        DraftError,
        build_topic_prompt,
        clean_llm_markdown,
        validate_deck_markdown,
    )
    from .llm import get_llm_client, query_llm

    client = get_llm_client(provider, base_url=base_url, api_key=api_key)
    quiet_console = Console(file=io.StringIO(), record=True)
    result = query_llm(
        client, provider, build_topic_prompt(topic, slides),
        quiet_console, model, think=think,
    )
    if not result:
        raise DraftError(_llm_failure_hint(quiet_console))
    deck_markdown = clean_llm_markdown(result)
    validate_deck_markdown(deck_markdown)
    return deck_markdown.rstrip() + "\n"


def _llm_failure_hint(console: Any) -> str:
    """A user-facing reason for an empty LLM result, from the captured
    console output (the LLM layer prints the provider's error there)."""
    lines = [
        ln.strip()
        for ln in str(console.export_text()).splitlines()
        if ln.strip()
    ]
    hint = lines[-1] if lines else ""
    if hint:
        return f"The LLM returned no content. Last message: {hint}"
    return "The LLM returned no content. Try again."


# Drafting from a topic is cheap relative to a render, but it is still an LLM
# call on an open endpoint — give it its own per-IP budget.
DEMO_DRAFTS_PER_HOUR = 10
_DEMO_DRAFT_HITS: dict[str, list[float]] = {}


def _demo_draft_rate_ok(ip: str, now: float | None = None) -> bool:
    """True if this IP may start another topic draft (and record the hit)."""
    t = now if now is not None else time.time()
    with _LOCK:
        hits = [h for h in _DEMO_DRAFT_HITS.get(ip, []) if t - h < 3600]
        if len(hits) >= DEMO_DRAFTS_PER_HOUR:
            _DEMO_DRAFT_HITS[ip] = hits
            return False
        hits.append(t)
        _DEMO_DRAFT_HITS[ip] = hits
        return True


SETTINGS_TEMPLATE = """\
# SlideStream settings (~/.slidestream.yaml)
# Uncomment and edit what you use. Keys can reference environment variables
# as ${VAR}. Full docs: https://github.com/michael-borck/slide-stream
providers:
  llm:
    provider: gemini            # gemini | openai | claude | groq | ollama | none
    model: gemini-2.0-flash
    # Ollama (incl. behind an authenticating proxy) or any OpenAI-compatible
    # server:
    # base_url: https://ollama.example.org
    # api_key: "${OLLAMA_TOKEN}"   # sent as a Bearer token
    # think: false                 # reasoning models: keep answers clean
  tts:
    provider: voicebox          # default; or gtts (free), kokoro (offline),
                                # chatterbox, elevenlabs, openai
    # base_url: https://voice.example.org        # voicebox/chatterbox server
    # api_key: "${VOICEBOX_TOKEN}"               # if the server needs auth
    # engine: kokoro            # voicebox: kokoro|chatterbox|qwen|luxtts|tada
    # profile_id: "<id from POST /profiles>"     # a stored voicebox voice...
    # voice_sample: /path/to/you.wav  # ...or clone this clip per run, then
    #                                 # delete it from the server
    # accent: australian        # gtts: australian|british|american|...
  images:
    provider: text              # text (no AI) | gemini | dalle3 | swarmui |
                                # local | pexels | unsplash
    # base_url: https://swarmui.example.org
    # model: juggernautXL_v9
  avatar:
    provider: none              # none | static | puppet | wan-s2v |
                                # sadtalker | wav2lip | comfyui | d-id
    # source: teddy             # built-in mascot, or a photo/video path
    # base_url: https://comfyui.example.org
    # api_key: "${COMFYUI_TOKEN}"   # if the ComfyUI server needs auth
    # wan-s2v animates mascots AND human head shots (no face detector);
    # sadtalker/wav2lip/comfyui are human-faces-only.
settings:
  strict: false
  narration:
    target_seconds: 45
"""


def create_app(config: dict[str, Any] | None = None, token: str | None = None,
               max_workers: int = 1, demo: bool | None = None,
               local: bool | None = None):
    """Build the FastAPI app. Requires the ``[serve]`` extra.

    ``demo`` (or the ``SLIDESTREAM_DEMO`` env var) shows a banner in the UI
    inviting users to install locally for full control over the LLM, image,
    and video generation — used on the hosted VPS instance.

    ``local`` (or ``SLIDESTREAM_LOCAL=1``) is desktop/laptop mode: no token,
    no demo limits, and a Settings page that edits ~/.slidestream.yaml. Used
    by the Tauri desktop shell.
    """
    try:
        from fastapi import (
            Depends,
            FastAPI,
            File,
            Form,
            Header,
            HTTPException,
            Request,
            UploadFile,
        )
        from fastapi.responses import FileResponse, HTMLResponse, JSONResponse
        from starlette.background import BackgroundTask
        from starlette.responses import Response
    except ImportError as e:  # pragma: no cover
        raise RuntimeError(
            "The web UI needs extra packages. Install with: "
            'pip install "slide-stream[serve]"'
        ) from e

    base_config = config if config is not None else load_config()
    if local is None:
        local = os.getenv("SLIDESTREAM_LOCAL", "").lower() in ("1", "true", "yes")
    local_mode = bool(local)
    auth_token = "" if local_mode else (token or os.getenv("SLIDESTREAM_TOKEN") or "")
    if demo is None:
        demo = os.getenv("SLIDESTREAM_DEMO", "").lower() in ("1", "true", "yes")
    demo_mode = bool(demo) and not local_mode
    # Only honor X-Forwarded-For when the operator says a reverse proxy sits
    # in front (deploy/docker-compose.yml sets this); otherwise the header is
    # client-supplied and would let anyone dodge the demo rate limit.
    trusted_proxy = os.getenv(
        "SLIDESTREAM_TRUSTED_PROXY", ""
    ).lower() in ("1", "true", "yes")
    executor = ThreadPoolExecutor(max_workers=max_workers)
    jobs_root = Path(tempfile.mkdtemp(prefix="slidestream_serve_"))

    app = FastAPI(title="SlideStream")

    if local_mode:

        @app.middleware("http")
        async def local_guard(
            request: Request,
            call_next: Callable[[Request], Awaitable[Response]],
        ) -> Response:
            # Local (desktop) mode has no token, so block CSRF/DNS-rebinding:
            # state-changing requests must target localhost and come from a
            # local page (or the Tauri shell's webview).
            if request.method in ("POST", "PUT", "PATCH", "DELETE"):
                host = (
                    (request.headers.get("host") or "")
                    .rsplit(":", 1)[0].strip("[]").lower()
                )
                if host not in _LOCAL_HOSTS:
                    return JSONResponse(
                        {"detail": "Requests must be addressed to localhost"},
                        status_code=403,
                    )
                origin = request.headers.get("origin")
                if origin and not _local_origin_ok(origin):
                    return JSONResponse(
                        {"detail": "Cross-origin requests are not allowed"},
                        status_code=403,
                    )
            return await call_next(request)

    def require_token(authorization: str | None = Header(default=None)) -> None:
        # Demo mode is friction-free: no token, guarded by rate/slide limits
        # instead. A token only gates private/full instances.
        if demo_mode or not auth_token:
            return
        expected = f"Bearer {auth_token}"
        if authorization is None or not secrets.compare_digest(
            authorization.encode(), expected.encode()
        ):
            raise HTTPException(status_code=401, detail="Invalid or missing token")

    def client_ip(request: Request) -> str:
        # Behind a trusted reverse proxy the real IP is the value the proxy
        # appended to X-Forwarded-For (the rightmost one — earlier entries
        # are whatever the client chose to send).
        if trusted_proxy:
            fwd = request.headers.get("x-forwarded-for")
            if fwd:
                return fwd.split(",")[-1].strip()
        return request.client.host if request.client else "unknown"

    @app.get("/", response_class=HTMLResponse)
    def index() -> str:
        return INDEX_HTML

    @app.get("/api/config")
    def api_config() -> dict[str, Any]:
        # Public so the UI can bootstrap: token/demo, and the choices this
        # server actually supports (built-in avatars; accents only if the
        # configured TTS provider offers them — currently gTTS; stock voice
        # names from providers.tts.voice_choices for the teaser's picker).
        from .avatars import avatar_names
        from .providers.tts import GTTS_ACCENTS

        tts_config = base_config.get("providers", {}).get("tts", {})
        voices = [
            Path(str(v)).stem for v in (tts_config.get("voice_choices") or [])
        ]
        return {
            "version": __version__,
            "auth_required": bool(auth_token) and not demo_mode,
            "demo": demo_mode,
            "local": local_mode,
            "limits": (
                {
                    "max_slides": DEMO_MAX_SLIDES,
                    "jobs_per_hour": DEMO_JOBS_PER_HOUR,
                    "job_ttl_minutes": JOB_TTL_SECONDS // 60,
                }
                if demo_mode
                else None
            ),
            "avatars": avatar_names(),
            "voices": voices,
            "accents": list(GTTS_ACCENTS) if tts_config.get("provider") == "gtts" else [],
            # The UI offers a PowerPoint output; AI presenter notes need an LLM.
            "llm": base_config.get("providers", {}).get("llm", {}).get("provider", "none") != "none",
        }

    @app.get("/api/avatars/{name}")
    def avatar_image(name: str):
        # Bundled mascot art for the teaser's thumbnail picker. Public: the
        # images ship with the package and carry nothing sensitive, and an
        # <img> tag cannot send the instance token.
        from .avatars import avatar_names, resolve_avatar

        if name not in avatar_names():
            raise HTTPException(status_code=404, detail="Unknown avatar")
        path = resolve_avatar(name)
        if not path or not Path(path).is_file():
            raise HTTPException(status_code=404, detail="No image")
        return FileResponse(path)

    settings_path = Path.home() / ".slidestream.yaml"

    @app.post("/api/quit")
    def quit_app() -> dict[str, Any]:
        # Desktop mode only: the Tauri shell calls this when the window
        # closes so the sidecar server (and its render subprocesses' parent)
        # exits cleanly even if the process kill only reaches the launcher.
        if not local_mode:
            raise HTTPException(status_code=404, detail="Not available")
        threading.Timer(0.2, _shutdown).start()
        return {"ok": True}

    @app.get("/api/settings")
    def get_settings() -> dict[str, Any]:
        # Desktop mode only: read the user's config for the Settings page.
        if not local_mode:
            raise HTTPException(status_code=404, detail="Not available")
        text = ""
        if settings_path.exists():
            text = settings_path.read_text(encoding="utf-8")
        return {"path": str(settings_path), "yaml": text,
                "template": SETTINGS_TEMPLATE}

    @app.put("/api/settings")
    async def put_settings(request: Request) -> dict[str, Any]:
        if not local_mode:
            raise HTTPException(status_code=404, detail="Not available")
        body = await request.json()
        text = body.get("yaml", "")
        try:
            parsed = yaml.safe_load(text) if text.strip() else None
            if parsed is not None and not isinstance(parsed, dict):
                raise ValueError("top level must be a mapping")
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Invalid YAML: {e}") from e
        # Owner-only: the settings file typically holds API keys. os.open's
        # mode only applies on creation, so chmod covers pre-existing files.
        fd = os.open(settings_path, os.O_WRONLY | os.O_CREAT | os.O_TRUNC, 0o600)
        with os.fdopen(fd, "w", encoding="utf-8") as f:
            f.write(text)
        os.chmod(settings_path, 0o600)
        return {"ok": True, "path": str(settings_path)}

    async def save_upload(
        upload: UploadFile, dest: Path, max_bytes: int, kind: str
    ) -> None:
        # Stream to disk in chunks with a hard cap — never buffer in RAM.
        size = 0
        with open(dest, "wb") as f:
            while chunk := await upload.read(1024 * 1024):
                size += len(chunk)
                if size > max_bytes:
                    raise HTTPException(
                        status_code=413,
                        detail=f"{kind} exceeds the "
                        f"{max_bytes // (1024 * 1024)} MB limit",
                    )
                f.write(chunk)

    @app.post("/api/jobs")
    async def create_job(
        request: Request,
        deck: UploadFile = File(...),
        voice: UploadFile | None = File(default=None),
        photo: UploadFile | None = File(default=None),
        narration_seconds: str | None = Form(default=None),
        image_provider: str | None = Form(default=None),
        avatar: str | None = Form(default=None),
        avatar_name: str | None = Form(default=None),
        avatar_slides: str | None = Form(default=None),
        reuse_avatar: str | None = Form(default=None),
        transition: str | None = Form(default=None),
        accent: str | None = Form(default=None),
        voice_name: str | None = Form(default=None),
        output: str | None = Form(default=None),
        notes: str | None = Form(default=None),
        _: None = Depends(require_token),
    ) -> JSONResponse:
        _reap_expired_jobs()
        suffix = Path(deck.filename or "deck.md").suffix.lower()
        if suffix not in (".md", ".pptx", ".txt", ".qmd"):
            raise HTTPException(
                status_code=400,
                detail="Deck must be .md, .pptx, .txt or .qmd",
            )
        mode = "pptx" if (output or "video").lower() == "pptx" else "video"
        notes_mode = (notes or "").lower() if notes else None
        if notes_mode not in (None, "fill", "all"):
            raise HTTPException(status_code=400, detail="notes must be 'fill' or 'all'")

        # Reject obviously oversized requests before touching the body.
        declared = request.headers.get("content-length")
        max_request = MAX_DECK_BYTES + MAX_VOICE_BYTES + MAX_PHOTO_BYTES + 1024 * 1024
        if declared and declared.isdigit() and int(declared) > max_request:
            raise HTTPException(status_code=413, detail="Upload too large")

        if demo_mode and not _demo_rate_ok(client_ip(request)):
            raise HTTPException(
                status_code=429,
                detail=f"Demo limit: {DEMO_JOBS_PER_HOUR} videos per hour. "
                "Install locally for unlimited renders: pip install slide-stream",
            )

        job_id = uuid.uuid4().hex
        workdir = jobs_root / job_id
        workdir.mkdir(parents=True)
        demo_notice: str | None = None
        try:
            # Text formats are parsed as markdown downstream, so they are
            # stored under the canonical .md name whatever they uploaded as.
            deck_path = workdir / (
                "deck.pptx" if suffix == ".pptx" else "deck.md"
            )
            await save_upload(deck, deck_path, MAX_DECK_BYTES, "Deck")

            if demo_mode:
                n = _count_slides(deck_path)
                if n is None:
                    # Fail closed: an unparseable deck must not dodge the cap.
                    raise HTTPException(
                        status_code=400, detail="Could not parse the deck"
                    )
                if n == 0:
                    raise HTTPException(
                        status_code=400,
                        detail="No slides found. Start each slide with a "
                        "'# ' heading (or separate them with '---' lines), "
                        "then upload again.",
                    )
                if n > DEMO_MAX_SLIDES:
                    # Rather than reject bigger decks, render their first
                    # DEMO_MAX_SLIDES slides and say so — a taste of the full
                    # product beats a dead end.
                    if not _trim_demo_deck(deck_path, DEMO_MAX_SLIDES):
                        raise HTTPException(
                            status_code=400,
                            detail=f"Demo limit: this deck has {n} slides and "
                            f"could not be trimmed to {DEMO_MAX_SLIDES}.",
                        )
                    demo_notice = (
                        f"Your deck has {n} slides — the demo renders the "
                        f"first {DEMO_MAX_SLIDES}. Download the desktop app "
                        "to narrate or enhance the full presentation."
                    )

            voice_path = None
            photo_path = None
            if not demo_mode:
                # The open teaser accepts no biometrics: no voice clones, no
                # face photos — stock voices and mascots only. The UI hides
                # these inputs in demo mode; this holds even for API clients.
                if voice is not None and voice.filename:
                    voice_path = (
                        workdir / f"voice{Path(voice.filename).suffix or '.wav'}"
                    )
                    await save_upload(
                        voice, voice_path, MAX_VOICE_BYTES, "Voice sample"
                    )
                if photo is not None and photo.filename:
                    from .providers.avatar import _source_kind

                    photo_path = (
                        workdir / f"photo{Path(photo.filename).suffix or '.png'}"
                    )
                    await save_upload(photo, photo_path, MAX_PHOTO_BYTES, "Photo")
                    if _source_kind(str(photo_path)) == "image":
                        problem = _validate_photo_upload(photo_path)
                        if problem:
                            raise HTTPException(status_code=400, detail=problem)

            options = {
                "narration_seconds": narration_seconds,
                "image_provider": image_provider,
                "avatar": (avatar or "true").lower() != "false",
                "avatar_name": avatar_name,
                "avatar_slides": avatar_slides,
                "reuse_avatar": (reuse_avatar or "").lower() == "true"
                if reuse_avatar is not None else None,
                "transition_seconds": transition,
                "accent": accent,
                "voice_name": voice_name,
            }
            if demo_mode and not avatar_slides:
                # Teaser economics: a wan-s2v clip costs minutes of GPU per
                # slide, so demo presenters appear on the opening and closing
                # slides unless a client explicitly asks for more.
                options["avatar_slides"] = "first,last"
            # Desktop mode re-reads ~/.slidestream.yaml per job so Settings
            # edits apply without restarting the app.
            job_base = load_config() if local_mode else base_config
            job_yaml = _build_job_config(
                job_base, workdir, options, voice_path, photo_path
            )
        except Exception:
            # Never leave a rejected job's uploads (or its key-bearing
            # job.yaml) on disk.
            shutil.rmtree(workdir, ignore_errors=True)
            raise

        job = Job(id=job_id, workdir=workdir, created_at=time.time(),
                  download_token=secrets.token_urlsafe(24))
        with _LOCK:
            _JOBS[job_id] = job
        if demo_mode and mode == "video":
            # Teaser mascot videos run the robust two-pass: enrich first
            # (images generated once, per-slide failures isolated), then
            # create FROM the enhanced deck.
            executor.submit(_run_video_pipeline, job, deck_path, job_yaml,
                            voice_path, photo_path)
        else:
            executor.submit(_run_job, job, deck_path, job_yaml, voice_path,
                            photo_path, mode, notes_mode)
        return JSONResponse({"job_id": job_id, "status": job.status,
                             "token": job.download_token,
                             "notice": demo_notice})

    @app.post("/api/check")
    async def check_deck(
        deck: UploadFile = File(...),
        voice: UploadFile | None = File(default=None),
        photo: UploadFile | None = File(default=None),
        narration_seconds: str | None = Form(default=None),
        image_provider: str | None = Form(default=None),
        avatar: str | None = Form(default=None),
        avatar_name: str | None = Form(default=None),
        accent: str | None = Form(default=None),
        voice_name: str | None = Form(default=None),
        output: str | None = Form(default=None),
        _: None = Depends(require_token),
    ) -> JSONResponse:
        """Offline preflight (the 'doctor'): assess the deck + resolved config
        and return warnings + estimates as JSON, without rendering anything."""
        suffix = Path(deck.filename or "deck.md").suffix.lower()
        if suffix not in (".md", ".pptx", ".txt", ".qmd"):
            raise HTTPException(
                status_code=400,
                detail="Deck must be .md, .pptx, .txt or .qmd",
            )

        workdir = jobs_root / ("check_" + uuid.uuid4().hex)
        workdir.mkdir(parents=True)
        try:
            deck_path = workdir / f"deck{suffix}"
            await save_upload(deck, deck_path, MAX_DECK_BYTES, "Deck")
            voice_path = None
            if voice is not None and voice.filename:
                voice_path = workdir / f"voice{Path(voice.filename).suffix or '.wav'}"
                await save_upload(voice, voice_path, MAX_VOICE_BYTES, "Voice sample")
            photo_path = None
            if photo is not None and photo.filename:
                photo_path = workdir / f"photo{Path(photo.filename).suffix or '.png'}"
                await save_upload(photo, photo_path, MAX_PHOTO_BYTES, "Photo")

            options = {
                "narration_seconds": narration_seconds,
                "image_provider": image_provider,
                "avatar": (avatar or "true").lower() != "false",
                "avatar_name": avatar_name,
                "accent": accent,
                "voice_name": voice_name,
            }
            job_base = load_config() if local_mode else base_config
            job_yaml = _build_job_config(job_base, workdir, options, voice_path, photo_path)
            cfg = yaml.safe_load(job_yaml.read_text(encoding="utf-8")) or {}

            slides = _parse_deck_slides(deck_path)
            if not slides:
                raise HTTPException(status_code=400, detail="Could not parse the deck")

            from .doctor import run_doctor

            avatar_enabled = (
                cfg.get("providers", {}).get("avatar", {}).get("provider", "none") != "none"
            )
            report = run_doctor(slides, cfg, {
                "mode": "pptx" if (output or "video").lower() == "pptx" else "create",
                "input_ext": suffix,
                "verbatim_notes": False,
                "script_blocks": None,
                "avatar_enabled": avatar_enabled,
                "narration_seconds": float(narration_seconds) if narration_seconds else None,
                "output_path": None,
            })
            return JSONResponse({
                "blockers": report.blockers,
                "warnings": report.warnings,
                "findings": [
                    {"group": f.group, "severity": f.severity, "message": f.message}
                    for f in report.findings
                ],
                "estimates": report.estimates,
            })
        finally:
            shutil.rmtree(workdir, ignore_errors=True)

    @app.post("/api/draft-deck")
    async def draft_deck_from_topic(
        request: Request,
        _: None = Depends(require_token),
    ) -> JSONResponse:
        """Draft a deck from a typed idea/topic. Stateless (no project): the
        caller decides what to do with the Markdown — the teaser posts it
        straight back as a job deck."""
        body = await request.json()
        topic = str(body.get("topic", "")).strip()
        if len(topic) < 3:
            raise HTTPException(
                status_code=400, detail="Give the topic a few more words."
            )
        if len(topic) > 2000:
            raise HTTPException(
                status_code=400, detail="Topic is too long (2000 characters max)."
            )
        slides: int | None = None
        raw_slides = body.get("slides")
        if raw_slides not in (None, ""):
            try:
                slides = int(raw_slides)
            except (TypeError, ValueError):
                raise HTTPException(
                    status_code=400, detail="slides must be a number"
                )
            if not 1 <= slides <= 20:
                raise HTTPException(
                    status_code=400, detail="slides must be between 1 and 20"
                )
        if demo_mode:
            slides = min(slides or DEMO_MAX_SLIDES, DEMO_MAX_SLIDES)
            if not _demo_draft_rate_ok(client_ip(request)):
                raise HTTPException(
                    status_code=429,
                    detail=f"Demo limit: {DEMO_DRAFTS_PER_HOUR} drafts per "
                    "hour. Install locally for unlimited drafting: "
                    "pip install slide-stream",
                )
        job_base = load_config() if local_mode else base_config
        llm = job_base.get("providers", {}).get("llm", {})
        provider = llm.get("provider", "none")
        if provider == "none":
            raise HTTPException(
                status_code=400,
                detail="Drafting needs an LLM provider configured in Settings "
                "(e.g. claude, openai, gemini).",
            )

        from starlette.concurrency import run_in_threadpool

        from .draft import DraftError

        try:
            markdown = await run_in_threadpool(
                _do_topic_draft, topic, slides, provider, llm.get("model"),
                llm.get("base_url"), llm.get("api_key"), llm.get("think"),
            )
        except (DraftError, ValueError, ImportError) as e:
            raise HTTPException(status_code=400, detail=str(e))
        return JSONResponse({"markdown": markdown, "slides": slides})

    @app.get("/api/jobs/{job_id}")
    def job_status(job_id: str, _: None = Depends(require_token)) -> dict[str, Any]:
        _reap_expired_jobs()
        job = _JOBS.get(job_id)
        if job is None:
            raise HTTPException(status_code=404, detail="Unknown job")
        out: dict[str, Any] = {"job_id": job.id, "status": job.status,
                               "error": job.error}
        warnings = _job_warnings(job)
        if warnings:
            out["warnings"] = warnings
        trace = _job_trace(job)
        if trace:
            out["trace"] = trace
        if job.status in ("queued", "running"):
            # Demo-safe derived progress (slide x/y + whitelisted stage label)
            # and a heartbeat: seconds since the render last printed anything.
            progress = _job_progress(job)
            if progress:
                out["progress"] = progress
            out["idle"] = max(0, int(time.time() - (job.updated_at or job.created_at)))
        if demo_mode:
            # The open demo gets only the coarse status the UI shows anyway:
            # raw render logs can leak paths and exception text, and the
            # download token must stay knowable only to the job's creator.
            out["log"] = ""
        else:
            out["log"] = job.log[-4000:]
            out["token"] = job.download_token
        return out

    @app.get("/api/jobs/{job_id}/result")
    def job_result(job_id: str, t: str | None = None,
                   authorization: str | None = Header(default=None)):
        _reap_expired_jobs()
        job = _JOBS.get(job_id)
        if job is None:
            raise HTTPException(status_code=404, detail="Result not ready")
        # A browser download link can't set an Authorization header, so accept
        # the per-job download token via ?t= — never the long-lived instance
        # token, which must not land in proxy logs or browser history.
        header_ok = bool(
            auth_token
            and authorization is not None
            and secrets.compare_digest(
                authorization.encode(), f"Bearer {auth_token}".encode()
            )
        )
        token_ok = bool(
            t and job.download_token
            and secrets.compare_digest(t.encode(), job.download_token.encode())
        )
        if not (header_ok or token_ok):
            raise HTTPException(
                status_code=401, detail="Invalid or missing download token"
            )
        if job.status != "done" or job.output_path is None:
            raise HTTPException(status_code=404, detail="Result not ready")
        background: BackgroundTask | None = None
        if demo_mode:
            # Demo promises "nothing stored": one download, then the whole
            # job dir goes away.
            workdir = job.workdir

            def _wipe() -> None:
                with _LOCK:
                    _JOBS.pop(job_id, None)
                if workdir is not None:
                    shutil.rmtree(workdir, ignore_errors=True)

            background = BackgroundTask(_wipe)
        return FileResponse(str(job.output_path), media_type=job.media_type,
                            filename=job.download_name, background=background)

    # ---- Project workflow: draft -> edit -> enrich -> render --------------
    # A stateful session (unlike the one-shot jobs above) so a document can
    # become a deck, be edited, enriched, and rendered without re-uploading
    # between steps. Gated to non-demo mode: it needs an LLM and persists state,
    # so it's the desktop / self-hosted experience, not the open demo.
    projects_root = jobs_root / "projects"
    projects_root.mkdir(exist_ok=True)

    def _project_or_401(project_id: str, token: str | None) -> Project:
        if demo_mode:
            raise HTTPException(
                status_code=403,
                detail="The project workflow is available in the app / "
                "self-hosted mode, not the open demo.",
            )
        _reap_expired_projects()
        project = _PROJECTS.get(project_id)
        if project is None:
            raise HTTPException(status_code=404, detail="Unknown project")
        if not (
            token
            and secrets.compare_digest(token.encode(), project.token.encode())
        ):
            raise HTTPException(
                status_code=401, detail="Invalid or missing project token"
            )
        return project

    @app.post("/api/projects")
    async def create_project(
        deck: UploadFile | None = File(default=None),
        _: None = Depends(require_token),
    ) -> JSONResponse:
        if demo_mode:
            raise HTTPException(
                status_code=403,
                detail="The project workflow is available in the app / "
                "self-hosted mode, not the open demo.",
            )
        _reap_expired_projects()
        pid = uuid.uuid4().hex
        workdir = projects_root / pid
        workdir.mkdir(parents=True)
        try:
            if deck is not None and deck.filename:
                suffix = Path(deck.filename).suffix.lower()
                if suffix not in (".md", ".pptx"):
                    raise HTTPException(
                        status_code=400, detail="Deck must be .md or .pptx"
                    )
                await save_upload(
                    deck, workdir / f"deck{suffix}", MAX_DECK_BYTES, "Deck"
                )
        except Exception:
            shutil.rmtree(workdir, ignore_errors=True)
            raise
        project = Project(
            id=pid, workdir=workdir, created_at=time.time(),
            token=secrets.token_urlsafe(24),
        )
        with _LOCK:
            _PROJECTS[pid] = project
        return JSONResponse(
            {"project_id": pid, "token": project.token,
             "state": _project_state(project)}
        )

    @app.get("/api/projects/{project_id}")
    def project_state_endpoint(
        project_id: str,
        x_project_token: str | None = Header(default=None),
        _: None = Depends(require_token),
    ) -> dict[str, Any]:
        return _project_state(_project_or_401(project_id, x_project_token))

    @app.post("/api/projects/{project_id}/draft")
    async def project_draft(
        project_id: str,
        source: UploadFile = File(...),
        slides: str | None = Form(default=None),
        x_project_token: str | None = Header(default=None),
        _: None = Depends(require_token),
    ) -> JSONResponse:
        project = _project_or_401(project_id, x_project_token)
        job_base = load_config() if local_mode else base_config
        llm = job_base.get("providers", {}).get("llm", {})
        provider = llm.get("provider", "none")
        if provider == "none":
            raise HTTPException(
                status_code=400,
                detail="draft needs an LLM provider configured in Settings "
                "(e.g. claude, openai, gemini).",
            )
        n: int | None = None
        if slides:
            try:
                n = int(slides)
            except ValueError:
                raise HTTPException(status_code=400, detail="slides must be a number")
            if n < 1:
                raise HTTPException(status_code=400, detail="slides must be positive")

        from .draft import SUPPORTED_SUFFIXES

        suffix = Path(source.filename or "source").suffix.lower()
        if suffix not in SUPPORTED_SUFFIXES:
            raise HTTPException(
                status_code=400,
                detail=f"Unsupported source type. Use: {', '.join(SUPPORTED_SUFFIXES)}",
            )
        src_path = project.workdir / f"source{suffix}"
        await save_upload(source, src_path, MAX_DECK_BYTES, "Document")

        from starlette.concurrency import run_in_threadpool

        from .draft import DraftError

        try:
            markdown = await run_in_threadpool(
                _do_draft, src_path, n, provider, llm.get("model"),
                llm.get("base_url"), llm.get("api_key"), llm.get("think"),
            )
        except (DraftError, ValueError, ImportError) as e:
            raise HTTPException(status_code=400, detail=str(e))
        finally:
            src_path.unlink(missing_ok=True)

        (project.workdir / "deck.md").write_text(markdown, encoding="utf-8")
        # A freshly drafted deck supersedes any previously uploaded .pptx.
        (project.workdir / "deck.pptx").unlink(missing_ok=True)
        return JSONResponse(
            {"markdown": markdown, "state": _project_state(project)}
        )

    @app.put("/api/projects/{project_id}/deck")
    async def save_deck(
        project_id: str,
        request: Request,
        x_project_token: str | None = Header(default=None),
        _: None = Depends(require_token),
    ) -> JSONResponse:
        project = _project_or_401(project_id, x_project_token)
        body = await request.json()
        markdown = body.get("markdown", "")
        if not isinstance(markdown, str) or not markdown.strip():
            raise HTTPException(
                status_code=400, detail="markdown must be a non-empty string"
            )
        if len(markdown.encode("utf-8")) > MAX_DECK_BYTES:
            raise HTTPException(status_code=413, detail="Deck too large")

        from .parser import parse_markdown

        if not parse_markdown(markdown):
            raise HTTPException(
                status_code=400,
                detail="No slides found (need at least one '# ' heading).",
            )
        (project.workdir / "deck.md").write_text(markdown, encoding="utf-8")
        (project.workdir / "deck.pptx").unlink(missing_ok=True)
        return JSONResponse({"state": _project_state(project)})

    def _prepare_project_job(
        project: Project, options: dict[str, Any],
        voice_path: Path | None, photo_path: Path | None, include_images: bool,
    ) -> tuple[Job, Path, Path, Path | None, Path | None]:
        """Copy the project deck (and, for renders, its enriched images) into a
        fresh job workdir so the canonical project is never consumed. Returns the
        job plus the paths the runner needs."""
        deck = _project_deck(project)
        if deck is None:
            raise HTTPException(status_code=400, detail="Project has no deck yet")
        job_id = uuid.uuid4().hex
        workdir = jobs_root / job_id
        workdir.mkdir(parents=True)
        deck_copy = workdir / deck.name
        shutil.copyfile(deck, deck_copy)
        # Bring the enriched images alongside the deck so its ``images/slide_N``
        # references resolve during the render.
        if include_images:
            proj_images = project.workdir / "images"
            if proj_images.is_dir():
                dest = workdir / "images"
                dest.mkdir(exist_ok=True)
                for img in proj_images.iterdir():
                    if img.is_file():
                        shutil.copyfile(img, dest / img.name)
        # Voice/photo were uploaded into a temp area; move them under the job.
        moved_voice = moved_photo = None
        if voice_path is not None:
            moved_voice = workdir / voice_path.name
            shutil.move(str(voice_path), moved_voice)
        if photo_path is not None:
            moved_photo = workdir / photo_path.name
            shutil.move(str(photo_path), moved_photo)
        job_base = load_config() if local_mode else base_config
        job_yaml = _build_job_config(
            job_base, workdir, options, moved_voice, moved_photo
        )
        job = Job(
            id=job_id, workdir=workdir, created_at=time.time(),
            download_token=secrets.token_urlsafe(24),
        )
        with _LOCK:
            _JOBS[job_id] = job
        return job, deck_copy, job_yaml, moved_voice, moved_photo

    def _job_response(job: Job) -> JSONResponse:
        return JSONResponse(
            {"job_id": job.id, "status": job.status, "token": job.download_token}
        )

    @app.get("/api/projects/{project_id}/images/{name}")
    def project_image(
        project_id: str, name: str, t: str | None = None,
        x_project_token: str | None = Header(default=None),
        _: None = Depends(require_token),
    ):
        # An <img> tag can't set a header, so accept the project token via ?t=.
        project = _project_or_401(project_id, x_project_token or t)
        safe = Path(name).name  # strip any path components (traversal guard)
        img = project.workdir / "images" / safe
        if not img.is_file():
            raise HTTPException(status_code=404, detail="No such image")
        return FileResponse(str(img))

    @app.post("/api/projects/{project_id}/enrich")
    async def enrich_project(
        project_id: str,
        image_provider: str | None = Form(default=None),
        notes: str | None = Form(default=None),
        x_project_token: str | None = Header(default=None),
        _: None = Depends(require_token),
    ) -> JSONResponse:
        _reap_expired_jobs()
        project = _project_or_401(project_id, x_project_token)
        notes_mode = (notes or "").lower() if notes else None
        if notes_mode not in (None, "fill", "all"):
            raise HTTPException(status_code=400, detail="notes must be 'fill' or 'all'")
        options = {"image_provider": image_provider, "avatar": False}
        job, deck_copy, job_yaml, _v, _p = _prepare_project_job(
            project, options, None, None, include_images=False
        )
        executor.submit(
            _run_project_enrich, job, project, deck_copy, job_yaml, notes_mode
        )
        return _job_response(job)

    @app.post("/api/projects/{project_id}/render")
    async def render_project(
        project_id: str,
        voice: UploadFile | None = File(default=None),
        photo: UploadFile | None = File(default=None),
        narration_seconds: str | None = Form(default=None),
        image_provider: str | None = Form(default=None),
        avatar: str | None = Form(default=None),
        avatar_name: str | None = Form(default=None),
        avatar_slides: str | None = Form(default=None),
        reuse_avatar: str | None = Form(default=None),
        transition: str | None = Form(default=None),
        accent: str | None = Form(default=None),
        x_project_token: str | None = Header(default=None),
        _: None = Depends(require_token),
    ) -> JSONResponse:
        _reap_expired_jobs()
        project = _project_or_401(project_id, x_project_token)
        if _project_deck(project) is None:
            raise HTTPException(
                status_code=400, detail="Project has no deck to render yet"
            )
        staging = project.workdir / ("render_" + uuid.uuid4().hex)
        staging.mkdir()
        voice_path = photo_path = None
        try:
            if voice is not None and voice.filename:
                voice_path = staging / f"voice{Path(voice.filename).suffix or '.wav'}"
                await save_upload(voice, voice_path, MAX_VOICE_BYTES, "Voice sample")
            if photo is not None and photo.filename:
                photo_path = staging / f"photo{Path(photo.filename).suffix or '.png'}"
                await save_upload(photo, photo_path, MAX_PHOTO_BYTES, "Photo")
                from .providers.avatar import _source_kind

                if _source_kind(str(photo_path)) == "image":
                    problem = _validate_photo_upload(photo_path)
                    if problem:
                        raise HTTPException(status_code=400, detail=problem)
            options = {
                "narration_seconds": narration_seconds,
                "image_provider": image_provider,
                "avatar": (avatar or "true").lower() != "false",
                "avatar_name": avatar_name,
                "avatar_slides": avatar_slides,
                "reuse_avatar": (reuse_avatar or "").lower() == "true"
                if reuse_avatar is not None else None,
                "transition_seconds": transition,
                "accent": accent,
            }
            job, deck_copy, job_yaml, mv, mp = _prepare_project_job(
                project, options, voice_path, photo_path, include_images=True
            )
            executor.submit(_run_job, job, deck_copy, job_yaml, mv, mp, "video", None)
            return _job_response(job)
        finally:
            shutil.rmtree(staging, ignore_errors=True)

    return app


# Single-page UI. Remembers the voice sample + photo in the browser (IndexedDB)
# so the lecturer never re-picks them; the server stores neither at rest.
INDEX_HTML = """<!doctype html>
<html><head><meta charset="utf-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<title>SlideStream — slides in, narrated video out</title>
<link rel="icon" href="data:image/svg+xml,%3Csvg xmlns='http://www.w3.org/2000/svg' viewBox='0 0 100 100'%3E%3Ctext y='.9em' font-size='90'%3E%F0%9F%8E%AC%3C/text%3E%3C/svg%3E">
<link rel="preconnect" href="https://fonts.googleapis.com">
<link href="https://fonts.googleapis.com/css2?family=Fraunces:ital,wght@0,500;0,600;1,500&family=Inter:wght@400;500;600&display=swap" rel="stylesheet">
<style>
:root{--bg:#faf9f6;--ink:#1e1b18;--muted:#6b645c;--line:#e7e2d9;--accent:#c2410c;
 --accent-soft:#c2410c14;--card:#ffffff;--radius:14px}
@media (prefers-color-scheme: dark){
 :root{--bg:#17150f;--ink:#f2ede3;--muted:#a89f92;--line:#2c2820;--card:#201d16;
  --accent-soft:#c2410c26}}
*{box-sizing:border-box}
body{margin:0;font-family:Inter,system-ui,sans-serif;background:var(--bg);color:var(--ink);line-height:1.55}
.wrap{max-width:660px;margin:0 auto;padding:2.2rem 1.25rem 3rem}
h1{font-family:Fraunces,Georgia,serif;font-weight:600;font-size:1.9rem;margin:0;letter-spacing:-.01em}
h1 em{font-style:italic;color:var(--accent)}
.tag{color:var(--muted);margin:.3rem 0 1.6rem;font-size:.98rem}
.card{background:var(--card);border:1px solid var(--line);border-radius:var(--radius);
 padding:1.4rem 1.5rem;box-shadow:0 1px 2px rgba(0,0,0,.03)}
label{display:block;margin:1.05rem 0 .3rem;font-weight:600;font-size:.92rem}
label:first-child{margin-top:0}
input,select,textarea{font:inherit;color:inherit;width:100%;padding:.55rem .7rem;border:1px solid var(--line);
 border-radius:9px;background:var(--bg)}
input[type=file]{padding:.45rem .5rem;cursor:pointer}
input[type=file]::file-selector-button{font:inherit;font-weight:600;margin-right:.8rem;
 padding:.4rem .9rem;border:0;border-radius:7px;background:var(--accent-soft);color:var(--accent);cursor:pointer}
input:focus,select:focus,textarea:focus{outline:2px solid var(--accent);outline-offset:1px;border-color:transparent}
button{font:inherit;font-weight:600;width:100%;margin-top:1.4rem;padding:.75rem 1.2rem;
 border:0;border-radius:10px;background:var(--accent);color:#fff;cursor:pointer;transition:transform .06s}
button:active{transform:translateY(1px)}
button:disabled{opacity:.55;cursor:not-allowed}
.ghost{background:transparent;border:1px solid var(--line);color:var(--ink)}
.ghost:hover{border-color:var(--accent);color:var(--accent)}
.row{display:flex;gap:.55rem;align-items:center;margin-top:1rem}
.row input{width:auto;accent-color:var(--accent)} .row label{margin:0}
.muted{color:var(--muted);font-size:.85rem;margin:.35rem 0 0}
.banner{display:none;background:var(--accent-soft);border:1px solid var(--accent);
 border-radius:var(--radius);padding:.85rem 1.1rem;margin-bottom:1.2rem;font-size:.92rem}
.banner a{color:var(--accent);font-weight:600;text-decoration:none}
.banner code{font-family:ui-monospace,Menlo,monospace;font-size:.88em;
 background:rgba(120,120,120,.14);padding:.08em .4em;border-radius:5px}
.req{font-size:.72rem;font-weight:600;color:var(--accent);background:var(--accent-soft);
 padding:.1rem .5rem;border-radius:99px;vertical-align:middle}
details{margin-top:1.3rem;border-top:1px solid var(--line);padding-top:1rem}
details summary{cursor:pointer;font-weight:600;font-size:.95rem;color:var(--muted);list-style:none}
details summary::before{content:"▸ ";color:var(--accent)}
details[open] summary::before{content:"▾ "}
details summary:hover{color:var(--ink)}
/* Stepper */
.steps{display:flex;gap:.3rem;margin:0 0 1.3rem;flex-wrap:wrap}
.steps button{flex:1;min-width:110px;margin:0;padding:.5rem .4rem;background:transparent;
 border:1px solid var(--line);border-radius:9px;color:var(--muted);font-size:.85rem;font-weight:600;
 display:flex;gap:.4rem;align-items:center;justify-content:center}
.steps button .n{display:inline-flex;width:1.35rem;height:1.35rem;flex:none;align-items:center;
 justify-content:center;border-radius:99px;background:var(--accent-soft);color:var(--accent);font-size:.78rem}
.steps button.active{border-color:var(--accent);color:var(--ink)}
.steps button.done .n{background:var(--accent);color:#fff}
.step{display:none}.step.on{display:block}
.opt{border:1px solid var(--line);border-radius:11px;padding:1rem 1.1rem;margin-top:.9rem}
.opt:first-of-type{margin-top:0}
.opt h3{margin:0 0 .3rem;font-size:1rem}
.modes{display:grid;gap:.55rem}
.mode{display:flex;gap:.6rem;align-items:flex-start;border:1px solid var(--line);
 border-radius:11px;padding:.75rem .85rem;cursor:pointer;margin:0;font-weight:400}
.mode:hover{border-color:var(--accent)}
.mode input{width:auto;flex:none;margin-top:.25rem;accent-color:var(--accent)}
.mode strong{font-size:.95rem;display:block}
.mode small{color:var(--muted)}
.mode:has(input:checked){border-color:var(--accent);background:var(--accent-soft)}
.picker{display:flex;gap:.6rem;flex-wrap:wrap;margin:.2rem 0 .3rem}
.pick{flex:1;min-width:88px;border:1px solid var(--line);border-radius:10px;padding:.5rem .4rem;
 margin:0;font-weight:400;display:flex;flex-direction:column;align-items:center;gap:.3rem;
 cursor:pointer;text-align:center;font-size:.85rem}
.pick:hover{border-color:var(--accent)}
.pick img{width:64px;height:64px;object-fit:cover;border-radius:50%}
.pick input{accent-color:var(--accent)}
.pick:has(input:checked){border-color:var(--accent);background:var(--accent-soft)}
.err{display:none;margin:.2rem 0 .6rem;border:1px solid rgba(220,38,38,.45);
 background:rgba(220,38,38,.07);color:#dc2626;border-radius:9px;padding:.55rem .8rem;font-size:.88rem}
.dz{display:flex;align-items:center;justify-content:center;gap:.5rem;text-align:center;
 border:1.5px dashed var(--line);border-radius:9px;padding:.65rem .8rem;margin-top:.3rem;
 color:var(--muted);cursor:pointer;font-size:.88rem}
.dz:hover,.dz.drag{border-color:var(--accent);color:var(--accent)}
.dz input{display:none}
.dz.has{border-style:solid;color:var(--ink)}
.dz:focus-visible{outline:2px solid var(--accent);outline-offset:1px}
.elapsed{color:var(--muted);font-size:.85rem;font-weight:400;margin-left:.35rem}
.spinner{display:inline-block;width:13px;height:13px;border:2px solid var(--accent);
 border-right-color:transparent;border-radius:50%;animation:ssspin .8s linear infinite;
 vertical-align:-2px;margin-right:.45rem}
@keyframes ssspin{to{transform:rotate(360deg)}}
#progWrap{margin:.5rem 0 .2rem}
#progTrack{height:8px;border-radius:99px;background:var(--line);overflow:hidden}
#progBar{height:100%;width:0;border-radius:99px;background:var(--accent);transition:width .6s}
#progLabel{margin:.3rem 0 0}
#deckEditor{min-height:300px;font-family:ui-monospace,Menlo,monospace;font-size:.8rem;line-height:1.5;resize:vertical}
.sc{font-size:.85rem;color:var(--muted);margin-top:.5rem}
.navbtns{display:flex;gap:.6rem;margin-top:1.4rem}
.navbtns button{margin:0;flex:1}
#status{margin:1rem 0 .4rem;font-weight:500}
#status a{color:var(--accent);font-weight:600}
.badge{display:inline-block;padding:.12rem .65rem;border-radius:99px;font-size:.8rem;
 background:var(--accent-soft);color:var(--accent);font-weight:600}
#logWrap{margin-top:.8rem;border:1px solid var(--line);border-radius:9px;padding:.5rem .8rem .6rem;font-size:.85rem}
#logWrap summary{cursor:pointer;color:var(--muted);font-weight:600;font-size:.85rem}
#logWrap summary:hover{color:var(--accent)}
#trace{margin-top:.5rem}
#trace div{margin:.12rem 0;color:var(--ink)}
#log{white-space:pre-wrap;background:var(--bg);border:1px solid var(--line);
 padding:.6rem .7rem;border-radius:9px;font-family:ui-monospace,Menlo,monospace;
 font-size:.72rem;max-height:220px;overflow:auto;color:var(--muted);margin-top:.5rem}
#report{display:none;margin-top:1rem;border:1px solid var(--line);border-radius:9px;
 padding:.7rem 1rem .9rem;font-size:.88rem}
#report.on{display:block}
#report h4{margin:.7rem 0 .25rem;font-size:.72rem;text-transform:uppercase;
 letter-spacing:.07em;color:var(--accent)}
#report h4:first-child{margin-top:.2rem}
#report .f{margin:.16rem 0} #report .est{color:var(--muted)}
#report .sum{margin-top:.7rem;padding-top:.5rem;border-top:1px solid var(--line);font-weight:600}
.thumbs{display:grid;grid-template-columns:repeat(auto-fill,minmax(90px,1fr));gap:.4rem;margin-top:.9rem}
.thumbs img{width:100%;aspect-ratio:16/9;object-fit:cover;border-radius:6px;border:1px solid var(--line)}
footer{margin-top:2rem;padding-top:1.1rem;border-top:1px solid var(--line);
 font-size:.85rem;color:var(--muted);display:flex;gap:1.2rem;flex-wrap:wrap}
footer a{color:var(--muted);text-decoration:none;font-weight:500}
footer a:hover{color:var(--accent)}
.head{display:flex;align-items:baseline;justify-content:space-between}
#gear{display:none;font:inherit;width:auto;margin:0;padding:.4rem .8rem;
 background:transparent;border:1px solid var(--line);border-radius:9px;color:var(--muted)}
#gear:hover{color:var(--accent);border-color:var(--accent)}
#settings{display:none;position:fixed;inset:0;background:rgba(0,0,0,.35);z-index:20}
#settings.on{display:block}
.panel{position:absolute;top:0;right:0;bottom:0;width:min(560px,94vw);background:var(--card);
 border-left:1px solid var(--line);padding:1.3rem 1.4rem;overflow:auto;display:flex;flex-direction:column}
.panel h2{font-family:Fraunces,Georgia,serif;font-weight:600;margin:0 0 .2rem;font-size:1.3rem}
.panel textarea{flex:1;min-height:320px;font-family:ui-monospace,Menlo,monospace;font-size:.8rem;
 line-height:1.5;padding:.7rem;resize:vertical;margin:.8rem 0}
.panel .btns{display:flex;gap:.6rem}
.panel .btns button{margin:0;width:auto;padding:.55rem 1rem}
#saveMsg{font-size:.85rem;margin:.5rem 0 0}
</style></head><body><div class="wrap">
<div class="head"><h1>🎬 Slide<em>Stream</em></h1>
<button id="gear" title="Providers &amp; settings">⚙ Settings</button></div>
<p class="tag">A document or a deck in — a narrated video out. In your voice, or a friendly mascot's.</p>

<div id="settings"><div class="panel">
<h2>Settings</h2>
<p class="muted" id="setPath">Edits your ~/.slidestream.yaml — providers, servers, keys.</p>
<textarea id="setYaml" spellcheck="false" placeholder="# empty — click 'Insert template' to start"></textarea>
<div class="btns">
 <button id="setSave">Save</button>
 <button id="setTpl" class="ghost">Insert template</button>
 <button id="setClose" class="ghost">Close</button>
</div>
<p id="saveMsg"></p>
</div></div>

<div id="demo" class="banner">
 <strong>Want more?</strong> The free <strong>desktop app</strong> unlocks
 <strong>your own cloned voice</strong>, <strong>your own headshot or video as the presenter</strong>,
 unlimited slides, drafting from documents, your own AI providers — with full privacy.
 <a id="dlBanner" href="https://github.com/michael-borck/slide-stream/releases/latest">⬇ Get the desktop app</a> &middot;
 <code>pip install slide-stream</code> &middot;
 <a href="https://slidestream.eduserver.au">learn more</a>
 <span class="muted" id="limits" style="display:block"></span>
 <span class="muted" style="display:block">Everything you see — images, voice, animated presenter — is generated on one
 modest self-hosted machine. Slower than a cloud render farm, and proof of
 what consumer-grade hardware can do.</span>
</div>

<div class="steps demo-hide" id="steps">
 <button data-step="deck" class="active"><span class="n">1</span>Deck</button>
 <button data-step="configure" disabled><span class="n">2</span>Configure</button>
 <button data-step="render" disabled><span class="n">3</span>Render</button>
</div>

<div class="card">
<div id="tokrow" style="display:none"><label for="token">Access token</label>
 <input id="token" type="password" placeholder="paste your token" autocomplete="off">
 <p class="muted">Stored in this browser only.</p></div>

<!-- STEP 1: DECK (source + review/edit in one) -->
<div class="step on" id="step-deck">
 <p class="err" id="err-deck" role="alert"></p>
 <div class="modes" id="inputPicker" style="display:none">
  <label class="mode"><input type="radio" name="inmode" value="idea" id="inmodeIdea">
   <span><strong>✍️ Start from an idea</strong>
   <small>Type a topic — AI drafts the slides</small></span></label>
  <label class="mode" id="inmodeDoc"><input type="radio" name="inmode" value="doc">
   <span><strong>📝 Draft from a document</strong>
   <small>A PDF, Word doc or PPT becomes the slides</small></span></label>
  <label class="mode"><input type="radio" name="inmode" value="upload">
   <span><strong>📄 Use an existing deck</strong>
   <small>.md, .pptx, .txt or .qmd — review and tweak below</small></span></label>
 </div>
 <div class="opt" id="ideaOpt" style="display:none">
  <label for="ideaText">Your topic <span style="font-weight:400;color:var(--muted)">(an idea, a title, a rough outline)</span></label>
  <textarea id="ideaText" rows="3" placeholder="e.g. Why sleep matters for learning — for first-year students"></textarea>
  <div id="ideaSlidesWrap" class="demo-hide">
   <label for="ideaSlides">Number of slides <span style="font-weight:400;color:var(--muted)">(blank = let the AI decide)</span></label>
   <input id="ideaSlides" type="number" min="1" max="20" placeholder="e.g. 8">
  </div>
  <button id="ideaGo">Draft the deck</button>
  <p class="muted">AI-generated — review the facts before presenting.</p>
 </div>
 <div class="opt" id="draftOpt" style="display:none">
  <label for="draftFile">Document <span style="font-weight:400;color:var(--muted)">(.pdf, .docx, .pptx, .txt, .md)</span></label>
  <input id="draftFile" type="file" accept=".pdf,.docx,.pptx,.txt,.md">
  <label for="draftSlides">Number of slides <span style="font-weight:400;color:var(--muted)">(blank = let the AI decide)</span></label>
  <input id="draftSlides" type="number" min="1" placeholder="e.g. 10">
  <button id="draftGo">Generate deck</button>
  <p class="muted">AI-generated — review the facts before presenting.</p>
 </div>
 <div class="opt" id="deckOpt" style="display:none">
  <label for="deckFile">Deck file <span class="req">.md, .pptx, .txt or .qmd</span></label>
  <input id="deckFile" type="file" accept=".md,.pptx,.txt,.qmd">
 </div>
 <p class="sc" id="draftStatus" style="display:none"></p>
 <div id="deckEditWrap" style="display:none">
  <label for="deckEditor">Review your deck <span style="font-weight:400;color:var(--muted)">(Markdown — one '# ' heading per slide; narration and images are built from this)</span></label>
  <textarea id="deckEditor" spellcheck="false"></textarea>
  <p class="sc" id="deckCount"></p>
  <p class="sc" id="trimHint" style="display:none;color:var(--accent)"></p>
 </div>
 <div id="deckPptxNote" style="display:none">
  <p>🖼️ <strong>PowerPoint deck loaded.</strong> Slides are read straight from your file — editing is disabled. To change the content, edit the deck in PowerPoint and re-upload.</p>
 </div>
 <div id="teaserConfig" style="display:none"></div>
 <div class="navbtns" id="navDeck">
  <button id="deckNext" disabled>Save &amp; continue →</button>
 </div>
</div>

<!-- STEP 3: CONFIGURE -->
<div class="step" id="step-configure">
 <label>Make</label>
 <div class="modes" id="modesDiv">
  <label class="mode"><input type="radio" name="outmode" value="video_plain" checked>
   <span><strong>🎬 Narrated video — slides as they are</strong>
   <small>No AI images; your deck, narrated.</small></span></label>
  <label class="mode"><input type="radio" name="outmode" value="video_rich">
   <span><strong>🖼️ Narrated video — enhanced slides</strong>
   <small>Every slide gets AI artwork, then narration on top.</small></span></label>
  <label class="mode"><input type="radio" name="outmode" value="deck">
   <span><strong>📄 Enhanced slides only (.pptx)</strong>
   <small>Images + optional AI speaker notes added to your deck. No narration.</small></span></label>
 </div>
 <div id="imgWrap" class="demo-hide">
  <label for="imageProvider">Slide images</label>
  <select id="imageProvider">
   <option value="">Use my configured provider</option>
   <option value="text">Text cards (no AI, always works)</option>
   <option value="dalle3">DALL·E 3 (OpenAI)</option>
   <option value="gemini">Imagen (Gemini)</option>
   <option value="pexels">Pexels stock photos</option>
   <option value="unsplash">Unsplash stock photos</option>
  </select>
 </div>
 <div id="deckOnlyNote" style="display:none">
  <p class="muted">Narration and presenter options don't apply to a .pptx export.</p>
 </div>
 <div id="notesWrap" style="display:none" class="demo-hide">
  <label for="notes">AI speaker notes</label>
  <select id="notes">
   <option value="">None</option>
   <option value="all">Write for every slide</option>
   <option value="fill">Only where notes are missing</option>
  </select>
 </div>
 <div id="videoOnly">
  <div id="secsWrap" class="demo-hide">
   <label for="secs">Seconds of narration per slide</label>
   <input id="secs" type="number" min="10" placeholder="e.g. 30">
  </div>
  <div id="transitionWrap" class="demo-hide">
   <label for="transition">Slide transition <span style="font-weight:400;color:var(--muted)">(crossfade seconds; 0 = hard cut)</span></label>
   <input id="transition" type="number" min="0" step="0.1" placeholder="0">
  </div>
  <details id="extras" open>
  <summary id="extrasSummary">Voice &amp; presenter <span style="font-weight:400">(optional)</span></summary>
  <div id="voiceWrap" class="demo-hide">
   <label for="voice">Your voice <span style="font-weight:400;color:var(--muted)">(a 10–30s sample clones it for this render only — mp3, wav, m4a, or even a video; we extract the audio)</span></label>
   <input id="voice" type="file" accept="audio/*,video/*">
  </div>
  <div id="mascotSelectWrap" class="demo-hide">
   <label for="avatarName">Mascot presenter</label>
   <select id="avatarName"><option value="">None</option></select>
  </div>
  <p class="muted">A friendly character presents in the corner<span class="demo-hide"> — or upload a photo or short video of yourself below (desktop app)</span>.</p>
  <details id="teaserPickers" style="display:none" open>
   <summary>Presenter &amp; voice</summary>
  </details>
  <div id="avatarSlidesWrap" class="demo-hide">
   <label for="avatarSlides">Show presenter on</label>
   <select id="avatarSlides">
    <option value="">Every slide</option>
    <option value="first,last">First &amp; last slide</option>
    <option value="first">First slide only</option>
    <option value="every:3">Every 3rd slide</option>
    <option value="none">No slides</option>
   </select>
   <p class="muted">Narration always plays; the talking head appears only on these slides — handy when a slide's corner is busy, or to spare a slow GPU.</p>
  </div>
  <div id="photoWrap" class="demo-hide">
   <label for="photo">Your photo or short video <span style="font-weight:400;color:var(--muted)">(front-facing)</span></label>
   <input id="photo" type="file" accept="image/*,video/*">
   <p class="muted" id="remembered"></p>
  </div>
  <div id="animateWrap" class="demo-hide">
   <div class="row"><input id="avatar" type="checkbox" checked><label for="avatar">Animate the presenter</label></div>
   <p class="muted">On: your presenter talks with AI lip-sync — mascots and photos alike (driven by the narration audio). Off: the presenter appears as a still image in the corner.</p>
   <p class="muted">More animation engines (D-ID, SadTalker, Wav2Lip) are available in the desktop app.</p>
  </div>
  <label for="accent" id="accentRow" style="display:none">Accent</label>
  <select id="accent" aria-label="Voice accent" style="display:none"><option value="">— default —</option></select>
  </details>
 </div>
 <div class="navbtns">
  <button class="ghost" data-goto="deck">← Back</button>
  <button data-goto="render">Continue →</button>
 </div>
</div>

<!-- STEP 4: RENDER -->
<div class="step" id="step-render">
 <p class="muted" id="renderHint">Preview the plan, then render. The video uses the images your provider generates.</p>
 <button id="check" class="ghost demo-hide">Check deck first</button>
 <button id="go">Create video</button>
 <p id="status"></p><p id="notice" class="muted"></p>
 <div id="progWrap" style="display:none">
  <div id="progTrack"><div id="progBar"></div></div>
  <p class="muted" id="progLabel"></p>
 </div>
 <div id="doneAdvert" class="banner" style="display:none">
  <strong>Like it?</strong> That whole video — images, voice, animated
  presenter — was generated on the same modest self-hosted machine. The free
  <strong>desktop app</strong> unlocks your own cloned voice, your own
  headshot, unlimited slides — and it is faster on your own hardware.
  <a id="doneDl" href="https://github.com/michael-borck/slide-stream/releases/latest">⬇ Get the desktop app</a>
  <span class="muted" id="linkExpiry" style="display:block"></span>
 </div>
 <details id="logWrap" style="display:none">
  <summary>Render log &amp; trace</summary>
  <div id="trace"></div>
  <div id="log"></div>
 </details>
 <div id="report"></div><div id="thumbs" class="thumbs"></div>
 <div class="navbtns"><button class="ghost" data-goto="configure">← Back</button></div>
</div>
</div>

<footer>
 <a href="https://slidestream.eduserver.au">About</a>
 <a href="https://github.com/michael-borck/slide-stream/blob/main/docs/USER_GUIDE.md">Docs</a>
 <a href="https://pypi.org/project/slide-stream/">pip install slide-stream</a>
 <a href="https://github.com/michael-borck/slide-stream">GitHub</a>
 <span class="muted" id="appVersion"></span>
</footer>
</div>
<script>
const $=id=>document.getElementById(id);
let cfg={},demo=false,canDraft=false;
let projectId=null,projectToken=null;
// The working deck, represented client-side so Check/Render/Export can reuse it:
// {name, file (File|Blob|null), markdown (string|null), isPptx}
let deck=null;

// Platform-detected desktop download in the demo banner (stable asset names).
(()=>{const ua=(navigator.userAgent||"").toLowerCase();
 const f=ua.includes("mac")?"SlideStream-macos-apple-silicon.dmg":
   ua.includes("win")?"SlideStream-windows-setup.exe":
   ua.includes("linux")?"SlideStream-linux.AppImage":null;
 if(f)$("dlBanner").href=
  "https://github.com/michael-borck/slide-stream/releases/latest/download/"+f})();

$("token").value=localStorage.getItem("ss_token")||"";
$("token").oninput=e=>localStorage.setItem("ss_token",e.target.value);
const auth=()=>{const h={};const t=$("token").value;if(t)h.Authorization="Bearer "+t;return h};
const pauth=()=>{const h=auth();if(projectToken)h["X-Project-Token"]=projectToken;return h};

// --- Inline errors (no alert() dialogs) --------------------------------------
function showErr(step,msg){const el=$("err-"+step);if(!el)return;
 el.textContent=msg;el.classList.add("on")}
function clearErrs(){document.querySelectorAll(".err").forEach(e=>e.classList.remove("on"))}

// --- Remembered choices (localStorage; voice/photo live in IndexedDB) --------
const REMEMBER=["imageProvider","secs","transition","accent","avatarName","avatarSlides","notes"];
const sget=k=>localStorage.getItem("ss_"+k);
const sset=(k,v)=>{if(v)localStorage.setItem("ss_"+k,v);else localStorage.removeItem("ss_"+k)};
function saveSettings(){REMEMBER.forEach(k=>{const el=$(k);if(el)sset(k,el.value)});
 const r=document.querySelector('input[name="outmode"]:checked');if(r)sset("outmode",r.value)}
function loadSettings(){if(demo)return; // the teaser is stateless by design
 REMEMBER.forEach(k=>{const v=sget(k),el=$(k);
  // Only apply when the option exists (accent/mascot options load async).
  if(v&&el&&(el.tagName==="SELECT"&&el.querySelector('option[value="'+v+'"]')||el.tagName==="INPUT"))el.value=v});
 const m=sget("outmode");
 if(m){const r=document.querySelector('input[name="outmode"][value="'+m+'"]');
  if(r&&!r.checked)r.checked=true}}
loadSettings();
document.addEventListener("change",saveSettings);

// --- Drag-and-drop upload zones ----------------------------------------------
// Wraps each file input in a clickable drop target that shows the file name;
// the input itself stays hidden but is what handlers keep reading.
function dropzone(id,hint){const inp=$(id);if(!inp)return;
 const box=document.createElement("div");box.className="dz";
 box.tabIndex=0;box.setAttribute("role","button");box.setAttribute("aria-label",hint);
 const label=document.createElement("span");label.textContent=hint;
 inp.parentNode.insertBefore(box,inp);box.append(inp,label);
 const set=f=>{label.textContent=f?("📄 "+f.name):hint;
  box.classList.toggle("has",!!f)};
 inp.addEventListener("change",()=>set(inp.files[0]));
 box.addEventListener("click",()=>inp.click());
 box.addEventListener("keydown",e=>{if(e.key==="Enter"||e.key===" "){e.preventDefault();inp.click()}});
 ["dragenter","dragover"].forEach(ev=>box.addEventListener(ev,e=>{
  e.preventDefault();box.classList.add("drag")}));
 ["dragleave","dragend"].forEach(ev=>box.addEventListener(ev,()=>box.classList.remove("drag")));
 box.addEventListener("drop",e=>{e.preventDefault();box.classList.remove("drag");
  const dt=e.dataTransfer;if(!dt||!dt.files.length)return;
  try{const transfer=new DataTransfer();[...dt.files].forEach(f=>transfer.items.add(f));
   inp.files=transfer.files}
  catch(_){/* older browsers: drop only highlights; click still works */}
  set(inp.files[0]||dt.files[0])});
}

fetch("/api/config").then(r=>r.json()).then(c=>{
 cfg=c;demo=!!c.demo;
 if(c.version)$("appVersion").textContent="SlideStream v"+c.version;
 if(c.auth_required)$("tokrow").style.display="block";
  if(c.demo){$("demo").style.display="block";
   if(c.limits)$("limits").textContent=
    "Hosted demo: first "+c.limits.max_slides+" slides per video, "+
    c.limits.jobs_per_hour+" videos per hour, nothing stored.";}
 if(c.local)$("gear").style.display="inline-block";
 // Draft needs an LLM and the (non-demo) project workflow.
 canDraft=!demo&&!!c.llm;
 setupInputModes();
  (c.avatars||[]).forEach(a=>{const o=document.createElement("option");o.value=a;o.textContent=a;$("avatarName").appendChild(o)});
  if((c.accents||[]).length){$("accentRow").style.display="block";$("accent").style.display="block";
   c.accents.forEach(a=>{const o=document.createElement("option");o.value=a;o.textContent=a;$("accent").appendChild(o)})}
  if(demo){ // hosted teaser: two happy paths, no personal voice or photos
   document.querySelectorAll(".demo-hide").forEach(e=>e.style.display="none");
   const radio=v=>document.querySelector('input[name="outmode"][value="'+v+'"]');
   radio("video_plain").closest(".mode").style.display="none";
   // Enhance first, mascot second — clicking the mascot card reveals the
   // presenter accordion right below it.
   radio("deck").checked=true;
   const modes=$("modesDiv");
   modes.append(radio("deck").closest(".mode"), radio("video_rich").closest(".mode"));
   radio("deck").closest(".mode").querySelector("span").innerHTML=
    "<strong>🖼️ Enhance the slides (.pptx)</strong><small>AI images + freshly written speaker notes. No narration.</small>";
   radio("video_rich").closest(".mode").querySelector("span").innerHTML=
    "<strong>🎬 Talking-mascot video</strong><small>Narrated by a character you pick below, with AI images.</small>";
   $("teaserConfig").style.display="";
   $("teaserConfig").append($("modesDiv"),$("teaserPickers"),$("navDeck"));
   $("extras").style.display="none"; // the old presenter wrapper is redundant here
   $("deckNext").textContent="Generate ⚡";
   // The old configure page is a leftover in the teaser — the render step's
   // Back returns to the deck page instead.
   const back=document.querySelector('#step-render [data-goto="configure"]');
   if(back){back.dataset.goto="deck";back.textContent="← Edit deck / start over"}
   if(cfg.avatars&&cfg.avatars.length)$("avatarName").value=cfg.avatars[0];
   const pickers=$("teaserPickers");
   if(cfg.avatars&&cfg.avatars.length){
    const mLab=document.createElement("label");mLab.textContent="Pick your presenter";
    const grid=document.createElement("div");grid.className="picker";
    cfg.avatars.forEach((a,i)=>{const lab=document.createElement("label");lab.className="pick";
     lab.innerHTML='<input type="radio" name="teaserAvatar" value="'+a+'"'+(i===0?" checked":"")+
      '><img src="/api/avatars/'+encodeURIComponent(a)+'" alt=""><span>'+a+"</span>";
     lab.querySelector("input").onchange=()=>{$("avatarName").value=a};
     grid.appendChild(lab)});
    pickers.append(mLab,grid);
   }
   if((cfg.voices||[]).length){
    const vLab=document.createElement("label");vLab.textContent="Pick a voice";
    const grid=document.createElement("div");grid.className="picker";
    cfg.voices.forEach((v,i)=>{const lab=document.createElement("label");lab.className="pick";
     lab.innerHTML='<input type="radio" name="teaserVoice" value="'+v+'"'+(i===0?" checked":"")+
      '><span>🎙️ '+v+"</span>";
     grid.appendChild(lab)});
    pickers.append(vLab,grid);
   }  }
  loadSettings();updateMode();
 }).catch(()=>{setInputMode("upload");}); // config failed: still let people upload

// --- Step navigation --------------------------------------------------------
const ORDER=["deck","configure","render"];
let reached={deck:true};
function go(step){
 ORDER.forEach(s=>{$("step-"+s).classList.toggle("on",s===step)});
 clearErrs();
 document.querySelectorAll("#steps button").forEach(b=>{
  const s=b.dataset.step;b.classList.toggle("active",s===step);
  b.classList.toggle("done",reached[s]&&s!==step);
  b.disabled=!reached[s]});
 window.scrollTo({top:0,behavior:"smooth"});
}
function reach(step){reached[step]=true}
document.querySelectorAll("#steps button").forEach(b=>{
 b.onclick=()=>{if(reached[b.dataset.step])go(b.dataset.step)}});
document.querySelectorAll("[data-goto]").forEach(b=>{
 b.onclick=()=>{const t=b.dataset.goto;reach(t);go(t)}});

// --- IndexedDB: remember voice + photo across renders (client-side only) -----
let db;const openDB=()=>new Promise(r=>{const q=indexedDB.open("ss",1);
 q.onupgradeneeded=()=>q.result.createObjectStore("files");q.onsuccess=()=>{db=q.result;r()}});
const put=(k,v)=>new Promise(r=>{db.transaction("files","readwrite").objectStore("files").put(v,k).onsuccess=r});
const get=k=>new Promise(r=>{const q=db.transaction("files").objectStore("files").get(k);q.onsuccess=()=>r(q.result)});
let savedVoice,savedPhoto;
openDB().then(async()=>{savedVoice=await get("voice");savedPhoto=await get("photo");
 const b=[];if(savedVoice)b.push("voice: "+savedVoice.name);if(savedPhoto)b.push("photo: "+savedPhoto.name);
 $("remembered").textContent=b.length?("Remembered "+b.join(", ")+" — leave the fields empty to reuse."):""});
async function fileOrSaved(input,key,saved){const f=input.files[0];
 if(f){await put(key,f);return f}return saved||null}

// --- Deck helpers -----------------------------------------------------------
function countSlides(md){return (md.match(/^#\\s+\\S/gm)||[]).length}
function refreshDeckCount(){const n=countSlides($("deckEditor").value);
 $("deckCount").textContent=n?(n+" slide"+(n===1?"":"s")):"No slides yet — add a '# ' heading per slide.";
 const cap=demo&&cfg.limits?cfg.limits.max_slides:0;
 if(demo&&cap&&n>cap){$("trimHint").style.display="";
  $("trimHint").textContent=
  "The demo uses the first "+cap+" slides — get the desktop app for all "+n+"."}
 else $("trimHint").style.display="none"}
$("deckEditor").oninput=refreshDeckCount;

// --- Output mode (three ways to finish) --------------------------------------
function outMode(){const el=document.querySelector('input[name="outmode"]:checked');
 return el?el.value:"video_plain"}
function updateMode(){const m=outMode(),deck=m==="deck",rich=m==="video_rich";
 $("imgWrap").style.display=(deck||rich)?"":"none";
 // Teaser grouping: mascot + voice belong to the narration path only.
 if(demo)$("teaserPickers").style.display=rich?"":"none";
 if(deck){ // only the image provider + notes apply to a .pptx export
   $("deckOnlyNote").style.display="";
   $("notesWrap").style.display=demo?"none":"";
   $("videoOnly").style.display="none";
  $("renderHint").textContent="Preview the plan, then export your enhanced deck.";
 }else{
  $("deckOnlyNote").style.display="none";$("notesWrap").style.display="none";
  $("videoOnly").style.display="";
  $("renderHint").textContent=rich?
   "Preview the plan, then render. Each slide gets AI artwork from your chosen provider.":
   "Preview the plan, then render. Your slides stay as they are — no AI images.";
 }
 $("go").textContent=deck?"Export enhanced deck (.zip)":"Create video"}
document.querySelectorAll('input[name="outmode"]').forEach(r=>{r.onchange=updateMode});
updateMode();
// The bytes to send to the stateless /api/check (and, in demo, /api/jobs).
function deckBlob(){
 if(deck&&deck.isPptx&&deck.file)return deck.file;
 return new Blob([$("deckEditor").value],{type:"text/markdown"});
}
function deckName(){return deck&&deck.isPptx?(deck.name||"deck.pptx"):"deck.md"}

async function ensureProject(){
 if(demo||projectId)return;
 const r=await fetch("/api/projects",{method:"POST",headers:auth()});
 if(!r.ok)throw new Error(await r.text());
 const j=await r.json();projectId=j.project_id;projectToken=j.token;
}

// --- Deck step: load a deck, review it, continue -----------------------------
// --- Deck input modes (exclusive: idea / document / file) ---------------------
// One chooser, one visible input set — an idea never "modifies" an upload.
function setInputMode(m){
 $("ideaOpt").style.display=m==="idea"?"":"none";
 $("draftOpt").style.display=m==="doc"?"":"none";
 $("deckOpt").style.display=m==="upload"?"":"none";
}
function setupInputModes(){
 const picker=$("inputPicker");
 const hasLLM=!!cfg.llm, canDoc=canDraft; // doc drafting needs the project workflow
 $("inmodeDoc").style.display=canDoc?"":"none";
 document.getElementById("inmodeIdea").closest(".mode").style.display=hasLLM?"":"none";
 if(!hasLLM&&!canDoc){ // only one path exists: skip the chooser entirely
  picker.style.display="none";setInputMode("upload");return}
 picker.style.display="";
 const def=hasLLM?"idea":(canDoc?"doc":"upload");
 const el=document.querySelector('input[name="inmode"][value="'+def+'"]');
 if(el)el.checked=true;
 setInputMode(def);
}
document.querySelectorAll('input[name="inmode"]').forEach(r=>{
 r.onchange=()=>setInputMode(r.value)});

function revealEditor({pptx=false,markdown="",slideCount=null}={}){
 const ed=$("deckEditor");
 if(pptx){
  ed.disabled=true;ed.value="";
  ed.placeholder="PowerPoint content can't be edited here — slides are read straight from your file. Edit the deck in PowerPoint and re-upload to change it.";
  $("deckPptxNote").style.display="block";
  $("deckCount").textContent=slideCount?
   slideCount+" slide"+(slideCount===1?"":"s")+" (read from the .pptx)":
   "Slides are read from the .pptx file.";
  $("trimHint").style.display="none";
  $("deckNext").textContent="Continue →";
 }else{
  ed.disabled=false;ed.placeholder="";
  ed.value=markdown;
  $("deckPptxNote").style.display="none";
  $("deckNext").textContent="Save & continue →";
  refreshDeckCount();
 }
 $("deckEditWrap").style.display="";
 $("deckNext").disabled=false;
 $("deckEditWrap").scrollIntoView({behavior:"smooth",block:"nearest"});
}

$("draftGo").onclick=async()=>{
 const f=$("draftFile").files[0];
 if(!f){showErr("deck","Pick a document first (or drop one onto the box above).");return}
 $("draftGo").disabled=true;$("draftGo").textContent="Drafting…";
 draftBegin();
 try{
  await ensureProject();
  const fd=new FormData();fd.append("source",f);
  if($("draftSlides").value)fd.append("slides",$("draftSlides").value);
  const r=await fetch("/api/projects/"+projectId+"/draft",{method:"POST",headers:pauth(),body:fd});
  if(!r.ok){const j=await r.json().catch(()=>({}));
   draftEnd(false);showErr("deck","Draft failed: "+(j.detail||r.statusText||"error"));return}
  const j=await r.json();
  deck={name:"deck.md",file:null,markdown:j.markdown,isPptx:false};
  draftEnd(true);
  revealEditor({markdown:j.markdown});
 }catch(e){draftEnd(false);showErr("deck","Draft failed: "+e.message)}
 finally{$("draftGo").disabled=false;$("draftGo").textContent="Generate deck"}
};

// AI drafting: the request is async (LLM runs server-side, typically 5–20s),
// so surface an obvious waiting state where the deck will land — including on
// the Generate button, which stays disabled until the draft exists.
function draftBegin(){
 $("draftStatus").style.display="";
 $("draftStatus").innerHTML=
  '<span class="spinner"></span>Drafting your deck — this can take a little while…';
 $("deckNext").disabled=true;
 $("deckNext").textContent="Waiting for the draft…";
}
function draftEnd(ok){
 $("draftStatus").style.display="none";
 $("deckNext").disabled=!ok;
 $("deckNext").textContent=demo?"Generate ⚡":"Save & continue →";
}

$("ideaGo").onclick=async()=>{
 const topic=$("ideaText").value.trim();
 if(topic.length<3){showErr("deck","Type a topic or a few sentences first.");return}
 $("ideaGo").disabled=true;$("ideaGo").textContent="Drafting…";
 draftBegin();
 try{
  const body={topic:topic};
  const n=demo?5:parseInt($("ideaSlides").value,10);
  if(n)body.slides=n;
  const r=await fetch("/api/draft-deck",{
   method:"POST",
   headers:Object.assign({"Content-Type":"application/json"},auth()),
   body:JSON.stringify(body)});
  if(!r.ok){const j=await r.json().catch(()=>({}));
   draftEnd(false);showErr("deck","Draft failed: "+(j.detail||"error"));return}
  const j=await r.json();
  deck={name:"deck.md",file:null,markdown:j.markdown,isPptx:false};
  draftEnd(true);
  revealEditor({markdown:j.markdown});
 }catch(e){draftEnd(false);showErr("deck","Draft failed: "+e.message)}
 finally{$("ideaGo").disabled=false;$("ideaGo").textContent="Draft the deck"}
};

// Loading starts as soon as a file is picked — no extra button.
$("deckFile").onchange=async()=>{
 const f=$("deckFile").files[0];
 if(!f)return;
 const isPptx=/\\.pptx$/i.test(f.name);
 try{
  let slideCount=null;
  if(!demo){
   // Create a project seeded with the uploaded deck.
   projectId=null;projectToken=null;
   const fd=new FormData();fd.append("deck",f);
   const r=await fetch("/api/projects",{method:"POST",headers:auth(),body:fd});
   if(!r.ok){showErr("deck","Upload failed: "+(await r.text()));return}
   const j=await r.json();projectId=j.project_id;projectToken=j.token;
   slideCount=j.state&&j.state.slide_count;
  }
  deck={name:f.name,file:f,markdown:null,isPptx};
  if(isPptx){
   revealEditor({pptx:true,slideCount});
  }else{
   deck.markdown=await f.text();
   if(!countSlides(deck.markdown)){
    deck=null;
    showErr("deck","No slides found — start each slide with a '# ' heading (or separate slides with '---' lines), then upload again.");
    return}
   revealEditor({markdown:deck.markdown});
  }
 }catch(e){showErr("deck","Upload failed: "+e.message)}
};

// Save edits and continue (desktop) — or jump straight to Generate (teaser).
$("deckNext").onclick=async()=>{
 if(!deck){showErr("deck","Upload a deck — or draft one from a document — first.");return}
 if(demo){reach("render");go("render");$("go").click();return}
 if(deck.isPptx){reach("configure");go("configure");return}
 const md=$("deckEditor").value;
 if(!countSlides(md)){$("deckCount").textContent="Add at least one '# ' heading first.";return}
 deck.markdown=md;
 if(!demo){
  $("deckNext").disabled=true;
  try{
   const r=await fetch("/api/projects/"+projectId+"/deck",{method:"PUT",
    headers:Object.assign({"Content-Type":"application/json"},pauth()),
    body:JSON.stringify({markdown:md})});
   if(!r.ok){const j=await r.json().catch(()=>({}));
    showErr("deck","Save failed: "+(j.detail||"error"));return}
  }finally{$("deckNext").disabled=false}
 }
 reach("configure");go("configure");
};

// --- Configure / Render form fields -----------------------------------------
async function renderOptions(fd){
 // Narration/presenter options only apply to video renders, not .pptx export.
 if(outMode()==="deck")return fd;
 if(demo){
  // Teaser paths carry no PII: stock server voice (picked by name), mascot
  // presenter only.
  fd.append("avatar","true");
  if($("avatarName").value)fd.append("avatar_name",$("avatarName").value);
  const voice=document.querySelector('input[name="teaserVoice"]:checked');
  if(voice)fd.append("voice_name",voice.value);
  return fd;
 }
 const voice=await fileOrSaved($("voice"),"voice",savedVoice);
 const photo=await fileOrSaved($("photo"),"photo",savedPhoto);
 if(voice)fd.append("voice",voice);if(photo)fd.append("photo",photo);
 fd.append("avatar",$("avatar").checked?"true":"false");
 if($("avatarName").value)fd.append("avatar_name",$("avatarName").value);
 if($("avatarSlides").value)fd.append("avatar_slides",$("avatarSlides").value);
 if($("accent").value)fd.append("accent",$("accent").value);
 if($("secs").value)fd.append("narration_seconds",$("secs").value);
 if($("transition").value)fd.append("transition",$("transition").value);
 return fd;
}

// --- Doctor check -----------------------------------------------------------
const ICON={ok:"✅",warn:"⚠️",blocker:"❌"};
const esc=s=>{const d=document.createElement("div");d.textContent=s;return d.innerHTML};
function renderReport(rep){
 const groups={};rep.findings.forEach(f=>{(groups[f.group]=groups[f.group]||[]).push(f)});
 let h="";Object.keys(groups).forEach(g=>{h+="<h4>"+esc(g)+"</h4>";
  groups[g].forEach(f=>{h+='<div class="f">'+(ICON[f.severity]||"")+" "+esc(f.message)+"</div>"})});
 if((rep.estimates||[]).length){h+="<h4>Estimates</h4>";
  rep.estimates.forEach(e=>{h+='<div class="f est">• '+esc(e)+"</div>"})}
 h+='<div class="sum">'+(rep.blockers?("❌ "+rep.blockers+" blocker(s) · "):"")+
  (rep.warnings?("⚠️ "+rep.warnings+" warning(s)"):"✅ no warnings")+"</div>";
 $("report").innerHTML=h;$("report").classList.add("on")}
$("check").onclick=async()=>{
 $("check").disabled=true;$("status").textContent="Checking…";$("report").classList.remove("on");
 try{const fd=await renderOptions(new FormData());
  fd.append("output",outMode()==="deck"?"pptx":"video");
  fd.append("deck",deckBlob(),deckName());
  const r=await fetch("/api/check",{method:"POST",headers:auth(),body:fd});
  if(!r.ok){$("status").textContent="Error: "+(await r.text());return}
  $("status").textContent="";renderReport(await r.json())}
 finally{$("check").disabled=false}};

// --- Render / export (one submit, three output modes) ------------------------
function startOver(){ // back to step 1, deck kept, last job's panel cleared
 finishedJob=false;$("status").textContent="";$("notice").textContent="";
 hideProg();$("log").textContent="";$("logWrap").style.display="none";
 $("report").classList.remove("on");$("thumbs").innerHTML="";
 reach("deck");go("deck");updateMode()}

$("go").onclick=async()=>{
 if(finishedJob){startOver();return}
 const m=outMode(),deckOut=m==="deck";
 $("go").disabled=true; // no double submits while a job is in flight
 $("status").textContent="Uploading…";$("log").textContent="";$("report").classList.remove("on");
 $("notice").textContent="";hideProg();
 try{
  const fd=await renderOptions(new FormData());
  // "Slides as they are" must stay image-free regardless of the server's
  // configured provider; the other modes honour the chosen one.
  if(m==="video_plain")fd.append("image_provider","text");
  else if($("imageProvider").value)fd.append("image_provider",$("imageProvider").value);
  if(deckOut){ // teaser: always replace speaker notes (when the server has an LLM)
   const notesVal=demo?(cfg.llm?"all":""):$("notes").value;
   if(notesVal)fd.append("notes",notesVal);
  }
  let res;
  if(!demo){
   if(deck&&!deck.isPptx){ // persist any last edit
    await fetch("/api/projects/"+projectId+"/deck",{method:"PUT",
     headers:Object.assign({"Content-Type":"application/json"},pauth()),
     body:JSON.stringify({markdown:$("deckEditor").value})});
   }
   const path=deckOut?"/api/projects/"+projectId+"/enrich"
                     :"/api/projects/"+projectId+"/render";
   res=await fetch(path,{method:"POST",headers:pauth(),body:fd});
  }else{
   fd.append("deck",deckBlob(),deckName());
   fd.append("output",deckOut?"pptx":"video");
   res=await fetch("/api/jobs",{method:"POST",headers:auth(),body:fd});
  }
  if(!res.ok){$("status").textContent="Error: "+(await res.text());
   $("go").disabled=false;return}
  const {job_id,token,notice}=await res.json();
  if(notice){$("notice").textContent="ℹ️ "+notice;$("notice").style.display="block"}
  startProgress();poll(job_id,token,deckOut?"pptx":"video")}
 catch(e){$("status").textContent="Error: "+e.message;$("go").disabled=false}};

// --- Job polling: status badge + elapsed time + gentle stage hints -----------
// Renders take minutes (an AI image + narration per slide), so show that
// time is passing and roughly where the job tends to be. The hints are
// heuristics — the server only reports a coarse status.
let finishedJob=false; // done-state: the primary button becomes "Start over"

const STAGES=[[0,"waiting for a worker slot…"],
 [8,"generating slides & narration…"],
 [45,"still working — AI images and talking heads take a while…"],
 [150,"almost there — encoding the final file…"]];
let t0=0,tick=null;
function fmtDur(sec){sec=Math.floor(sec);const m=Math.floor(sec/60);
 return m?m+"m "+(sec%60)+"s":sec+"s"}
function stageFor(sec){let msg=STAGES[0][1];
 STAGES.forEach(([s,m])=>{if(sec>=s)msg=m});return msg}
function renderProgress(){if(!t0)return;
 const sec=(Date.now()-t0)/1000;
 const el=document.querySelector("#status .elapsed");
 if(el)el.textContent=fmtDur(sec)+" · "+stageFor(sec)}
function startProgress(){t0=Date.now();renderProgress();
 if(tick)clearInterval(tick);tick=setInterval(renderProgress,1000)}
function stopProgress(){if(tick){clearInterval(tick);tick=null}}
function hideProg(){$("progWrap").style.display="none";
 $("progBar").style.width="0";$("progLabel").textContent=""}
async function poll(id,tok,kind){
 let j;
 try{const r=await fetch("/api/jobs/"+id,{headers:auth()});j=await r.json()}
 catch(e){ // transient network hiccup: keep polling rather than dying
  setTimeout(()=>poll(id,tok,kind),4000);return}
 const hasTrace=j.trace&&j.trace.length;
 $("logWrap").style.display=(j.log||hasTrace)?"block":"none";
 $("log").textContent=j.log||"";
 $("trace").innerHTML=hasTrace?j.trace.map(t=>"<div>"+esc(t)+"</div>").join(""):"";
 if(j.status==="done"){stopProgress();t0=0;$("go").disabled=false;hideProg();
  finishedJob=true;$("go").textContent="↻ Start over";
  if(j.warnings&&j.warnings.length){
   $("notice").textContent="⚠️ "+j.warnings.join(" ");$("notice").style.display="block"}
  const lbl=kind==="pptx"?"download deck (.zip)":"download video";
  $("status").innerHTML='<span class="badge">done</span> <a href="/api/jobs/'+id+
   '/result?t='+encodeURIComponent(tok||j.token||"")+'" download>⬇ '+lbl+"</a>";
  if(demo){
   const a=$("status").querySelector("a");
   if(a)a.onclick=()=>{$("status").innerHTML=
    '<span class="badge">done</span> <span class="elapsed">downloaded — the job is wiped (demo)</span>'};
   const doneAd=$("doneAdvert");doneAd.style.display="block";
   const doneDl=$("doneDl");if(doneDl&&$("dlBanner"))doneDl.href=$("dlBanner").href;
   const ttl=cfg.limits&&cfg.limits.job_ttl_minutes?cfg.limits.job_ttl_minutes:60;
   $("linkExpiry").textContent=
    "Download links work once and expire after about "+ttl+" minutes (demo).";
  }
  if(kind==="pptx"&&!demo&&projectId)showThumbs();return}
 if(j.status==="error"){stopProgress();t0=0;$("go").disabled=false;hideProg();
  $("status").textContent="Failed: "+(j.error||"see log");return}
 // Still working. Prefer real progress from the render log (slide x/y, stage,
 // "last update Ns ago"); fall back to elapsed-time heuristics before the
 // first slide line appears.
 const sec=t0?(Date.now()-t0)/1000:0;
 const p=j.progress;
 if(p&&(p.slide||p.stage)){
  $("progWrap").style.display="block";
  if(p.slide&&p.slides)
   $("progBar").style.width=
    Math.min(100,Math.round(100*(p.slide-0.5)/p.slides))+"%";
  const lbl=[];
  if(p.slide)lbl.push("slide "+p.slide+" of "+(p.slides||p.slide));
  if(p.stage)lbl.push(p.stage);
  if(typeof j.idle==="number")lbl.push(j.idle<15?"working…":"last update "+j.idle+"s ago");
  $("progLabel").textContent=lbl.join(" · ");
  $("status").innerHTML='<span class="badge">'+j.status+"</span>"+
   ' <span class="elapsed">'+fmtDur(sec)+"</span>";
 }else{
  $("status").innerHTML='<span class="badge">'+j.status+"</span>"+
   ' <span class="elapsed">'+fmtDur(sec)+" — "+stageFor(sec)+"</span>";
 }
 setTimeout(()=>poll(id,tok,kind),2500)}
// After an enrich, show the images it added to the project.
async function showThumbs(){
 try{const r=await fetch("/api/projects/"+projectId,{headers:pauth()});
  if(!r.ok)return;const st=await r.json();
  $("thumbs").innerHTML=(st.images||[]).map(n=>
   '<img alt="'+n+'" src="/api/projects/'+projectId+'/images/'+encodeURIComponent(n)+
   '?t='+encodeURIComponent(projectToken||"")+'">').join("")}
 catch(e){}}

// --- Settings (desktop/local mode): edit ~/.slidestream.yaml in-app ---------
let setTemplate="";
$("gear").onclick=async()=>{
 const r=await fetch("/api/settings");if(!r.ok)return;
 const s=await r.json();setTemplate=s.template||"";
 $("setYaml").value=s.yaml||"";$("setPath").textContent="Edits "+s.path+" — providers, servers, keys.";
 $("saveMsg").textContent="";$("settings").classList.add("on")};
$("setClose").onclick=()=>$("settings").classList.remove("on");
$("settings").onclick=e=>{if(e.target.id==="settings")$("settings").classList.remove("on")};
$("setTpl").onclick=()=>{if(!$("setYaml").value.trim()||confirm("Replace current contents with the template?"))
 $("setYaml").value=setTemplate};
$("setSave").onclick=async()=>{
 const r=await fetch("/api/settings",{method:"PUT",headers:{"Content-Type":"application/json"},
  body:JSON.stringify({yaml:$("setYaml").value})});
 const j=await r.json().catch(()=>({}));
 $("saveMsg").textContent=r.ok?"✓ Saved — applies to your next video.":("✗ "+(j.detail||"Save failed"));
 $("saveMsg").style.color=r.ok?"":"var(--accent)"};

// --- Upload drop zones --------------------------------------------------------
dropzone("draftFile","Drop a PDF / Word / PowerPoint / txt — or click to browse");
dropzone("deckFile","Drop a .md or .pptx — or click to browse");
dropzone("voice","Drop a voice clip (mp3, wav, m4a — even a video) — or click to browse");
dropzone("photo","Drop a photo or short video — or click to browse");
</script></body></html>
"""
