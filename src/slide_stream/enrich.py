"""Deck enrichment: add an image to each slide and write a new deck.

This is the ``enrich`` output track (ported/adapted from slide-vision): the
same slide input, but instead of a narrated video the output is a new,
editable deck — a Markdown file plus an ``images/`` folder, and optionally a
PowerPoint. Run ``create`` on the result to narrate it, or use ``create``
directly (it enriches internally) for a one-pass video.
"""

import zipfile
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path
from typing import Any

from rich.console import Console

from .narration import notes_narration_ready
from .providers.base import ImageProvider


def _slide_query(slide: dict[str, Any]) -> str:
    """A search/keyword query for a slide's image."""
    title = str(slide.get("title", "")).strip()
    if title:
        return title
    for item in slide.get("content", []):
        text = str(item).strip()
        if text:
            return text
    return "presentation slide"


def _generate_notes(slide: dict[str, Any], llm: dict[str, Any]) -> str:
    """AI presenter notes for one slide, reusing the narration writer so the
    notes read as a spoken script (and drive narration if the deck is later
    rendered — ``create`` reads .pptx speaker notes as its narration source).

    Always written from the slide's own content/title, never from any existing
    notes, so ``all`` mode genuinely regenerates rather than paraphrasing.
    """
    from rich.console import Console

    from .llm import query_llm
    from .narration import build_narration_prompt, target_words

    source = "content" if slide.get("content") else "title"
    wpm = llm.get("wpm", 150)
    words = target_words(llm.get("target_seconds"), wpm)
    prompt = build_narration_prompt(slide, source, words, wpm)
    text = query_llm(
        llm["client"], llm["provider"], prompt, Console(), llm.get("model"),
        think=llm.get("think"),
    )
    return (text or "").strip()


def enrich_deck(
    slides: list[dict[str, Any]],
    image_provider: ImageProvider,
    output_dir: Path,
    input_stem: str,
    *,
    also_pptx: bool = False,
    also_zip: bool = False,
    notes_mode: str | None = None,
    llm: dict[str, Any] | None = None,
) -> Path:
    """Write an enriched Markdown deck (and optional PPTX) into output_dir.

    Returns the output directory. Each slide gets an image from
    ``image_provider``; slides the ``local`` provider could not match are
    listed in ``prompts.md`` with ready-to-paste AI-image prompts.

    ``notes_mode`` adds presenter notes to the PowerPoint (requires
    ``also_pptx`` and an ``llm`` context — client/provider/model/wpm/
    target_seconds):
      - ``fill``: keep a slide's existing speaker notes; AI-write notes only
        for slides that have none.
      - ``all``: AI-write notes for every slide, replacing any existing ones.
    """
    if notes_mode and llm is None:
        raise ValueError("notes_mode requires an 'llm' context")

    images_dir = output_dir / "images"
    images_dir.mkdir(parents=True, exist_ok=True)

    console = Console()
    total = len(slides)
    enriched: list[dict[str, Any]] = []
    # The image server and the notes LLM are separate backends, so within a
    # slide they run CONCURRENTLY (wall time = max, not sum). Slides themselves
    # stay sequential: one slide in flight keeps the log trace accurate.
    with ThreadPoolExecutor(max_workers=2) as pool:
        for i, slide in enumerate(slides, 1):
            title = str(slide.get("title", "")).strip()
            # Same parseable per-slide marker the create flow prints — the web
            # UI's job trace and progress parser key off this exact text.
            console.print(f"Slide {i}/{total}: {title}")
            img_path = images_dir / f"slide_{i}.png"

            existing_notes = str(slide.get("notes", "")).strip()
            keep_existing = (
                notes_mode == "fill" and notes_narration_ready(existing_notes)
            )
            will_write = bool(notes_mode) and not keep_existing

            fb_before = getattr(image_provider, "fallback_count", 0)
            image_future = pool.submit(
                image_provider.generate_image,
                _slide_query(slide), str(img_path), slide=slide,
            )
            notes_future = (
                pool.submit(_generate_notes, slide, llm or {}) if will_write else None
            )
            image_future.result()
            # A provider fallback wrote a generic text card — worthless in an
            # editable deck (it duplicates the bullet text as a picture), so
            # the slide is left text-only instead and listed in prompts.md.
            fell_back = (
                getattr(image_provider, "fallback_count", fb_before) > fb_before
            )
            # The local provider reports whether it matched a real folder
            # image; other providers always produce an image (or their own
            # text fallback).
            matched = (
                getattr(image_provider, "matched_last", True) and not fell_back
            )
            if fell_back:
                img_path.unlink(missing_ok=True)
                console.print("  - Image generation failed — slide left text-only")

            wrote_notes = False
            notes = ""
            if notes_future is not None:
                notes = notes_future.result() or ""
                wrote_notes = True
            elif notes_mode == "fill":
                notes = existing_notes
            if wrote_notes:
                console.print("  - Wrote AI speaker notes")

            enriched.append(
                {
                    "index": i,
                    "title": title,
                    "content": [str(c).strip() for c in slide.get("content", []) if str(c).strip()],
                    "image": img_path.name,
                    "matched": matched,
                    "notes": notes,
                }
            )

    md_path = output_dir / f"{input_stem}.md"
    md_path.write_text(_build_markdown(enriched), encoding="utf-8")

    missing = [s for s in enriched if not s["matched"]]
    if missing:
        (output_dir / "prompts.md").write_text(_build_prompts(missing), encoding="utf-8")

    if also_pptx:
        _write_pptx(enriched, images_dir, output_dir / f"{input_stem}.pptx")

    if also_zip:
        zip_path = output_dir.parent / f"{output_dir.name}.zip"
        with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
            for f in output_dir.rglob("*"):
                if f.is_file():
                    zf.write(f, f.relative_to(output_dir.parent))

    return output_dir


def _build_markdown(slides: list[dict[str, Any]]) -> str:
    blocks: list[str] = []
    for slide in slides:
        lines = [f"# {slide['title']}" if slide["title"] else f"# Slide {slide['index']}"]
        lines.append("")
        if slide.get("matched", True):
            lines.append(f"![{slide['title']}](images/{slide['image']})")
        if slide["content"]:
            lines.append("")
            for item in slide["content"]:
                lines.append(f"- {item}")
        blocks.append("\n".join(lines))
    return "\n\n---\n\n".join(blocks) + "\n"


def _build_prompts(missing: list[dict[str, Any]]) -> str:
    lines = ["# Image Prompts", "",
             "Slides with no matching local image. Paste a prompt into an AI "
             "image tool (DALL-E, Midjourney, ...) and drop the result into "
             "the images/ folder.", ""]
    for slide in missing:
        preview = " ".join(slide["content"])[:300]
        lines += [
            f"## Slide {slide['index']}: {slide['title']}",
            "",
            f'A high-quality, professional illustration for a presentation slide '
            f'titled "{slide["title"]}".',
        ]
        if preview:
            lines.append(f"The slide covers: {preview}.")
        lines += ["Style: clean, modern, no text overlays.", "", "---", ""]
    return "\n".join(lines)


def _write_pptx(slides: list[dict[str, Any]], images_dir: Path, out_path: Path) -> None:
    """Build a PowerPoint deck that stays EDITABLE: the educator's title and
    bullet points are real text; the generated image sits beside them as a
    complement. Speaker notes (from --notes) ride along as usual."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    blank = prs.slide_layouts[6]
    for slide in slides:
        s = prs.slides.add_slide(blank)
        # Title across the top — real, editable text.
        title_box = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.9))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.text = slide["title"] or f"Slide {slide['index']}"
        tf.paragraphs[0].runs[0].font.size = Pt(30)

        # The educator's bullet points — the core of the slide — as editable
        # text on the left.
        bullets_box = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.6), Inches(5.4))
        btf = bullets_box.text_frame
        btf.word_wrap = True
        content = [str(c).strip() for c in slide.get("content", []) if str(c).strip()]
        for j, bullet in enumerate(content):
            para = btf.paragraphs[0] if j == 0 else btf.add_paragraph()
            para.text = f"• {bullet}"
            para.font.size = Pt(18)
            if j:
                para.space_before = Pt(8)

        # Generated image on the right as a complement; width fixed, height
        # keeps the source aspect ratio.
        img = images_dir / slide["image"]
        if img.is_file():
            s.shapes.add_picture(str(img), Inches(5.4), Inches(1.9), width=Inches(4.1))

        # Presenter notes (added by --notes). create reads these back as the
        # narration source, so an enriched .pptx round-trips into a video.
        note = str(slide.get("notes", "")).strip()
        if note:
            text_frame = s.notes_slide.notes_text_frame
            if text_frame is not None:
                text_frame.text = note
    prs.save(str(out_path))
