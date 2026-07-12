"""PowerPoint parsing functionality for Slide Stream."""

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER


def parse_powerpoint(file_path: str | Path) -> list[dict[str, Any]]:
    """Parse PowerPoint file into slide data.

    Extracts slide titles, bullet points, and speaker notes from .pptx files.

    Args:
        file_path: Path to the PowerPoint file (.pptx)

    Returns:
        List of slide dictionaries with 'title', 'content', 'notes', 'images',
        and 'has_real_title' keys. 'has_real_title' is True when a genuine
        title was found on the slide and False when the "Slide N" default was
        used, so downstream code can tell placeholders apart reliably.
    """
    try:
        prs = Presentation(str(file_path))
    except Exception as e:
        raise ValueError(f"Could not open PowerPoint file: {e}") from e

    slides = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_data: dict[str, Any] = {
            "title": f"Slide {slide_num}",  # Default title
            "content": [],
            "notes": "",
            "images": [],
            "has_real_title": False,
        }

        # First pass: look for a genuine TITLE/CENTER_TITLE placeholder, so a
        # caption textbox earlier in shape (z-)order cannot steal the title.
        title_shape = None
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if not shape.is_placeholder:  # type: ignore[attr-defined]
                continue
            ph_type = shape.placeholder_format.type  # type: ignore[attr-defined]
            if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                if shape.text.strip():  # type: ignore[attr-defined]
                    title_shape = shape
                    break

        # Extract slide content
        title_set = False
        for shape in slide.shapes:
            # Collect embedded pictures so image-only slides can still be
            # narrated (via a vision-capable LLM).
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    slide_data["images"].append({
                        "data": shape.image.blob,  # type: ignore[attr-defined]
                        "content_type": shape.image.content_type,  # type: ignore[attr-defined]
                    })
                except Exception:
                    pass  # unsupported/linked image formats: skip silently
                continue
            if not shape.has_text_frame:
                continue
            text = shape.text.strip()  # type: ignore[attr-defined]
            if not text:
                continue

            # The title placeholder found in the first pass wins regardless of
            # shape order; the heuristic (first short, single-line text) only
            # applies to layouts with no populated title slot. (Compare with
            # ==: shape proxies are rebuilt on each iteration, so identity
            # would fail; BaseShape equality compares the underlying element.)
            if title_shape is not None and shape == title_shape:
                slide_data["title"] = text
                title_set = True
                continue

            if (
                title_shape is None
                and not title_set
                and len(text) < 100
                and "\n" not in text
            ):
                slide_data["title"] = text
                title_set = True
                continue

            # Otherwise this is body content: extract bullet points / paragraphs.
            for paragraph in shape.text_frame.paragraphs:  # type: ignore[attr-defined]
                para_text = paragraph.text.strip()
                if para_text:
                    slide_data["content"].append(para_text)

        # Extract speaker notes. notes_text_frame resolves the BODY
        # placeholder, so header/footer/slide-number placeholders on the notes
        # slide cannot masquerade as notes; fall back to scanning shapes only
        # when no notes text frame is available.
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            notes_frame = notes_slide.notes_text_frame
            if notes_frame is not None:
                notes_text = notes_frame.text.strip()
                if notes_text and notes_text != "Click to add notes":
                    slide_data["notes"] = notes_text
            else:
                for shape in notes_slide.shapes:
                    if shape.has_text_frame:
                        notes_text = shape.text.strip()  # type: ignore[attr-defined]
                        if notes_text and notes_text != "Click to add notes":
                            slide_data["notes"] = notes_text
                            break

        # Keep any slide with body content, speaker notes, images, or a real
        # (non-default) title, so title-only / divider / image-only slides are
        # not silently dropped.
        slide_data["has_real_title"] = title_set
        if (
            slide_data["content"]
            or slide_data["notes"]
            or slide_data["images"]
            or title_set
        ):
            slides.append(slide_data)

    return slides


def format_powerpoint_content_for_llm(slide: dict[str, Any]) -> str:
    """Format PowerPoint slide data for LLM processing.

    Combines title, content, and speaker notes into a coherent prompt.

    Args:
        slide: Slide dictionary with title, content, and notes

    Returns:
        Formatted string for LLM input
    """
    parts = [f"Title: {slide['title']}"]

    if slide["content"]:
        parts.append("Content:")
        for item in slide["content"]:
            parts.append(f"• {item}")

    if slide["notes"]:
        parts.append(f"Speaker Notes: {slide['notes']}")

    return "\n".join(parts)
